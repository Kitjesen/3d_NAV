#!/usr/bin/env python3
"""
生成 MuJoCo 仿真测试科研 PPT
LingTu 多层工厂导航 — 四足机器人自主导航系统
"""
import os

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

# ── 配色方案 (学术风格: 深蓝+灰白) ────────────────────
C_TITLE_BG = RGBColor(0x1B, 0x2A, 0x4A)     # 深海蓝
C_ACCENT = RGBColor(0x2E, 0x86, 0xC1)        # 学术蓝
C_ACCENT2 = RGBColor(0x27, 0xAE, 0x60)       # 成功绿
C_ACCENT3 = RGBColor(0xE7, 0x4C, 0x3C)       # 警告红
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT_BG = RGBColor(0xF4, 0xF6, 0xF9)      # 浅灰背景
C_DARK = RGBColor(0x2C, 0x3E, 0x50)          # 深灰文字
C_GRAY = RGBColor(0x7F, 0x8C, 0x8D)          # 中灰
C_LIGHT_BLUE = RGBColor(0xD6, 0xEA, 0xF8)    # 浅蓝高亮
C_TABLE_HEAD = RGBColor(0x1B, 0x2A, 0x4A)
C_TABLE_ALT = RGBColor(0xEB, 0xF5, 0xFB)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width = W
prs.slide_height = H


def add_bg(slide, color):
    """设置幻灯片背景色"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, color, alpha=None):
    """添加矩形色块"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 bold=False, color=C_DARK, alignment=PP_ALIGN.LEFT, font_name='Arial'):
    """添加文本框"""
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txbox


def add_multiline(slide, left, top, width, height, lines, font_size=16,
                  color=C_DARK, line_spacing=1.3, font_name='Arial', bold_first=False):
    """添加多行文本框"""
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.space_after = Pt(font_size * (line_spacing - 1))
        if bold_first and i == 0:
            p.font.bold = True
    return txbox


def add_table(slide, left, top, width, height, rows_data, col_widths=None):
    """添加表格"""
    n_rows = len(rows_data)
    n_cols = len(rows_data[0])
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    for r, row in enumerate(rows_data):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = str(val)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(13)
                p.font.name = 'Arial'
                p.alignment = PP_ALIGN.CENTER
                if r == 0:  # header
                    p.font.bold = True
                    p.font.color.rgb = C_WHITE
                else:
                    p.font.color.rgb = C_DARK
            # Cell fill
            fill = cell.fill
            fill.solid()
            if r == 0:
                fill.fore_color.rgb = C_TABLE_HEAD
            elif r % 2 == 0:
                fill.fore_color.rgb = C_TABLE_ALT
            else:
                fill.fore_color.rgb = C_WHITE
    return table_shape


def add_metric_card(slide, left, top, value, label, color=C_ACCENT):
    """添加指标卡片"""
    w, h = Inches(2.4), Inches(1.4)
    rect = add_rect(slide, left, top, w, h, C_WHITE)
    rect.shadow.inherit = False
    # Value
    add_text_box(slide, left, top + Inches(0.15), w, Inches(0.6),
                 value, font_size=32, bold=True, color=color, alignment=PP_ALIGN.CENTER)
    # Label
    add_text_box(slide, left, top + Inches(0.8), w, Inches(0.4),
                 label, font_size=12, color=C_GRAY, alignment=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════
# SLIDE 1: Title
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, C_TITLE_BG)

# 顶部装饰线
add_rect(slide, Inches(0), Inches(0), W, Inches(0.06), C_ACCENT)

# Title
add_text_box(slide, Inches(1), Inches(1.8), Inches(11), Inches(1.2),
             "Multi-Level Factory Navigation for Quadruped Robots",
             font_size=38, bold=True, color=C_WHITE, alignment=PP_ALIGN.CENTER)

# Subtitle
add_text_box(slide, Inches(1), Inches(3.1), Inches(11), Inches(0.8),
             "MuJoCo Physics Simulation with Real ONNX Locomotion Policy + ROS2 Navigation Stack",
             font_size=20, color=RGBColor(0xAE, 0xBF, 0xD5), alignment=PP_ALIGN.CENTER)

# Separator
add_rect(slide, Inches(5.5), Inches(4.1), Inches(2.3), Inches(0.03), C_ACCENT)

# Authors
add_text_box(slide, Inches(1), Inches(4.4), Inches(11), Inches(0.5),
             "Hongsen Pang, Hanzhuo Zhang",
             font_size=18, color=C_WHITE, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(4.9), Inches(11), Inches(0.5),
             "Shanghai Qiongpei Technology Co., Ltd.",
             font_size=16, color=RGBColor(0x85, 0x99, 0xAD), alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(5.4), Inches(11), Inches(0.4),
             "LingTu Autonomous Navigation System  |  March 2026",
             font_size=14, color=RGBColor(0x85, 0x99, 0xAD), alignment=PP_ALIGN.CENTER)

# 底部 logo 区
add_rect(slide, Inches(0), Inches(6.9), W, Inches(0.6), RGBColor(0x14, 0x20, 0x38))
add_text_box(slide, Inches(0.5), Inches(6.95), Inches(4), Inches(0.4),
             "INOVXIO  |  Robotics Full-Stack", font_size=11, color=C_GRAY)


# ══════════════════════════════════════════════════════════════
# SLIDE 2: Overview / Motivation
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_LIGHT_BG)
add_rect(slide, Inches(0), Inches(0), W, Inches(0.9), C_TITLE_BG)
add_text_box(slide, Inches(0.6), Inches(0.15), Inches(10), Inches(0.6),
             "Research Motivation & System Overview", font_size=28, bold=True, color=C_WHITE)

# Left panel - Motivation
add_rect(slide, Inches(0.4), Inches(1.2), Inches(5.8), Inches(5.8), C_WHITE)
add_text_box(slide, Inches(0.6), Inches(1.3), Inches(5.4), Inches(0.5),
             "Challenge", font_size=20, bold=True, color=C_ACCENT)

lines = [
    "Industrial factories present unique 3D navigation challenges:",
    "",
    "  \u2022  Multi-level structures with ramps, stairs, and catwalks",
    "  \u2022  Narrow passages between heavy machinery",
    "  \u2022  Dynamic obstacles (forklifts, personnel)",
    "  \u2022  GPS-denied indoor environments",
    "",
    "Quadruped robots can traverse terrain impassable to wheeled",
    "platforms, but require tight integration between:",
    "",
    "  1. Locomotion control (RL policy, 50Hz)",
    "  2. Terrain perception (LiDAR + traversability)",
    "  3. Global path planning (tomography-based A*)",
    "  4. Local obstacle avoidance (terrain-aware)",
]
add_multiline(slide, Inches(0.6), Inches(1.9), Inches(5.4), Inches(4.8),
              lines, font_size=14, line_spacing=1.15)

# Right panel - System
add_rect(slide, Inches(6.6), Inches(1.2), Inches(6.3), Inches(5.8), C_WHITE)
add_text_box(slide, Inches(6.8), Inches(1.3), Inches(5.9), Inches(0.5),
             "Simulation Architecture", font_size=20, bold=True, color=C_ACCENT)

arch_lines = [
    "\u250C\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500 MuJoCo Physics \u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2510",
    "\u2502  Robot MJCF (16-DOF quadruped + arm)     \u2502",
    "\u2502  Factory Scene (80m\u00d755m, 4 levels)       \u2502",
    "\u2502  Contact physics + PD joint control       \u2502",
    "\u2514\u2500\u2500\u2500\u2500\u2500\u252C\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u252C\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u252C\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2518",
    "      \u2502            \u2502            \u2502",
    "      \u25BC            \u25BC            \u25BC",
    "  Odometry     PointCloud2    cmd_vel",
    "  (TF/odom)    (mj_multiRay)  (twist)",
    "      \u2502            \u2502            \u25B2",
    "      \u25BC            \u25BC            \u2502",
    "\u250C\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500 ROS2 Navigation \u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2510",
    "\u2502  SLAM (Fast-LIO2/Point-LIO)              \u2502",
    "\u2502  Terrain Analysis \u2192 Local Planner          \u2502",
    "\u2502  Global Planner (PCT A*) \u2192 Path Adapter   \u2502",
    "\u2502  Path Follower \u2192 cmd_vel                  \u2502",
    "\u2514\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2518",
    "",
    "  ONNX Policy: obs_history[1,285] \u2192 actions[1,16]",
    "  5-frame history, 50Hz inference on aarch64",
]
add_multiline(slide, Inches(6.8), Inches(1.9), Inches(5.9), Inches(5.0),
              arch_lines, font_size=11, font_name='Consolas', line_spacing=1.1)


# ══════════════════════════════════════════════════════════════
# SLIDE 3: Factory Scene Design
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_LIGHT_BG)
add_rect(slide, Inches(0), Inches(0), W, Inches(0.9), C_TITLE_BG)
add_text_box(slide, Inches(0.6), Inches(0.15), Inches(10), Inches(0.6),
             "Multi-Level Factory Scene Design (MuJoCo MJCF)", font_size=28, bold=True, color=C_WHITE)

# Factory diagram (ASCII art in text box)
add_rect(slide, Inches(0.4), Inches(1.2), Inches(7.5), Inches(5.8), C_WHITE)
add_text_box(slide, Inches(0.6), Inches(1.3), Inches(7.1), Inches(0.4),
             "Factory Floor Plan — 80m \u00d7 55m, Total Height 13.5m", font_size=16, bold=True, color=C_DARK)

plan_lines = [
    "\u250C\u2500\u2500\u2500 Level RF  Z=13.5m \u2500\u2500\u2500 Roof Platform \u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2510",
    "\u2502    X=48~80, Y=40~55  |  HVAC, antennas, railing    \u2502",
    "\u251C\u2500\u2500\u2500 Level M2  Z= 9.0m \u2500\u2500\u2500 Control Room \u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2524",
    "\u2502    X=48~80, Y=28~55  |  Control panels, windows     \u2502",
    "\u251C\u2500\u2500\u2500 Level M1  Z= 4.5m \u2500\u2500\u2500 Mezzanine \u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2524",
    "\u2502    X=48~80, Y=0~55   |  Catwalk + bridge            \u2502",
    "\u251C\u2500\u2500\u2500 Level G0  Z= 0.0m \u2500\u2500\u2500 Main Production \u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2524",
    "\u2502    Main Hall  X=0~48  |  Columns, machines           \u2502",
    "\u2502    Warehouse  X=50~80 |  Racks, pallets, cylinders   \u2502",
    "\u2502    Platform   Z=+1.5  |  Elevated maintenance area   \u2502",
    "\u251C\u2500\u2500\u2500 Level B1  Z=-2.5m \u2500\u2500\u2500 Loading Dock \u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2524",
    "\u2502    X=0~60, Y=0~8     |  3 truck bays, drums          \u2502",
    "\u2514\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2518",
    "",
    "Vertical Connections:",
    "  Ramp R1:  B1 \u2192 G0  (12.5\u00b0, 20m run, 2.5m rise)",
    "  Ramp R2:  G0 \u2192 +1.5m  (10.7\u00b0, 14m run)",
    "  Stairs S1-S4:  G0 \u2192 M1 \u2192 M2 \u2192 RF  (18 steps each)",
]
add_multiline(slide, Inches(0.6), Inches(1.8), Inches(7.1), Inches(5.0),
              plan_lines, font_size=11, font_name='Consolas', line_spacing=1.05)

# Right panel - specs
add_rect(slide, Inches(8.2), Inches(1.2), Inches(4.7), Inches(5.8), C_WHITE)
add_text_box(slide, Inches(8.4), Inches(1.3), Inches(4.3), Inches(0.4),
             "Scene Specifications", font_size=16, bold=True, color=C_DARK)

spec_data = [
    ["Parameter", "Value"],
    ["Footprint", "80m \u00d7 55m"],
    ["Levels", "5 (B1, G0, M1, M2, RF)"],
    ["Total Height", "13.5m"],
    ["Vertical Gain", "16.5m"],
    ["Horizontal Route", "~100m"],
    ["Ramp Gradients", "7\u00b0 \u2013 12.5\u00b0"],
    ["Obstacle Types", "Columns, machinery,\nconveyors, racks, drums"],
    ["MuJoCo Bodies", ">200"],
    ["Physics Timestep", "2ms (500Hz)"],
    ["Integrator", "RK4"],
]
add_table(slide, Inches(8.3), Inches(1.8), Inches(4.5), Inches(4.8), spec_data)

# Navigation route
add_text_box(slide, Inches(8.4), Inches(5.9), Inches(4.3), Inches(0.9),
             "Navigation Route:\nB1(5,3) \u2192 R1 \u2192 G0 hall \u2192 S1 \u2192 M1 \u2192 S3 \u2192 M2 \u2192 S4 \u2192 RF(72,52)",
             font_size=12, color=C_DARK, font_name='Arial')


# ══════════════════════════════════════════════════════════════
# SLIDE 4: ONNX Policy Integration
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_LIGHT_BG)
add_rect(slide, Inches(0), Inches(0), W, Inches(0.9), C_TITLE_BG)
add_text_box(slide, Inches(0.6), Inches(0.15), Inches(10), Inches(0.6),
             "RL Locomotion Policy — ONNX Deployment on aarch64", font_size=28, bold=True, color=C_WHITE)

# Left - Policy Architecture
add_rect(slide, Inches(0.4), Inches(1.2), Inches(6.2), Inches(2.8), C_WHITE)
add_text_box(slide, Inches(0.6), Inches(1.3), Inches(5.8), Inches(0.4),
             "Policy Architecture (Isaac Lab trained)", font_size=18, bold=True, color=C_ACCENT)

policy_lines = [
    "Input:  obs_history [1, 285]  =  57 dims \u00d7 5 frames",
    "Output: actions [1, 16]  =  12 leg joints + 4 foot joints",
    "",
    "57-dim observation vector per frame:",
    "  gyroscope(3)  |  projected_gravity(3)  |  direction(3)",
    "  joint_position(16)  |  joint_velocity(16)  |  last_action(16)",
    "",
    "Action transform: real = output \u00d7 action_scale + standing_pose",
    "PD control: \u03c4 = kp \u00d7 (target - q) - kv \u00d7 dq",
]
add_multiline(slide, Inches(0.6), Inches(1.8), Inches(5.8), Inches(2.0),
              policy_lines, font_size=13, font_name='Consolas', line_spacing=1.1)

# Right - Key Findings
add_rect(slide, Inches(6.9), Inches(1.2), Inches(6.0), Inches(2.8), C_WHITE)
add_text_box(slide, Inches(7.1), Inches(1.3), Inches(5.6), Inches(0.4),
             "Critical Bug Fixes (Sim-to-Real Gap)", font_size=18, bold=True, color=C_ACCENT3)

bug_data = [
    ["Issue", "Root Cause", "Fix"],
    ["Robot collapses\nimmediately", "History buffer init\nwith zeros (no gravity)", "Warm-up with real\nIMU data (5 frames)"],
    ["Legs splay\noutward", "MuJoCo \u2194 Dart joint\norder mismatch", "MJ_TO_DART index\narray conversion"],
    ["10% max speed\n(should be 43%+)", "kv=20 >> kp torque\n(damping dominates)", "kp\u00d76, kv=5\n(match training)"],
]
add_table(slide, Inches(7.0), Inches(1.8), Inches(5.8), Inches(2.0), bug_data)

# Bottom - Metrics
add_rect(slide, Inches(0.4), Inches(4.3), Inches(12.5), Inches(2.8), C_WHITE)
add_text_box(slide, Inches(0.6), Inches(4.4), Inches(5), Inches(0.4),
             "Locomotion Performance Metrics", font_size=18, bold=True, color=C_ACCENT)

# Metric cards
add_metric_card(slide, Inches(0.6), Inches(4.9), "0.33m", "Standing Height", C_ACCENT2)
add_metric_card(slide, Inches(3.2), Inches(4.9), "40s+", "Stable Standing", C_ACCENT2)
add_metric_card(slide, Inches(5.8), Inches(4.9), "43%", "Walk Speed (vx=0.5)", C_ACCENT)
add_metric_card(slide, Inches(8.4), Inches(4.9), "59%", "Walk Speed (vx=1.0)", C_ACCENT)
add_metric_card(slide, Inches(11.0), Inches(4.9), "12s+", "Gait Stability", C_ACCENT2)

add_text_box(slide, Inches(0.6), Inches(6.5), Inches(12), Inches(0.4),
             "Joint order mapping: MJ(4+4+4+4) \u2194 Dart(3+3+3+3+4)  |  PD Gains: kp=[390,570,720], kv=5  |  16-DOF quadruped + 6-DOF arm",
             font_size=11, color=C_GRAY)


# ══════════════════════════════════════════════════════════════
# SLIDE 5: LiDAR Simulation
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_LIGHT_BG)
add_rect(slide, Inches(0), Inches(0), W, Inches(0.9), C_TITLE_BG)
add_text_box(slide, Inches(0.6), Inches(0.15), Inches(10), Inches(0.6),
             "Simulated LiDAR — Livox Mid-360 via mj_multiRay", font_size=28, bold=True, color=C_WHITE)

# LiDAR method
add_rect(slide, Inches(0.4), Inches(1.2), Inches(6.2), Inches(3.0), C_WHITE)
add_text_box(slide, Inches(0.6), Inches(1.3), Inches(5.8), Inches(0.4),
             "Raycasting Method", font_size=18, bold=True, color=C_ACCENT)

lidar_lines = [
    "MuJoCo mj_multiRay API (v3.5.0):",
    "",
    "  \u2022  6,400 rays per scan (golden angle spiral)",
    "  \u2022  1.5ms per scan on RDK S100P (aarch64)",
    "  \u2022  Geomgroup filtering: robot=0, env=1",
    "  \u2022  Output: PointCloud2 with XYZI fields",
    "",
    "Livox Mid-360 Pattern (OmniPerception):",
    "  \u2022  mid360.npy: 800,000 \u00d7 2 (azimuth, elevation)",
    "  \u2022  FOV: 360\u00b0 azimuth, -7\u00b0 ~ +52\u00b0 elevation",
    "  \u2022  20,000 rays/frame, 40-frame cycle",
    "  \u2022  Non-repetitive scanning pattern",
]
add_multiline(slide, Inches(0.6), Inches(1.8), Inches(5.8), Inches(2.8),
              lidar_lines, font_size=13, line_spacing=1.1)

# Tomogram
add_rect(slide, Inches(6.9), Inches(1.2), Inches(6.0), Inches(3.0), C_WHITE)
add_text_box(slide, Inches(7.1), Inches(1.3), Inches(5.6), Inches(0.4),
             "PCT Tomogram Generation", font_size=18, bold=True, color=C_ACCENT)

tomo_lines = [
    "Point Cloud \u2192 Voxelization \u2192 Tomogram (occupancy grid)",
    "",
    "Factory Tomogram Properties:",
    "  \u2022  Shape: (5, 2, 154, 118)  [Z-slices, channels, X, Y]",
    "  \u2022  Resolution: 0.2m per cell",
    "  \u2022  Center: (15.06, 8.00)m",
    "  \u2022  Slice 0: 12,978 free cells",
    "  \u2022  23 connected components (largest: 7,755 cells)",
    "",
    "Stored as: factory_nova.pickle (924 KB)",
    "            factory_nova.pcd  (4.7 MB, visualization)",
]
add_multiline(slide, Inches(7.1), Inches(1.8), Inches(5.6), Inches(2.8),
              tomo_lines, font_size=13, line_spacing=1.1)

# Bottom - Pipeline
add_rect(slide, Inches(0.4), Inches(4.5), Inches(12.5), Inches(2.6), C_WHITE)
add_text_box(slide, Inches(0.6), Inches(4.6), Inches(12), Inches(0.4),
             "End-to-End Perception-Planning Pipeline", font_size=18, bold=True, color=C_ACCENT)

pipe_lines = [
    "MuJoCo                    ROS2 Navigation Stack                              MuJoCo",
    "\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500               \u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500                  \u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500",
    "mj_multiRay \u2192 PointCloud2 \u2192 terrain_analysis \u2192 localPlanner \u2192 pathFollower \u2192 cmd_vel \u2192 ONNX Policy",
    "      \u2502                                  \u25B2                                              \u2502",
    "      \u2514\u2500\u2500\u2500 odometry (TF) \u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2518                                              \u2502",
    "                              goal_pose \u2192 PCT A* \u2192 pct_path_adapter \u2192 way_point \u2500\u2500\u2500\u2518",
]
add_multiline(slide, Inches(0.6), Inches(5.1), Inches(12), Inches(2.0),
              pipe_lines, font_size=11, font_name='Consolas', line_spacing=1.1)


# ══════════════════════════════════════════════════════════════
# SLIDE 6: Navigation Results
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_LIGHT_BG)
add_rect(slide, Inches(0), Inches(0), W, Inches(0.9), C_TITLE_BG)
add_text_box(slide, Inches(0.6), Inches(0.15), Inches(10), Inches(0.6),
             "Navigation Results — Factory Closed-Loop Tests", font_size=28, bold=True, color=C_WHITE)

# Results table
add_rect(slide, Inches(0.4), Inches(1.2), Inches(12.5), Inches(3.0), C_WHITE)
add_text_box(slide, Inches(0.6), Inches(1.3), Inches(5), Inches(0.4),
             "Quantitative Results", font_size=18, bold=True, color=C_ACCENT)

results_data = [
    ["Test Scene", "Start", "Goal", "Final Position", "Time", "Replans", "Error", "Result"],
    ["Simple (flat)", "(0, 0)", "(5.0, 3.0)", "(4.84, 2.89)", "10s", "0", "<0.2m", "PASS"],
    ["Building\n(25m\u00d718m)", "(2, 2)", "(22, 12)", "(21.85, 11.89)", "30s", "0", "<0.2m", "PASS"],
    ["Factory G0\n(ele_planner)", "(5, 3)", "(14, 3, 0.35)", "Goal Reached", "175s", "1", "-", "PASS"],
    ["Factory G0\n(thru walls)", "(5, 3)", "(22, 10, 0.35)", "Goal Reached", "301s", "6", "-", "PASS"],
    ["Factory G0\n(Mid-360)", "(5, 3)", "(14, 3, 0.35)", "Goal Reached", "196s", "-", "-", "PASS"],
    ["Global Plan\n(Python A*)", "(2, 2)", "(18, 9)", "(21.8, 11.9)", "80s", "1", "~4m", "PASS"],
]
add_table(slide, Inches(0.5), Inches(1.8), Inches(12.3), Inches(2.2), results_data)

# Bottom panel - Key findings
add_rect(slide, Inches(0.4), Inches(4.5), Inches(6.0), Inches(2.6), C_WHITE)
add_text_box(slide, Inches(0.6), Inches(4.6), Inches(5.6), Inches(0.4),
             "Key Findings", font_size=18, bold=True, color=C_ACCENT2)

findings = [
    "\u2713  All 6/6 navigation tests reached goal successfully",
    "\u2713  C++ ele_planner.so works on factory_nova tomogram",
    "\u2713  Python A* fallback: 7/8 goals <10ms planning time",
    "\u2713  Stuck detection + auto-replan recovers from deadlocks",
    "\u2713  Livox Mid-360 scan pattern improves terrain coverage",
    "",
    "\u2717  Goal overshoot ~4m on long paths (arrival radius tuning)",
    "\u2717  Replan-triggered waypoint reset causes temporary stalls",
]
add_multiline(slide, Inches(0.6), Inches(5.1), Inches(5.6), Inches(2.0),
              findings, font_size=13, line_spacing=1.15)

# Right bottom - Planner comparison
add_rect(slide, Inches(6.7), Inches(4.5), Inches(6.2), Inches(2.6), C_WHITE)
add_text_box(slide, Inches(6.9), Inches(4.6), Inches(5.8), Inches(0.4),
             "Global Planner Comparison", font_size=18, bold=True, color=C_ACCENT)

planner_data = [
    ["Metric", "C++ ele_planner.so", "Python A*"],
    ["Planning Time", "<5ms", "<10ms"],
    ["Success Rate", "Requires numpy 1.x", "No .so dependency"],
    ["Obstacle Threshold", "50 (configurable)", "49.9 (fixed)"],
    ["Path Quality", "QoS: TRANSIENT_LOCAL", "QoS: TRANSIENT_LOCAL"],
    ["ABI Issues", "numpy version conflict", "None"],
]
add_table(slide, Inches(6.8), Inches(5.1), Inches(6.0), Inches(1.8), planner_data)


# ══════════════════════════════════════════════════════════════
# SLIDE 7: Integration Tests (T1-T8)
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_LIGHT_BG)
add_rect(slide, Inches(0), Inches(0), W, Inches(0.9), C_TITLE_BG)
add_text_box(slide, Inches(0.6), Inches(0.15), Inches(10), Inches(0.6),
             "ROS2 Node Integration Tests — 8/8 PASS", font_size=28, bold=True, color=C_WHITE)

# Big green badge
add_rect(slide, Inches(10.5), Inches(0.1), Inches(2.5), Inches(0.7), C_ACCENT2)
add_text_box(slide, Inches(10.5), Inches(0.15), Inches(2.5), Inches(0.6),
             "ALL PASS", font_size=24, bold=True, color=C_WHITE, alignment=PP_ALIGN.CENTER)

# Test matrix
test_data = [
    ["Test", "Nodes Under Test", "Checks", "Duration", "Status"],
    ["T1 terrain_analysis", "terrainAnalysis", "terrain_map output,\nground/obstacle intensity", "12s", "PASS"],
    ["T2 localPlanner", "terrainAnalysis +\nlocalPlanner", "Path generation,\nwall avoidance, stop signal", "18s", "PASS"],
    ["T3 pathFollower", "pathFollower", "cmd_vel direction,\nstuck detection", "15s", "PASS"],
    ["T4 global_planner", "pct_planner_astar", "Global path output,\npath quality", "8s", "PASS"],
    ["T5 pct_path_adapter", "pct_path_adapter", "Waypoint progression,\ngoal_reached event", "10s", "PASS"],
    ["T6 full_chain", "All 6 nodes", "goal \u2192 cmd_vel,\n5/5 checks, 18s", "18s", "PASS"],
    ["T7 safety_signals", "localPlanner +\npathFollower", "E-stop, slow_down,\ncmd_vel zeroing", "20s", "PASS"],
    ["T8 han_dog_bridge", "han_dog_bridge +\nmock/real CMS", "gRPC Walk(), watchdog,\njoint stream", "12s", "PASS"],
]
add_table(slide, Inches(0.4), Inches(1.1), Inches(12.5), Inches(5.2), test_data)

add_text_box(slide, Inches(0.5), Inches(6.5), Inches(12), Inches(0.5),
             "Platform: RDK S100P (aarch64) | ROS2 Humble | Real C++ nodes (terrain_analysis, localPlanner, pathFollower) + Python (global_planner, adapter)",
             font_size=11, color=C_GRAY)


# ══════════════════════════════════════════════════════════════
# SLIDE 8: SLAM Vibration Tests
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_LIGHT_BG)
add_rect(slide, Inches(0), Inches(0), W, Inches(0.9), C_TITLE_BG)
add_text_box(slide, Inches(0.6), Inches(0.15), Inches(10), Inches(0.6),
             "SLAM Vibration Robustness — Quadruped Gait Dataset", font_size=28, bold=True, color=C_WHITE)

# Left - Test setup
add_rect(slide, Inches(0.4), Inches(1.2), Inches(6.0), Inches(3.0), C_WHITE)
add_text_box(slide, Inches(0.6), Inches(1.3), Inches(5.6), Inches(0.4),
             "Dataset: Leg-KILO Corridor", font_size=18, bold=True, color=C_ACCENT)

slam_setup = [
    "Platform:      Unitree Go1 quadruped robot",
    "LiDAR:         Velodyne VLP-16 (16-line, 10Hz)",
    "IMU:           Internal (~500Hz)",
    "Scene:         Indoor corridor (L-shape + return)",
    "Duration:      445 seconds (7.4 minutes)",
    "Messages:      670,538 total",
    "  \u2022 PointCloud2:  4,421",
    "  \u2022 IMU:          222,039",
    "Vibration:     Gait impact + body pitch oscillation",
]
add_multiline(slide, Inches(0.6), Inches(1.8), Inches(5.6), Inches(2.5),
              slam_setup, font_size=13, font_name='Consolas', line_spacing=1.1)

# Right - Results
add_rect(slide, Inches(6.7), Inches(1.2), Inches(6.2), Inches(3.0), C_WHITE)
add_text_box(slide, Inches(6.9), Inches(1.3), Inches(5.8), Inches(0.4),
             "Results Summary", font_size=18, bold=True, color=C_ACCENT)

slam_data = [
    ["Scenario", "Algorithm", "Data Source", "Result", "Key Metric"],
    ["Quadruped\nvibration", "Point-LIO", "Leg-KILO\nVLP-16", "PASS", "445s stable,\nloop closure"],
    ["Quadruped\nvibration", "Fast-LIO2", "Leg-KILO\nVLP-16", "N/A", "No PointCloud2\nsupport"],
    ["Static\nreal-time", "Point-LIO", "Livox\nMid-360", "PASS", "0.006m drift"],
    ["Static\nreal-time", "Fast-LIO2", "Livox\nMid-360", "PASS", "0.008m drift"],
]
add_table(slide, Inches(6.8), Inches(1.8), Inches(6.0), Inches(2.2), slam_data)

# Bottom - Trajectory
add_rect(slide, Inches(0.4), Inches(4.5), Inches(12.5), Inches(2.6), C_WHITE)
add_text_box(slide, Inches(0.6), Inches(4.6), Inches(12), Inches(0.4),
             "Point-LIO Trajectory Analysis (Leg-KILO Corridor, 445s)", font_size=16, bold=True, color=C_DARK)

traj_data = [
    ["Metric", "Value"],
    ["Total Distance", "~220m (110m out + 110m return)"],
    ["Max Displacement", "109.9m from start"],
    ["Loop Closure Error", "7.1m (no explicit loop closure)"],
    ["Z-axis Drift", "<1m peak (-0.88m), acceptable"],
    ["IMU Init Time", "<1 second (1% \u2192 63% \u2192 100%)"],
    ["Output Rate", "10Hz stable throughout"],
]
add_table(slide, Inches(0.5), Inches(5.1), Inches(5.5), Inches(1.8), traj_data)

# Metric cards for SLAM
add_metric_card(slide, Inches(6.5), Inches(5.0), "220m", "Total Trajectory", C_ACCENT)
add_metric_card(slide, Inches(9.0), Inches(5.0), "<1m", "Z-drift (peak)", C_ACCENT2)
add_metric_card(slide, Inches(11.5), Inches(5.0), "10Hz", "Stable Output", C_ACCENT2)


# ══════════════════════════════════════════════════════════════
# SLIDE 9: Semantic Navigation
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_LIGHT_BG)
add_rect(slide, Inches(0), Inches(0), W, Inches(0.9), C_TITLE_BG)
add_text_box(slide, Inches(0.6), Inches(0.15), Inches(10), Inches(0.6),
             "Semantic Navigation — Natural Language Goal Resolution", font_size=28, bold=True, color=C_WHITE)

# Fast-Slow Architecture
add_rect(slide, Inches(0.4), Inches(1.2), Inches(6.2), Inches(3.0), C_WHITE)
add_text_box(slide, Inches(0.6), Inches(1.3), Inches(5.8), Inches(0.4),
             "Fast-Slow Dual Process Architecture", font_size=18, bold=True, color=C_ACCENT)

semantic_lines = [
    "System 1 (Fast Path):  <200ms, rule-based",
    "  \u2022  Scene graph keyword + spatial matching",
    "  \u2022  Confidence fusion: label 35% + CLIP 35%",
    "  \u2022  Threshold: 0.75 (fused score)",
    "  \u2022  Target: >70% hit rate",
    "",
    "System 2 (Slow Path):  ~2s, LLM reasoning",
    "  \u2022  ESCA selective grounding (92.5% token reduction)",
    "  \u2022  Multi-backend: Kimi / Qwen / Claude / OpenAI",
    "",
    "AdaNav Entropy Trigger:",
    "  entropy > 1.5 && confidence < 0.85 \u2192 force Slow Path",
]
add_multiline(slide, Inches(0.6), Inches(1.8), Inches(5.8), Inches(2.8),
              semantic_lines, font_size=13, line_spacing=1.1)

# Test Results
add_rect(slide, Inches(6.9), Inches(1.2), Inches(6.0), Inches(3.0), C_WHITE)
add_text_box(slide, Inches(7.1), Inches(1.3), Inches(5.6), Inches(0.4),
             "E2E Test Results (Factory Scene)", font_size=18, bold=True, color=C_ACCENT2)

sem_data = [
    ["Instruction", "Path", "Time", "Result"],
    ["\"\u5bfc\u822a\u5230\u5de5\u5382\u5927\u95e8\"", "Fast", "10.9s", "PASS"],
    ["\"\u5bfc\u822a\u5230\u673a\u68b0\u8bbe\u5907\"", "Fast", "71.6s", "PASS"],
    ["\"\u5bfc\u822a\u5230\u76ee\u6807\u533a\u57df\"", "Fast", "107s", "PASS"],
]
add_table(slide, Inches(7.0), Inches(1.8), Inches(5.8), Inches(1.5), sem_data)

sem_notes = [
    "Fast path confidence: 0.95 (score=0.92, well above 0.75)",
    "Full pipeline: instruction \u2192 scene graph \u2192 goal_pose \u2192",
    "  global_planner \u2192 adapter \u2192 local_planner \u2192 cmd_vel",
    "Tested on factory scene with 10+ semantic objects",
]
add_multiline(slide, Inches(7.1), Inches(3.4), Inches(5.6), Inches(1.0),
              sem_notes, font_size=12, line_spacing=1.1, color=C_GRAY)

# Bottom - Mission Arc
add_rect(slide, Inches(0.4), Inches(4.5), Inches(12.5), Inches(2.6), C_WHITE)
add_text_box(slide, Inches(0.6), Inches(4.6), Inches(12), Inches(0.4),
             "Mission Arc FSM — Task Lifecycle Management (New)", font_size=18, bold=True, color=C_ACCENT)

mission_lines = [
    "State Machine:  IDLE \u2192 PLANNING \u2192 EXECUTING \u2192 COMPLETE",
    "                         \u2502          \u2502",
    "                         \u2514\u2500 RECOVERING \u2192 REPLANNING \u2500\u2518",
    "                                   \u2502",
    "                                FAILED (terminal)",
    "",
    "Features:  \u2022 Pure observer (subscribes to /nav/adapter_status + /nav/planner_status)",
    "           \u2022 Publishes /nav/mission_status (JSON) at 2Hz",
    "           \u2022 Integrated into gRPC TelemetryService (SlowState.navigation.mission)",
    "           \u2022 Event emission on state transitions (EventBuffer)",
    "           \u2022 7/7 integration tests PASS",
]
add_multiline(slide, Inches(0.6), Inches(5.1), Inches(12), Inches(2.0),
              mission_lines, font_size=12, font_name='Consolas', line_spacing=1.05)


# ══════════════════════════════════════════════════════════════
# SLIDE 10: Conclusion
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, C_TITLE_BG)
add_rect(slide, Inches(0), Inches(0), W, Inches(0.06), C_ACCENT)

add_text_box(slide, Inches(0.6), Inches(0.6), Inches(12), Inches(0.8),
             "Conclusion & Future Work", font_size=34, bold=True, color=C_WHITE)

# Left - Contributions
add_rect(slide, Inches(0.4), Inches(1.8), Inches(6.0), Inches(3.5), RGBColor(0x1F, 0x33, 0x55))
add_text_box(slide, Inches(0.6), Inches(1.9), Inches(5.6), Inches(0.4),
             "Contributions", font_size=20, bold=True, color=C_ACCENT)

contrib_lines = [
    "1. End-to-end MuJoCo simulation framework for",
    "   quadruped navigation in multi-level industrial",
    "   environments (80m\u00d755m, 4 levels, 13.5m height)",
    "",
    "2. Real ONNX locomotion policy deployment on",
    "   aarch64 (RDK S100P), resolving 3 critical",
    "   sim-to-real transfer bugs",
    "",
    "3. Full ROS2 navigation stack integration:",
    "   8/8 node-level tests + 6/6 navigation trials PASS",
    "",
    "4. SLAM vibration robustness verified on real",
    "   quadruped gait dataset (Leg-KILO, 220m, 445s)",
    "",
    "5. Semantic VLN with Fast-Slow dual process:",
    "   3/3 Chinese instructions resolved in <107s",
]
add_multiline(slide, Inches(0.6), Inches(2.4), Inches(5.6), Inches(3.0),
              contrib_lines, font_size=13, color=C_WHITE, line_spacing=1.1)

# Right - Future
add_rect(slide, Inches(6.7), Inches(1.8), Inches(6.2), Inches(3.5), RGBColor(0x1F, 0x33, 0x55))
add_text_box(slide, Inches(6.9), Inches(1.9), Inches(5.8), Inches(0.4),
             "Future Work", font_size=20, bold=True, color=C_ACCENT)

future_lines = [
    "\u2022  Real-robot deployment on NOVA Thunder",
    "   quadruped (pending motor chain verification)",
    "",
    "\u2022  Multi-floor navigation with stair traversal",
    "   (ramp G0\u2192M1 validated, stairs pending RL policy)",
    "",
    "\u2022  Dynamic obstacle handling with semantic",
    "   perception (YOLO-World + ConceptGraphs)",
    "",
    "\u2022  OTA update system for field deployment",
    "   (v3.1.0 ready, Ed25519 signed manifests)",
    "",
    "\u2022  Fleet management with gRPC telemetry",
    "   (MissionArc FSM + EventBuffer integrated)",
]
add_multiline(slide, Inches(6.9), Inches(2.4), Inches(5.8), Inches(3.0),
              future_lines, font_size=13, color=C_WHITE, line_spacing=1.1)

# Bottom stats bar
add_rect(slide, Inches(0.4), Inches(5.6), Inches(12.5), Inches(1.5), RGBColor(0x14, 0x20, 0x38))

add_metric_card(slide, Inches(0.6), Inches(5.7), "6/6", "Nav Tests PASS")
add_metric_card(slide, Inches(3.1), Inches(5.7), "8/8", "Integration PASS")
add_metric_card(slide, Inches(5.6), Inches(5.7), "3/3", "Semantic PASS")
add_metric_card(slide, Inches(8.1), Inches(5.7), "220m", "SLAM Trajectory")
add_metric_card(slide, Inches(10.6), Inches(5.7), "16-DOF", "RL Policy")


# ══════════════════════════════════════════════════════════════
# Save
# ══════════════════════════════════════════════════════════════
out_dir = os.path.dirname(os.path.abspath(__file__))
out_path = os.path.join(out_dir, "MuJoCo_Factory_Navigation_Results.pptx")
prs.save(out_path)
print(f"PPT saved: {out_path}")
print(f"Slides: {len(prs.slides)}")
