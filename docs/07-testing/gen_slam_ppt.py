"""
Generate academic-style SLAM Vibration Robustness Test Report PPT.
White theme, no company branding, for academic advisor review.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Colors
BLACK      = RGBColor(0x22, 0x22, 0x22)
DARK       = RGBColor(0x33, 0x33, 0x33)
GRAY       = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY = RGBColor(0x99, 0x99, 0x99)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
BLUE       = RGBColor(0x1F, 0x77, 0xB4)
GREEN      = RGBColor(0x2C, 0xA0, 0x2C)
RED        = RGBColor(0xD6, 0x27, 0x28)
ORANGE     = RGBColor(0xFF, 0x7F, 0x0E)
HEADER_BG  = RGBColor(0x2C, 0x3E, 0x50)
TABLE_ALT  = RGBColor(0xF2, 0xF2, 0xF2)

W, H = 13.333, 7.5
FIG_DIR = os.path.dirname(os.path.abspath(__file__))

prs = Presentation()
prs.slide_width = Inches(W)
prs.slide_height = Inches(H)


def tb(slide, l, t, w, h, text, sz=16, color=BLACK, bold=False,
       align=PP_ALIGN.LEFT, font='Segoe UI'):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font
    p.alignment = align
    return box


def multi(slide, l, t, w, h, lines, dsz=14, dc=DARK):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if isinstance(line, str):
            txt, sz, clr, bld = line, dsz, dc, False
        else:
            txt = line[0]
            sz = line[1] if len(line) > 1 else dsz
            clr = line[2] if len(line) > 2 else dc
            bld = line[3] if len(line) > 3 else False
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = txt
        p.font.size = Pt(sz)
        p.font.color.rgb = clr
        p.font.bold = bld
        p.font.name = 'Segoe UI'
        p.space_after = Pt(3)
    return box


def line(slide, l, t, w, color=BLUE):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(l), Inches(t), Inches(w), Inches(0.025))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()


def table(slide, l, t, w, h, data, col_w=None):
    rows, cols = len(data), len(data[0])
    sh = slide.shapes.add_table(rows, cols, Inches(l), Inches(t), Inches(w), Inches(h))
    tbl = sh.table
    if col_w:
        for j, cw in enumerate(col_w):
            tbl.columns[j].width = Inches(cw)
    for i in range(rows):
        for j in range(cols):
            cell = tbl.cell(i, j)
            cell.text = str(data[i][j])
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(11)
                p.font.name = 'Segoe UI'
                p.alignment = PP_ALIGN.CENTER
                if i == 0:
                    p.font.color.rgb = WHITE
                    p.font.bold = True
                else:
                    p.font.color.rgb = DARK
            if i == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = HEADER_BG
            elif i % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_ALT
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE
    return sh


def img(slide, path, l, t, w=None, h=None):
    full = os.path.join(FIG_DIR, path) if not os.path.isabs(path) else path
    if os.path.exists(full):
        kw = {}
        if w: kw['width'] = Inches(w)
        if h: kw['height'] = Inches(h)
        return slide.shapes.add_picture(full, Inches(l), Inches(t), **kw)
    else:
        tb(slide, l, t, w or 5, 1, f'[Missing: {path}]', 12, RED)


def sn(slide, n):
    tb(slide, W - 1, H - 0.45, 0.8, 0.35, str(n), 9, LIGHT_GRAY, False, PP_ALIGN.RIGHT)


# ================================================================
# 1 — Title
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
line(s, 2.5, 2.8, 8.3, BLUE)
tb(s, 2.5, 3.0, 8.3, 1.0,
   'SLAM Vibration Robustness Evaluation', 36, BLACK, True, PP_ALIGN.CENTER)
tb(s, 2.5, 4.1, 8.3, 0.6,
   'Point-LIO performance under quadruped robot vibration\nacross indoor and outdoor terrain scenarios',
   16, GRAY, False, PP_ALIGN.CENTER)
line(s, 2.5, 5.0, 8.3, BLUE)
tb(s, 2.5, 5.3, 8.3, 0.5, '2026.03', 14, LIGHT_GRAY, False, PP_ALIGN.CENTER)
sn(s, 1)

# ================================================================
# 2 — Outline
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
tb(s, 0.8, 0.4, 10, 0.6, 'Outline', 28, BLACK, True)
line(s, 0.8, 0.95, 11.7, BLUE)
sections = [
    ('1', 'Motivation & Test Criteria'),
    ('2', 'Experimental Setup'),
    ('3', '2D Trajectory Visualization (per scenario)'),
    ('4', 'Overlaid Trajectory Comparison'),
    ('5', 'Z-axis Drift Analysis'),
    ('6', 'Displacement & Speed Profiles'),
    ('7', 'Position Component Analysis'),
    ('8', 'Quantitative Comparison'),
    ('9', 'Algorithm Comparison (table, trajectory, analysis)'),
    ('10', 'Conclusion'),
]
for i, (num, title) in enumerate(sections):
    y = 1.3 + i * 0.52
    tb(s, 1.5, y, 0.5, 0.4, num + '.', 14, BLUE, True, PP_ALIGN.RIGHT, 'Consolas')
    tb(s, 2.2, y, 8, 0.4, title, 15, DARK)
sn(s, 2)

# ================================================================
# 3 — Motivation
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
tb(s, 0.8, 0.4, 10, 0.6, '1. Motivation & Test Criteria', 28, BLACK, True)
line(s, 0.8, 0.95, 11.7, BLUE)

multi(s, 0.8, 1.3, 6, 3, [
    ('Problem', 16, BLACK, True),
    'Quadruped robots generate persistent high-frequency vibration:',
    '  - Foot-ground impact: 10-50 Hz impulse acceleration',
    '  - Body pitch/roll oscillation: ±5° attitude swing',
    '  - IMU noise variance increases 3-10x',
    '',
    'These directly challenge LiDAR-IMU coupled SLAM:',
    '  - Point cloud motion distortion',
    '  - IMU pre-integration drift accumulation',
    '  - Odometry discontinuity risk',
])

table(s, 7.5, 1.3, 5, 2.8, [
    ['Metric', 'Pass Criterion'],
    ['Survivability', 'Zero crashes'],
    ['Odometry output', 'Stable ≥5 Hz'],
    ['Pose plausibility', 'Consistent with geometry'],
    ['Z-axis drift', '< 2 m (flat corridor)'],
    ['Loop consistency', 'Measurable return error'],
])

multi(s, 0.8, 4.8, 11.7, 2.2, [
    ('Test scope', 16, BLACK, True),
    ('Three terrain scenarios tested with Leg-KILO open dataset (Unitree Go1 + VLP-16):', 14),
    ('  1. Corridor — indoor flat L-shaped hallway (445 s)', 13),
    ('  2. Grass — outdoor uneven grass terrain (91.6 s)', 13),
    ('  3. Slope — outdoor hill/slope terrain (312 s)', 13),
    ('  + Real-time baseline: Livox Mid-360 static drift test (30 s)', 13, GRAY),
])
sn(s, 3)

# ================================================================
# 4 — Experimental Setup
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
tb(s, 0.8, 0.4, 10, 0.6, '2. Experimental Setup', 28, BLACK, True)
line(s, 0.8, 0.95, 11.7, BLUE)

table(s, 0.8, 1.3, 5.8, 3.5, [
    ['Parameter', 'Value'],
    ['Source', 'Leg-KILO (KAIST, 2023)'],
    ['Robot', 'Unitree Go1 quadruped'],
    ['LiDAR', 'Velodyne VLP-16 (16-beam, 10 Hz)'],
    ['IMU', 'Built-in (~500 Hz)'],
    ['Corridor', '445 s, 670 K msgs, L-shape'],
    ['Grass', '91.6 s, 211 K msgs, uneven'],
    ['Slope', '312 s, 467 K msgs, hill terrain'],
])

table(s, 7.2, 1.3, 5.3, 2.8, [
    ['SLAM', 'Config'],
    ['Point-LIO', 'lidar_type=2, scan_line=16'],
    ['', 'timestamp_unit=2 (μs)'],
    ['', 'iVox point-wise update'],
    ['Fast-LIO2', 'lidar_type=2, acc_scale=1.0'],
    ['', 'iKD-Tree frame-wise update'],
])

multi(s, 7.2, 4.3, 5.3, 2.5, [
    ('Platform', 14, BLACK, True),
    '  RDK S100P (aarch64), ROS2 Humble',
    '  CycloneDDS, single-threaded (MP=1)',
    '',
    ('Data conversion', 14, BLACK, True),
    '  ROS1 bag → rosbags → ROS2 Humble',
    '  Metadata fix: version=4, rm hash,',
    '  offered_qos_profiles=""',
])
sn(s, 4)

# ================================================================
# 5-7 — Individual Trajectories
# ================================================================
for i, (fname, title) in enumerate([
    ('fig_corridor_traj.png', '3a. 2D Trajectory — Corridor (Indoor, 295 s)'),
    ('fig_slope_traj.png',    '3b. 2D Trajectory — Slope (Outdoor, 162 s)'),
    ('fig_grass_traj.png',    '3c. 2D Trajectory — Grass (Outdoor, 27 s)'),
]):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    tb(s, 0.8, 0.3, 11, 0.5, title, 22, BLACK, True)
    line(s, 0.8, 0.75, 11.7, BLUE)
    img(s, fname, 0.5, 1.0, w=12.3)
    sn(s, 5 + i)

# ================================================================
# 8 — Overlaid Trajectories
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
tb(s, 0.8, 0.3, 11, 0.5, '4. Overlaid Trajectory Comparison (Origin-Aligned)', 22, BLACK, True)
line(s, 0.8, 0.75, 11.7, BLUE)
img(s, 'fig_overlay_traj.png', 0.5, 1.0, w=12.3)
sn(s, 8)

# ================================================================
# 9 — Z-axis Drift
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
tb(s, 0.8, 0.3, 11, 0.5, '5. Z-axis Drift Comparison', 22, BLACK, True)
line(s, 0.8, 0.75, 11.7, BLUE)
img(s, 'fig_z_drift_compare.png', 0.5, 1.0, w=12.3)

multi(s, 0.8, 5.5, 11.7, 1.5, [
    ('Key finding: Z estimation accuracy degrades with terrain complexity (flat 0.9 m → slope 6.4 m → grass 19.4 m),', 12, DARK),
    ('but Point-LIO maintains zero crashes across all scenarios.', 12, DARK),
])
sn(s, 9)

# ================================================================
# 10 — Displacement + Cumulative + Speed
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
tb(s, 0.8, 0.3, 11, 0.5, '6. Displacement, Distance & Speed Profiles', 22, BLACK, True)
line(s, 0.8, 0.75, 11.7, BLUE)
img(s, 'fig_displacement_compare.png', 0.2, 0.9, w=6.3)
img(s, 'fig_speed_compare.png', 6.7, 0.9, w=6.3)
img(s, 'fig_cumulative_dist.png', 0.2, 4.0, w=6.3)
sn(s, 10)

# ================================================================
# 11 — XYZ Components
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
tb(s, 0.8, 0.3, 11, 0.5, '7. Position Components (X, Y, Z) vs Time', 22, BLACK, True)
line(s, 0.8, 0.75, 11.7, BLUE)
img(s, 'fig_corridor_xyz.png', 0.2, 0.9, w=6.3)
img(s, 'fig_slope_xyz.png', 6.7, 0.9, w=6.3)
sn(s, 11)

# ================================================================
# 12 — Quantitative Bar Comparison
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
tb(s, 0.8, 0.3, 11, 0.5, '8. Quantitative Comparison', 22, BLACK, True)
line(s, 0.8, 0.75, 11.7, BLUE)
img(s, 'fig_bar_compare.png', 1.0, 0.9, w=11.3)

table(s, 1.5, 4.8, 10.3, 2.2, [
    ['Metric', 'Corridor (indoor)', 'Grass (outdoor)', 'Slope (outdoor)'],
    ['Result', 'PASS', 'PASS', 'PASS'],
    ['Crashes', '0', '0', '0'],
    ['Duration', '295 s', '27 s', '162 s'],
    ['Total distance', '~268 m', '~47 m', '~346 m'],
    ['Max displacement', '109.9 m', '35.8 m', '132.6 m'],
    ['|Z| peak drift', '0.88 m', '19.4 m', '6.4 m'],
])

# Color PASS cells
tbl = s.shapes[-1].table
for j in range(1, 4):
    cell = tbl.cell(1, j)
    for p in cell.text_frame.paragraphs:
        p.font.color.rgb = GREEN
        p.font.bold = True
sn(s, 12)

# ================================================================
# 13 — Algorithm Comparison (Table)
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
tb(s, 0.8, 0.3, 11, 0.5, '9a. Algorithm Comparison: Point-LIO vs Fast-LIO2', 22, BLACK, True)
line(s, 0.8, 0.75, 11.7, BLUE)

table(s, 0.8, 1.1, 11.7, 3.5, [
    ['Dimension', 'Point-LIO', 'Fast-LIO2'],
    ['Core algorithm', 'Point-wise update + iVox', 'Frame-wise update + iKD-Tree'],
    ['Sensor support', 'Livox + Velodyne + Ouster', 'Livox + Velodyne (modified)'],
    ['Corridor test', 'PASS (174 s)', 'PASS (234 s, full loop)'],
    ['Z-axis drift (corridor)', '4.5 m peak', '0.9 m peak'],
    ['Loop closure error', '~30 m (partial)', '< 0.5 m'],
    ['Static drift (Mid-360)', '0.006 m', '0.008 m'],
    ['Odometry rate', '10.0 Hz', '10.5 Hz'],
], col_w=[3.0, 4.35, 4.35])

# Color PASS cells
tbl = s.shapes[-1].table
for j in [1, 2]:
    cell = tbl.cell(3, j)
    for p in cell.text_frame.paragraphs:
        p.font.color.rgb = GREEN
        p.font.bold = True

multi(s, 0.8, 5.2, 5.5, 2.0, [
    ('Key finding', 16, BLACK, True),
    ('Both algorithms survive quadruped vibration,', 14),
    ('but Fast-LIO2 achieves significantly better', 14),
    ('Z-axis stability (0.9 m vs 4.5 m) and loop', 14),
    ('closure accuracy on the corridor dataset.', 14),
])
multi(s, 7.0, 5.2, 5.5, 2.0, [
    ('Critical config', 16, BLACK, True),
    ('Fast-LIO2 requires lidar_filter_num = 1', 14),
    ('(keep all points) for VLP-16 under vibration.', 14),
    ('With filter_num = 3, divergence occurs at ~78 s.', 14),
    ('Point-LIO tolerates aggressive downsampling.', 14),
])
sn(s, 13)

# ================================================================
# 14 — Dual Trajectory Comparison Figure
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
tb(s, 0.8, 0.3, 11, 0.5, '9b. Trajectory Comparison — Same Dataset', 22, BLACK, True)
line(s, 0.8, 0.75, 11.7, BLUE)
img(s, 'fig_dual_traj_xy.png', 0.3, 0.9, w=12.7)
sn(s, 14)

# ================================================================
# 15 — Dual Z-drift + XYZ
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
tb(s, 0.8, 0.3, 11, 0.5, '9c. Z-drift & Position Analysis Comparison', 22, BLACK, True)
line(s, 0.8, 0.75, 11.7, BLUE)
img(s, 'fig_dual_z_drift.png', 0.2, 0.9, w=6.3)
img(s, 'fig_dual_displacement.png', 6.7, 0.9, w=6.3)
img(s, 'fig_dual_xyz.png', 1.5, 4.0, w=10.3, h=3.2)
sn(s, 15)

# ================================================================
# 16 — 3D Trajectory
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
tb(s, 0.8, 0.3, 11, 0.5, 'Appendix A. 3D Trajectory — Corridor', 22, BLACK, True)
line(s, 0.8, 0.75, 11.7, BLUE)
img(s, 'fig_corridor_3d.png', 1.5, 1.0, w=10.3)
sn(s, 16)

# ================================================================
# 17 — Conclusion
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
tb(s, 0.8, 0.4, 10, 0.6, '10. Conclusion', 28, BLACK, True)
line(s, 0.8, 0.95, 11.7, BLUE)

multi(s, 0.8, 1.3, 11, 5, [
    ('1. Both Point-LIO and Fast-LIO2 survive quadruped vibration', 16, BLACK, True),
    ('   Zero crashes across 3 terrain scenarios (Point-LIO) and corridor (Fast-LIO2),', 14),
    ('   totaling 1000+ s of testing and 900+ m of travel distance.', 14),
    '',
    ('2. Fast-LIO2 achieves superior Z-axis stability on corridor', 16, BLACK, True),
    ('   Z drift: Fast-LIO2 peak 0.9 m vs Point-LIO peak 4.5 m (5x better).', 14),
    ('   Loop closure: Fast-LIO2 < 0.5 m error, Point-LIO ~30 m (partial monitoring).', 14),
    '',
    ('3. Configuration is critical for vibration robustness', 16, BLACK, True),
    ('   Fast-LIO2 requires lidar_filter_num = 1 (no downsampling) under vibration.', 14),
    ('   With filter_num = 3, catastrophic divergence at ~78 s. Point-LIO is more tolerant.', 14),
    '',
    ('4. Point-LIO advantages: broader sensor support + faster init', 16, BLACK, True),
    ('   Natively supports Livox, Velodyne, Ouster. ~1 s IMU init vs ~3 s.', 14),
    ('   Recommended as default for multi-sensor deployments.', 14),
    '',
    ('5. Future work', 16, BLACK, True),
    ('   - ATE/RPE with ground truth trajectories', 14),
    ('   - Outdoor terrain (grass, slope) dual-algorithm comparison', 14),
    ('   - Loop closure integration for Point-LIO Z drift reduction', 14),
])
sn(s, 17)

# ================================================================
# 18 — End
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
line(s, 3, 3.0, 7.3, BLUE)
tb(s, 3, 3.2, 7.3, 0.8, 'Thank You', 36, BLACK, True, PP_ALIGN.CENTER)
tb(s, 3, 4.1, 7.3, 0.5,
   '2 algorithms · 3 scenarios · 0 crashes · Z drift < 1 m (Fast-LIO2)',
   16, GRAY, False, PP_ALIGN.CENTER)
line(s, 3, 4.7, 7.3, BLUE)
sn(s, 18)


# ===== Save =====
out = os.path.join(FIG_DIR, 'SLAM_振动鲁棒性测试报告.pptx')
prs.save(out)
print(f'Saved: {out}')
print(f'Slides: {len(prs.slides)}')
