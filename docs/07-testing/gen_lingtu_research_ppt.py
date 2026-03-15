"""
Comprehensive Research PPT — LingTu Autonomous Navigation System
=================================================================
Covers: system architecture, MuJoCo factory simulation, ONNX policy,
SLAM vibration robustness (Point-LIO / Fast-LIO2), integration tests T1-T8,
semantic navigation Fast-Slow dual process, quantitative results.

Output: D:\inovxio\docs\ppt\LingTu_Navigation_Research_2026.pptx
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

try:
    from PIL import Image
except ImportError:
    Image = None

# ──────────────────── paths ────────────────────
LINGTU = Path(r"D:\inovxio\brain\lingtu")
TESTING = LINGTU / "docs" / "07-testing"
TOOLS = LINGTU / "tools"
OUT = Path(r"D:\inovxio\docs\ppt\LingTu_Navigation_Research_2026.pptx")

# ─── SLAM figures ───
FIG_3SCENARIO    = TESTING / "fig_3scenario_comparison.png"
FIG_CORRIDOR     = TESTING / "fig_corridor_traj.png"
FIG_CORRIDOR_3D  = TESTING / "fig_corridor_3d.png"
FIG_CORRIDOR_XYZ = TESTING / "fig_corridor_xyz.png"
FIG_DUAL_TRAJ    = TESTING / "fig_dual_traj_xy.png"
FIG_DUAL_ZDRIFT  = TESTING / "fig_dual_z_drift.png"
FIG_DUAL_XYZ     = TESTING / "fig_dual_xyz.png"
FIG_DUAL_DISP    = TESTING / "fig_dual_displacement.png"
FIG_GRASS_TRAJ   = TESTING / "fig_grass_trajectory_2d.png"
FIG_GRASS_PANEL  = TESTING / "fig_grass_analysis_panels.png"
FIG_SLOPE_TRAJ   = TESTING / "fig_slope_trajectory_2d.png"
FIG_SLOPE_PANEL  = TESTING / "fig_slope_analysis_panels.png"
FIG_OVERLAY      = TESTING / "fig_overlay_traj.png"
FIG_ZDRIFT_CMP   = TESTING / "fig_z_drift_compare.png"
FIG_DISP_CMP     = TESTING / "fig_displacement_compare.png"
FIG_SPEED_CMP    = TESTING / "fig_speed_compare.png"
FIG_BAR_CMP      = TESTING / "fig_bar_compare.png"
FIG_COMPARE_DASH = TESTING / "fig_comparison_dashboard.png"
FIG_CUMULATIVE   = TESTING / "fig_cumulative_dist.png"
FIG_SLAM_ANALYSIS = TESTING / "slam_trajectory_analysis.png"
FIG_SLOPE_XYZ    = TESTING / "fig_slope_xyz.png"
FIG_ANALYSIS_PNL = TESTING / "fig_analysis_panels.png"

# ─── Navigation / MuJoCo figures ───
FIG_FACTORY_3D   = TOOLS / "factory_3d_overview.png"
FIG_E2E_FACTORY  = TOOLS / "e2e_factory_traj.png"
FIG_E2E_BUILDING = TOOLS / "e2e_building_traj.png"
FIG_E2E_PAPER    = TOOLS / "e2e_paper_figure.png"
FIG_E2E_FULL     = TOOLS / "e2e_pct_full_traj.png"
FIG_E2E_TERRAIN  = TOOLS / "e2e_terrain_compare.png"
FIG_E2E_PCD      = TOOLS / "e2e_pcd_terrain_traj.png"
FIG_SITL_TRAJ    = TOOLS / "factory_sitl_traj.png"
FIG_SITL         = TOOLS / "sitl_traj.png"
FIG_NOVA_STAND   = TOOLS / "nova_standing.png"
FIG_RVIZ_DATA    = TOOLS / "rviz_with_data.png"
FIG_RVIZ_VIRT    = TOOLS / "rviz_virt.png"
FIG_SIM_BUILD2   = TOOLS / "sim_building2_traj.png"
FIG_SIM_BUILD2V  = TOOLS / "sim_building2_viz_traj.png"
FIG_NAV_PCD      = TOOLS / "nav_on_pcd.png"
FIG_SPIRAL_OV    = TOOLS / "spiral_plan_overview.png"
FIG_SPIRAL_3D    = TOOLS / "spiral_plan_3d.png"
FIG_SPIRAL_FL    = TOOLS / "spiral_floors.png"
FIG_SPIRAL_HM    = TOOLS / "spiral_heightmap.png"
FIG_SPIRAL_ZPROF = TOOLS / "spiral_zprofile.png"
FIG_TRAJ_3D_MF   = TOOLS / "traj_3d_multifloor.png"
FIG_FINAL_F1     = TOOLS / "final_f1.png"
FIG_FINAL_F2     = TOOLS / "final_f2.png"
FIG_FINAL_F3     = TOOLS / "final_f3.png"
FIG_FINAL_F4     = TOOLS / "final_f4.png"
FIG_E2E_CMP_STR  = TOOLS / "e2e_compare_straight_vs_astar.png"
FIG_RVIZ_GOAL    = TOOLS / "rviz_frame_goal.png"
FIG_RVIZ_RF      = TOOLS / "rviz_frame_rf.png"
FIG_RVIZ_STAIR   = TOOLS / "rviz_frame_stair.png"
FIG_RVIZ_START   = TOOLS / "rviz_frame_start.png"

# ──────────────────── colors & fonts ────────────────────
SLIDE_W, SLIDE_H = 13.333, 7.5
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
BG_DARK    = RGBColor(0x0F, 0x17, 0x2A)
TITLE_CLR  = RGBColor(0x16, 0x25, 0x42)
TEXT_CLR   = RGBColor(0x28, 0x28, 0x28)
SUB_CLR    = RGBColor(0x5A, 0x5A, 0x5A)
ACCENT     = RGBColor(0x0B, 0x72, 0x85)
ACCENT2    = RGBColor(0x1E, 0x88, 0xE5)
GREEN      = RGBColor(0x2B, 0x8A, 0x3E)
RED        = RGBColor(0xC9, 0x2A, 0x2A)
ORANGE     = RGBColor(0xE6, 0x77, 0x00)
CARD_FILL  = RGBColor(0xF3, 0xF4, 0xF6)
CARD_LINE  = RGBColor(0xD2, 0xD6, 0xDC)
FORMULA_BG = RGBColor(0xEE, 0xF2, 0xF7)
FORMULA_LN = RGBColor(0xB8, 0xC4, 0xD4)
FONT = "Microsoft YaHei"
MONO = "Consolas"

# ──────────────────── helpers ────────────────────
def _fit(path, max_w, max_h):
    if Image is None:
        return max_w, max_h
    try:
        with Image.open(path) as img:
            iw, ih = img.size
    except Exception:
        return max_w, max_h
    scale = min(max_w * 96 / iw, max_h * 96 / ih)
    return iw * scale / 96, ih * scale / 96


def set_bg(slide, color=WHITE):
    f = slide.background.fill; f.solid(); f.fore_color.rgb = color


def add_shape(slide, x, y, w, h, fill, line=None, rounded=False):
    st = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if rounded else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    s = slide.shapes.add_shape(st, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line: s.line.color.rgb = line
    else: s.line.fill.background()
    return s


def add_text(slide, x, y, w, h, text, size=18, color=TEXT_CLR, bold=False,
             font=FONT, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text; p.alignment = align
    p.font.size = Pt(size); p.font.name = font
    p.font.color.rgb = color; p.font.bold = bold
    return box


def add_lines(slide, x, y, w, h, lines, font=FONT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame; tf.word_wrap = True; tf.clear()
    for idx, (text, size, color, bold) in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = text; p.font.size = Pt(size); p.font.name = font
        p.font.color.rgb = color; p.font.bold = bold; p.space_after = Pt(4)
    return box


def add_bullets(slide, x, y, w, h, items, size=15):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame; tf.word_wrap = True; tf.clear()
    for idx, item in enumerate(items):
        if isinstance(item, tuple) and len(item) == 2:
            text, color = item
        else:
            text = item[0] if isinstance(item, tuple) else item
            color = TEXT_CLR
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = text; p.font.size = Pt(size); p.font.name = FONT
        p.font.color.rgb = color; p.bullet = True; p.space_after = Pt(6)
    return box


def add_card(slide, x, y, w, h, title=None):
    add_shape(slide, x, y, w, h, CARD_FILL, CARD_LINE, rounded=True)
    if title:
        add_text(slide, x + 0.15, y + 0.12, w - 0.3, 0.26, title,
                 size=14, color=ACCENT, bold=True)


def add_formula_box(slide, x, y, w, h, lines):
    add_shape(slide, x, y, w, h, FORMULA_BG, FORMULA_LN, rounded=True)
    add_lines(slide, x + 0.15, y + 0.12, w - 0.3, h - 0.24,
              [(l, 12, TITLE_CLR, False) for l in lines], font=MONO)


def add_image(slide, path, x, y, w, h):
    p = Path(path)
    if not p.exists():
        add_shape(slide, x, y, w, h, CARD_FILL, CARD_LINE, rounded=True)
        add_text(slide, x, y + h / 2 - 0.12, w, 0.24,
                 f"[missing: {p.name}]", size=10, color=SUB_CLR, align=PP_ALIGN.CENTER)
        return None
    fw, fh = _fit(p, w, h)
    px, py_ = x + (w - fw) / 2, y + (h - fh) / 2
    return slide.shapes.add_picture(str(p), Inches(px), Inches(py_),
                                    width=Inches(fw), height=Inches(fh))


def add_header(slide, title, subtitle=""):
    add_text(slide, 0.55, 0.28, 11.8, 0.42, title, size=28, color=TITLE_CLR, bold=True)
    if subtitle:
        add_text(slide, 0.58, 0.74, 11.5, 0.18, subtitle, size=12, color=SUB_CLR)
    add_shape(slide, 0.55, 0.99, 1.35, 0.03, ACCENT)


def add_caption(slide, x, y, w, text, align=PP_ALIGN.LEFT):
    add_text(slide, x, y, w, 0.22, text, size=10, color=SUB_CLR, align=align)


def add_code(slide, x, y, w, h, title, lines):
    add_card(slide, x, y, w, h, title)
    add_lines(slide, x + 0.15, y + 0.42, w - 0.3, h - 0.5,
              [(l, 11, TEXT_CLR, False) for l in lines], font=MONO)


def ns(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s); return s


# ──────────────────── BUILD ────────────────────
def build():
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    # ═══════════════ S1: TITLE ═══════════════
    s = ns(prs); set_bg(s, BG_DARK)
    add_text(s, 0.8, 0.9, 11.5, 0.6,
             "LingTu: Autonomous Navigation for Quadruped Robots",
             size=34, color=WHITE, bold=True)
    add_text(s, 0.8, 1.55, 11.5, 0.4,
             "MuJoCo Simulation  |  SLAM Robustness  |  Semantic VLN  |  Full-Stack Integration",
             size=18, color=RGBColor(0x88, 0xCC, 0xEE))
    add_text(s, 0.8, 2.15, 11.5, 0.3,
             "Jetson/RDK S100P  |  ROS2 Humble  |  Fast-LIO2 / Point-LIO  |  PCT Planner  |  gRPC Gateway",
             size=13, color=RGBColor(0xAA, 0xAA, 0xBB))

    add_image(s, FIG_FACTORY_3D, 0.8, 3.0, 5.5, 3.5)
    add_image(s, FIG_3SCENARIO, 6.8, 3.0, 5.7, 3.5)
    add_text(s, 0.8, 6.7, 6.0, 0.2,
             "Shanghai Qiongpei Technology Co., Ltd.", size=11,
             color=RGBColor(0x88, 0x88, 0x99))
    add_text(s, 8.0, 6.7, 4.5, 0.2, "2026-03-11", size=11,
             color=RGBColor(0x88, 0x88, 0x99), align=PP_ALIGN.RIGHT)

    # ═══════════════ S2: SYSTEM ARCHITECTURE ═══════════════
    s = ns(prs)
    add_header(s, "System Architecture",
               "LiDAR \u2192 SLAM \u2192 Terrain \u2192 Planning \u2192 Dog Board | Semantic Layer (optional)")

    add_card(s, 0.55, 1.3, 12.2, 2.2, "Full-Stack Pipeline")
    add_formula_box(s, 0.75, 1.78, 11.75, 1.5, [
        "Livox Mid-360 LiDAR",
        "  \u2192 SLAM (Fast-LIO2 / Point-LIO + PGO + ICP Localizer)",
        "    \u2192 Terrain Analysis (ground estimation, traversability)",
        "      \u2192 Global Planner (PCT A* on tomogram) \u2192 Path Adapter (waypoint tracking)",
        "        \u2192 Local Planner (obstacle avoidance) \u2192 Path Follower (cmd_vel generation)",
        "          \u2192 han_dog_bridge (gRPC Walk) \u2192 Brainstem (RL policy \u2192 motor commands)",
    ])

    add_card(s, 0.55, 3.75, 3.85, 3.3, "Perception Stack")
    add_bullets(s, 0.75, 4.22, 3.4, 2.5, [
        ("Fast-LIO2: iKD-Tree + iEKF", ACCENT),
        "  Livox CustomMsg, 10Hz odom output",
        ("Point-LIO: point-level iVox", ACCENT),
        "  Velodyne + Livox, vibration-robust",
        "PGO loop closure (optional)",
        "ICP relocalization (map-based)",
    ], size=12)

    add_card(s, 4.6, 3.75, 4.05, 3.3, "Planning Stack")
    add_bullets(s, 4.8, 4.22, 3.6, 2.5, [
        ("PCT Planner: tomography-based A*", ACCENT),
        "  4D (x, y, z, heading), ele_planner.so C++",
        "pct_path_adapter: waypoint tracking",
        "  max_index_jump=3, stuck detection",
        "Local Planner: terrain-aware",
        "  slopeWeight, obstacle avoidance",
        "Path Follower: PID + stuck detection",
        "  10s warn, 15s confirmed STUCK",
    ], size=12)

    add_card(s, 8.85, 3.75, 3.9, 3.3, "Semantic Layer")
    add_bullets(s, 9.05, 4.22, 3.45, 2.5, [
        ("Fast-Slow dual process", GREEN),
        "  Fast: scene graph match (<200ms)",
        "  Slow: LLM reasoning (~2s)",
        "YOLO-World + CLIP + ConceptGraphs",
        "AdaNav entropy trigger",
        "LERa failure recovery (3-step)",
        ("Chinese VLN instruction support", GREEN),
    ], size=12)

    # ═══════════════ S3: MUJOCO FACTORY ═══════════════
    s = ns(prs)
    add_header(s, "MuJoCo Factory Simulation",
               "80m \u00d7 55m multi-level industrial building | 5 floors (B1, G0, M1, M2, RF)")

    add_image(s, FIG_FACTORY_3D, 0.55, 1.3, 4.5, 3.0)
    add_caption(s, 0.55, 4.35, 4.5, "3D factory overview: 5 levels, 13.5m height",
                PP_ALIGN.CENTER)

    add_image(s, FIG_E2E_FACTORY, 5.25, 1.3, 3.8, 3.0)
    add_caption(s, 5.25, 4.35, 3.8, "Factory flat-terrain E2E navigation trajectory",
                PP_ALIGN.CENTER)

    add_image(s, FIG_SITL_TRAJ, 9.25, 1.3, 3.45, 3.0)
    add_caption(s, 9.25, 4.35, 3.45, "SITL trajectory on factory map",
                PP_ALIGN.CENTER)

    add_card(s, 0.55, 4.75, 5.7, 2.3, "Scene Parameters")
    add_formula_box(s, 0.75, 5.22, 5.25, 1.55, [
        "Factory:  80m x 55m x 13.5m",
        "Levels:   B1(-3.5) G0(0) M1(+3.5) M2(+7.0) RF(+10.5)",
        "Tomogram: (5, 2, 154, 118), res=0.2m/cell",
        "File:     factory_nova.pickle (924 KB)",
        "Center:   (15.06, 8.00) m",
    ])

    add_card(s, 6.45, 4.75, 6.3, 2.3, "ONNX Policy Integration")
    add_formula_box(s, 6.65, 5.22, 5.85, 1.55, [
        "Input:  [1, 285] = 57-dim x 5 history frames",
        "  gyro(3) + gravity(3) + cmd(3) + pos(16) + vel(16) + act(16)",
        "Output: [1, 16] joint actions",
        "PD:     kp = [390, 570, 720], kv = 5, damping = 0",
        "Joint:  MJ_TO_DART index mapping (4+4+4+4 -> 3+3+3+3+4)",
    ])

    # ═══════════════ S4: LiDAR SIM + NAVIGATION RESULTS ═══════════════
    s = ns(prs)
    add_header(s, "MuJoCo Navigation Results",
               "6/6 navigation tests PASS | LiDAR: mj_multiRay 6400 rays per scan")

    add_image(s, FIG_E2E_BUILDING, 0.55, 1.3, 4.0, 2.6)
    add_caption(s, 0.55, 3.95, 4.0, "Building scenario: 25m x 18m, A* path planning",
                PP_ALIGN.CENTER)
    add_image(s, FIG_E2E_FULL, 4.75, 1.3, 4.0, 2.6)
    add_caption(s, 4.75, 3.95, 4.0, "Full PCT planning trajectory on tomogram",
                PP_ALIGN.CENTER)
    add_image(s, FIG_E2E_TERRAIN, 8.95, 1.3, 3.8, 2.6)
    add_caption(s, 8.95, 3.95, 3.8, "Terrain analysis comparison",
                PP_ALIGN.CENTER)

    add_card(s, 0.55, 4.35, 6.0, 2.7, "Test Results (6/6 PASS)")
    add_lines(s, 0.75, 4.82, 5.55, 2.0, [
        ("Test                      | Scene     | Result", 12, TITLE_CLR, True),
        ("Factory flat terrain      | 80x55m    | PASS", 11, GREEN, False),
        ("Building 25x18m           | building  | PASS", 11, GREEN, False),
        ("G0 level navigation       | factory   | PASS", 11, GREEN, False),
        ("Building 2-9 (A* path)    | building  | PASS", 11, GREEN, False),
        ("PCD terrain navigation    | custom    | PASS", 11, GREEN, False),
        ("Multi-floor spiral        | 5-level   | PASS", 11, GREEN, False),
    ], font=MONO)

    add_card(s, 6.75, 4.35, 5.95, 2.7, "LiDAR Simulation")
    add_bullets(s, 6.95, 4.82, 5.5, 2.0, [
        "MuJoCo mj_multiRay API (v3.5.0)",
        "6,400 rays per scan, 1.5ms on RDK S100P",
        "Golden angle spiral (non-repetitive pattern)",
        "Livox Mid-360 OmniPerception emulation",
        "  800K sampling points, 40-frame window",
        ("Output: /nav/map_cloud PointCloud2 (odom frame)", ACCENT),
    ], size=12)

    # ═══════════════ S5: MULTI-FLOOR PLANNING ═══════════════
    s = ns(prs)
    add_header(s, "Multi-Floor Navigation & Path Planning",
               "PCT tomography-based A* planner on 4D (x, y, z, heading) grid")

    add_image(s, FIG_SPIRAL_OV, 0.55, 1.3, 4.2, 2.7)
    add_caption(s, 0.55, 4.05, 4.2, "Spiral ramp plan overview: B1 \u2192 RF", PP_ALIGN.CENTER)
    add_image(s, FIG_SPIRAL_3D, 4.95, 1.3, 3.8, 2.7)
    add_caption(s, 4.95, 4.05, 3.8, "3D path through 5 floors", PP_ALIGN.CENTER)
    add_image(s, FIG_TRAJ_3D_MF, 9.0, 1.3, 3.7, 2.7)
    add_caption(s, 9.0, 4.05, 3.7, "Multi-floor 3D trajectory result", PP_ALIGN.CENTER)

    add_card(s, 0.55, 4.35, 4.5, 2.7, "PCT Planner Architecture")
    add_formula_box(s, 0.75, 4.82, 4.05, 2.0, [
        "# Tomogram generation:",
        "scene_pcd \u2192 voxelize(0.2m) \u2192 Z-slice",
        "\u2192 obstacle_map + traversability_map",
        "\u2192 pickle (5, 2, W, H)",
        "",
        "# A* search:",
        "4D state: (x, y, z_level, heading)",
        "Heuristic: Euclidean + height penalty",
        "Output: nav_msgs/Path (/nav/global_path)",
    ])

    add_card(s, 5.25, 4.35, 3.7, 2.7, "Path Adapter")
    add_bullets(s, 5.45, 4.82, 3.25, 2.0, [
        "pct_path_adapter: waypoint tracking",
        "  path \u2192 downsample \u2192 way_point",
        "  progress: /nav/adapter_status JSON",
        "max_index_jump = 3 (skip protection)",
        "max_first_waypoint_dist = 10m",
        ("Events: path_received, waypoint_reached,", ACCENT),
        ("  goal_reached, stuck_detected", ACCENT),
    ], size=12)

    add_card(s, 9.15, 4.35, 3.55, 2.7, "Staircase Figures")
    add_image(s, FIG_SPIRAL_FL, 9.35, 5.0, 1.5, 1.7)
    add_image(s, FIG_SPIRAL_ZPROF, 10.95, 5.0, 1.5, 1.7)

    # ═══════════════ S6: INTEGRATION TESTS T1-T8 ═══════════════
    s = ns(prs)
    add_header(s, "Integration Tests: T1-T8 Full Coverage",
               "8/8 PASS — Real ROS2 node communication on S100P (192.168.66.190)")

    add_card(s, 0.55, 1.3, 12.2, 3.0, "Node-Level Verification Matrix")
    add_lines(s, 0.75, 1.78, 11.7, 2.3, [
        ("Test | Nodes Under Test                   | Verifies                        | Result", 12, TITLE_CLR, True),
        ("T1   | terrainAnalysis                     | PointCloud2 \u2192 terrain_map intensity | 3/3 PASS", 11, GREEN, False),
        ("T2   | terrainAnalysis + localPlanner      | terrain + waypoint \u2192 /path + /stop  | 4/4 PASS", 11, GREEN, False),
        ("T3   | pathFollower                        | /path \u2192 /nav/cmd_vel direction       | 3/3 PASS", 11, GREEN, False),
        ("T4   | global_planner (A*)                 | goal_pose \u2192 global_path             | 5/5 PASS", 11, GREEN, False),
        ("T5   | pct_path_adapter                    | global_path + odom \u2192 waypoint        | 4/4 PASS", 11, GREEN, False),
        ("T6   | ALL 6 nodes (full chain)            | goal_pose \u2192 cmd_vel closed loop      | 5/5 PASS", 11, GREEN, True),
        ("T7   | localPlanner + pathFollower          | stop/slow_down safety signals     | 4/4 PASS", 11, GREEN, False),
        ("T8   | han_dog_bridge                      | cmd_vel \u2192 gRPC Walk() + watchdog   | 8/8 PASS", 11, GREEN, False),
    ], font=MONO)

    add_card(s, 0.55, 4.55, 5.8, 2.5, "T6 Full-Chain (Most Critical)")
    add_formula_box(s, 0.75, 5.02, 5.35, 1.75, [
        "# 6 real C++ nodes running simultaneously:",
        "terrainAnalysis + localPlanner + pathFollower",
        "pct_planner_astar + pct_path_adapter + TF",
        "",
        "# Test harness (Python):",
        "pub: synthetic cloud + odom (from cmd_vel integral)",
        "pub: /nav/goal_pose \u2192 trigger planning",
        "sub: /nav/cmd_vel \u2192 verify movement toward goal",
        "Result: 5/5 PASS in 18s (building2_9.pickle)",
    ])

    add_card(s, 6.55, 4.55, 6.2, 2.5, "T8 Driver Layer: han_dog_bridge")
    add_bullets(s, 6.75, 5.02, 5.7, 1.75, [
        "cmd_vel \u2192 _twist_to_walk (normalize to [-1,1]) \u2192 gRPC Walk()",
        "IMU: quaternion Hamilton(w,x,y,z) \u2192 ROS(x,y,z,w) \u2192 /Odometry",
        "Joint: AllJoints(16 DOF) \u2192 /robot_state",
        ("Watchdog: 200ms no cmd_vel \u2192 Walk(0,0,0)", RED),
        "SLAM position reset: every 5s from /nav/odometry",
        "Tested against real brainstem CMS on S100P",
        ("8/8 checks PASS including motor enable + history stream", GREEN),
    ], size=12)

    # ═══════════════ S7: SLAM VIBRATION — OVERVIEW ═══════════════
    s = ns(prs)
    add_header(s, "SLAM Vibration Robustness: Overview",
               "Point-LIO vs Fast-LIO2 on quadruped vibration data (Leg-KILO corridor)")

    add_image(s, FIG_3SCENARIO, 0.55, 1.3, 6.5, 4.3)
    add_caption(s, 0.55, 5.65, 6.5,
                "3-scenario comparison: corridor (445s), slope (162s), grass (27s)",
                PP_ALIGN.CENTER)

    add_card(s, 7.25, 1.3, 5.45, 2.4, "Test Setup")
    add_bullets(s, 7.45, 1.78, 5.0, 1.7, [
        "Dataset: Leg-KILO (Unitree Go1 quadruped)",
        "LiDAR: Velodyne VLP-16 (16-line, 10Hz)",
        "IMU: ~500Hz onboard",
        "Duration: 445s (corridor), 162s (slope), 27s (grass)",
        ("Vibration: leg-ground impacts + body pitch/roll", RED),
        "Platform: RDK S100P (aarch64, ROS2 Humble)",
    ], size=12)

    add_card(s, 7.25, 3.95, 5.45, 2.7, "Key Results")
    add_lines(s, 7.45, 4.42, 5.0, 2.0, [
        ("Algorithm    | Data      | Result | Z-drift", 12, TITLE_CLR, True),
        ("Point-LIO   | VLP-16    | PASS   | -0.88m", 11, GREEN, False),
        ("Fast-LIO2   | VLP-16    | N/A    | (no PointCloud2)", 11, SUB_CLR, False),
        ("Point-LIO   | Mid-360   | PASS   | 0.006m", 11, GREEN, False),
        ("Fast-LIO2   | Mid-360   | PASS   | 0.008m", 11, GREEN, False),
        ("Point-LIO   | Slope     | PASS   | -6.2m", 11, GREEN, False),
        ("Point-LIO   | Grass     | PASS*  | -19.4m", 11, ORANGE, False),
    ], font=MONO)
    add_caption(s, 7.25, 6.55, 5.4, "* Grass: survives but Z poor (vegetation occlusion)")

    # ═══════════════ S8: SLAM CORRIDOR DETAIL ═══════════════
    s = ns(prs)
    add_header(s, "SLAM Corridor Analysis: 445s Quadruped Walk",
               "Point-LIO trajectory, XYZ components, and 3D visualization")

    add_image(s, FIG_CORRIDOR, 0.55, 1.3, 4.0, 2.7)
    add_caption(s, 0.55, 4.05, 4.0, "2D trajectory with time labels and direction arrows",
                PP_ALIGN.CENTER)
    add_image(s, FIG_CORRIDOR_XYZ, 4.75, 1.3, 4.0, 2.7)
    add_caption(s, 4.75, 4.05, 4.0, "Position XYZ components vs time",
                PP_ALIGN.CENTER)
    add_image(s, FIG_CORRIDOR_3D, 8.95, 1.3, 3.75, 2.7)
    add_caption(s, 8.95, 4.05, 3.75, "3D trajectory with time coloring",
                PP_ALIGN.CENTER)

    add_card(s, 0.55, 4.35, 6.2, 2.7, "Quantitative Metrics")
    add_formula_box(s, 0.75, 4.82, 5.75, 2.0, [
        "Total duration:     445 seconds (7.4 minutes)",
        "Total trajectory:   ~220m (110m out + 110m return)",
        "Max displacement:   109.9m from start",
        "Loop closure error: 7.1m (returns near origin)",
        "Z-axis drift:       peak -0.88m (within +/-2m spec)",
        "IMU init time:      < 1 second",
        "Output rate:        10Hz stable throughout",
        "Crashes:            0",
    ])

    add_card(s, 6.95, 4.35, 5.75, 2.7, "Trajectory Path")
    add_bullets(s, 7.15, 4.82, 5.3, 2.0, [
        "Phase 1: Straight corridor (0\u219218m, 20s)",
        "Phase 2: Right turn + horizontal (18m\u21920m Y, 30s)",
        "Phase 3: Long corridor run (0\u2192-103m X, 100s)",
        ("Phase 4: Return path (-103\u21920m, ~200s)", ACCENT),
        "Loop closure: 7.1m residual (no PGO active)",
        ("Z-drift: peak -0.88m at t=250s, recovers to -0.3m", GREEN),
    ], size=12)

    # ═══════════════ S9: DUAL SLAM COMPARISON ═══════════════
    s = ns(prs)
    add_header(s, "Point-LIO vs Fast-LIO2: Dual Algorithm Comparison",
               "Same platform (Livox Mid-360), same scene | Static + real-time baseline")

    add_image(s, FIG_DUAL_TRAJ, 0.55, 1.3, 6.0, 2.7)
    add_caption(s, 0.55, 4.05, 6.0,
                "Side-by-side 2D trajectories (Point-LIO left, Fast-LIO2 right)",
                PP_ALIGN.CENTER)
    add_image(s, FIG_DUAL_ZDRIFT, 6.75, 1.3, 5.95, 2.7)
    add_caption(s, 6.75, 4.05, 5.95,
                "Z-axis drift comparison with divergence marker",
                PP_ALIGN.CENTER)

    add_image(s, FIG_DUAL_XYZ, 0.55, 4.3, 6.0, 2.7)
    add_caption(s, 0.55, 7.05, 6.0, "XYZ position components comparison", PP_ALIGN.CENTER)

    add_card(s, 6.75, 4.3, 5.95, 2.7, "Comparison Summary")
    add_lines(s, 6.95, 4.78, 5.5, 1.3, [
        ("Metric              | Point-LIO  | Fast-LIO2", 12, TITLE_CLR, True),
        ("Static drift        | 0.006m     | 0.008m", 11, GREEN, False),
        ("Input format        | Any PCL2   | Livox only", 11, TEXT_CLR, False),
        ("Quadruped vibration | PASS       | N/A (VLP-16)", 11, TEXT_CLR, False),
        ("Init time           | <1s        | ~2s", 11, TEXT_CLR, False),
        ("Output rate         | 10Hz       | 10Hz", 11, TEXT_CLR, False),
    ], font=MONO)
    add_text(s, 6.95, 6.3, 5.5, 0.4,
             "Both algorithms show comparable accuracy in static tests. Point-LIO is more versatile (any PointCloud2 input), making it preferred for quadruped applications.",
             size=11, color=SUB_CLR)

    # ═══════════════ S10: OUTDOOR TERRAIN ═══════════════
    s = ns(prs)
    add_header(s, "Outdoor Terrain SLAM: Slope & Grass",
               "Point-LIO robustness on challenging outdoor environments")

    add_image(s, FIG_SLOPE_TRAJ, 0.55, 1.3, 3.0, 2.4)
    add_caption(s, 0.55, 3.75, 3.0, "Hill/slope terrain 2D trajectory", PP_ALIGN.CENTER)
    add_image(s, FIG_SLOPE_PANEL, 3.75, 1.3, 4.7, 2.4)
    add_caption(s, 3.75, 3.75, 4.7, "Slope 6-panel analysis (XYZ, disp, Z, speed)",
                PP_ALIGN.CENTER)

    add_image(s, FIG_GRASS_TRAJ, 0.55, 4.1, 3.0, 2.4)
    add_caption(s, 0.55, 6.55, 3.0, "Grass terrain 2D trajectory", PP_ALIGN.CENTER)
    add_image(s, FIG_GRASS_PANEL, 3.75, 4.1, 4.7, 2.4)
    add_caption(s, 3.75, 6.55, 4.7, "Grass 6-panel analysis (Z-drift = -19.4m)",
                PP_ALIGN.CENTER)

    add_card(s, 8.65, 1.3, 4.05, 2.75, "Slope (162s, 346m)")
    add_bullets(s, 8.85, 1.78, 3.6, 1.9, [
        "Terrain: hill with elevation changes",
        "Z return: -6.2m \u2192 ground level",
        ("Z-drift: acceptable for outdoor", GREEN),
        "Speed: variable (uphill slower)",
        ("Result: PASS", GREEN),
    ], size=12)

    add_card(s, 8.65, 4.3, 4.05, 2.75, "Grass (27s, 47m)")
    add_bullets(s, 8.85, 4.78, 3.6, 1.9, [
        "Terrain: grass with vegetation",
        ("Z-drift: -19.4m (vegetation occlusion)", RED),
        "LiDAR penetrates grass canopy",
        "Ground estimation unreliable",
        ("Result: PASS* (survives, Z poor)", ORANGE),
    ], size=12)

    # ═══════════════ S11: PERCEPTION LAYER ═══════════════
    s = ns(prs)
    add_header(s, "Perception Layer: From Detection to Scene Graph",
               "YOLO-World + HOV-SG CLIP + ConceptGraphs + BA-HSG Belief | 6 paper implementations")

    # ── Left: Detection & CLIP ──
    add_card(s, 0.55, 1.3, 4.0, 3.0, "Object Detection & Feature Extraction")
    add_bullets(s, 0.75, 1.78, 3.55, 2.3, [
        ("YOLO-World: open-vocabulary detector", ACCENT),
        "  Arbitrary categories (vs Odin: 80 fixed COCO)",
        "  + Laplacian blur filter for quadruped gait",
        ("HOV-SG 3-source CLIP fusion [RSS 2024]:", ACCENT),
        "  f = 0.25*f_g + 0.50*f_l + 0.25*f_m",
        "  f_g: global image, f_l: crop, f_m: masked crop",
        ("Quality-aware EMA update [NaviMind \u00a73.2.3]:", ACCENT),
        "  \u03b1 = min(0.5, 0.3 * clamp(s/0.8, 0.5, 1.5))",
    ], size=12)

    # ── Middle: Scene Graph ──
    add_card(s, 4.75, 1.3, 4.0, 3.0, "4-Level Hierarchical Scene Graph")
    add_formula_box(s, 4.95, 1.78, 3.55, 2.3, [
        "# ConceptGraphs [ICRA 2024] + Hydra [RSS 2022]",
        "Floor",
        "  \u2514\u2500 Room (DBSCAN \u03b5=3m + LLM naming)",
        "      \u2514\u2500 Group (semantic category)",
        "          \u2514\u2500 Object (label, pos_3d, clip_512d,",
        "               score, detections, extent)",
        "",
        "# Edges: hierarchical + spatial relations",
        "#   near, on, left_of, right_of, in_front_of",
    ])

    # ── Right: Belief & Multi-view ──
    add_card(s, 8.95, 1.3, 3.75, 3.0, "Belief Propagation & Verification")
    add_bullets(s, 9.15, 1.78, 3.3, 2.3, [
        ("BA-HSG [NaviMind C1]:", ACCENT),
        "  Beta(\u03b1,\u03b2) existence belief per node",
        "  3-round graph diffusion (\u03b5=0.005)",
        ("SG-Nav multi-view credibility:", ACCENT),
        "  EMA over 10 frames, S_thresh=0.8",
        "  False positive penalty = 0.2",
        ("DBSCAN feature refinement:", ACCENT),
        "  Every 5 detections, cluster CLIP history",
        "  Remove outlier features (noise rejection)",
    ], size=12)

    # ── Bottom Left: Safety-Aware ──
    add_card(s, 0.55, 4.55, 4.0, 2.5, "Safety-Aware Perception")
    add_bullets(s, 0.75, 5.02, 3.55, 1.75, [
        ("Differential thresholds by category:", ACCENT),
        "  Safety classes (person, fire_extinguisher):",
        "    association=0.3, lower bar for detection",
        "  Normal classes: association=0.5",
        ("Odin: no safety differentiation", RED),
        "  Same 80-class COCO for all objects",
        ("Impact: critical objects detected earlier", GREEN),
    ], size=12)

    # ── Bottom Middle: Paper References ──
    add_card(s, 4.75, 4.55, 4.0, 2.5, "Paper Implementations in Perception")
    add_lines(s, 4.95, 5.02, 3.55, 1.75, [
        ("Paper               | Module             | What", 11, TITLE_CLR, True),
        ("ConceptGraphs[ICRA] | instance_tracker   | Scene graph", 10, TEXT_CLR, False),
        ("HOV-SG [RSS 2024]   | clip_encoder       | 3-source CLIP", 10, TEXT_CLR, False),
        ("SG-Nav [NeurIPS'24] | sgnav_reasoner     | Multi-view EMA", 10, TEXT_CLR, False),
        ("FSR-VLN [Horizon'25]| topo_memory        | Jaccard edges", 10, TEXT_CLR, False),
        ("BA-HSG [Ours]       | instance_tracker   | Belief prop.", 10, GREEN, False),
        ("SPADE [IROS 2025]   | topology_graph     | IG exploration", 10, TEXT_CLR, False),
    ], font=MONO)

    # ── Bottom Right: vs Odin ──
    add_card(s, 8.95, 4.55, 3.75, 2.5, "vs Odin Perception")
    add_lines(s, 9.15, 5.02, 3.3, 1.75, [
        ("Capability     | Odin  | NaviMind", 11, TITLE_CLR, True),
        ("Vocabulary     | 80 cls| Open", 10, GREEN, False),
        ("CLIP features  | None  | 3-source", 10, GREEN, False),
        ("Scene graph    | None  | 4-level HSG", 10, GREEN, False),
        ("Multi-view     | None  | 10-frame EMA", 10, GREEN, False),
        ("Belief model   | None  | Beta(\u03b1,\u03b2)+diffusion", 10, GREEN, False),
        ("Safety-aware   | No    | Diff thresholds", 10, GREEN, False),
        ("3D projection  | Depth | LiDAR+Depth", 10, GREEN, False),
    ], font=MONO)

    # ═══════════════ S12: SEMANTIC NAVIGATION ═══════════════
    s = ns(prs)
    add_header(s, "Semantic Navigation: Fast-Slow Dual Process",
               "Chinese VLN instructions | YOLO-World + CLIP + ConceptGraphs scene graph")

    add_card(s, 0.55, 1.3, 5.8, 2.8, "Fast Path (System 1, <200ms)")
    add_formula_box(s, 0.75, 1.78, 5.35, 2.05, [
        "# Confidence fusion:",
        "fused = 0.35*label + 0.35*clip + 0.15*det + 0.15*spatial",
        "",
        "# Scene graph keyword matching:",
        "instruction -> jieba tokenize -> match object labels",
        "threshold: fused >= 0.75 -> FAST PATH resolved",
        "",
        "# Hit rate target: >70%",
        "# Actual: ~80% on factory scene",
    ])

    add_card(s, 6.55, 1.3, 6.15, 2.8, "Slow Path (System 2, ~2s)")
    add_formula_box(s, 6.75, 1.78, 5.7, 2.05, [
        "# ESCA selective grounding:",
        "200 scene objects -> keyword filter -> ~15 objects",
        "Token reduction: 92.5%",
        "",
        "# LLM reasoning (Kimi / Qwen fallback):",
        "H-CoT 4-step prompt -> spatial reasoning",
        "Output: target position + confidence",
        "",
        "# AdaNav trigger: entropy > 1.5 AND conf < 0.85",
    ])

    add_card(s, 0.55, 4.35, 5.8, 2.7, "E2E Test Results (3/3 PASS)")
    add_lines(s, 0.75, 4.82, 5.35, 1.6, [
        ("Instruction (Chinese)            | Path | Time   | Result", 12, TITLE_CLR, True),
        ('"Navigate to factory gate"       | Fast | 10.9s  | PASS', 11, GREEN, False),
        ('"Navigate to machinery"          | Fast | 71.6s  | PASS', 11, GREEN, False),
        ('"Navigate to target region"      | Fast | 107s   | PASS', 11, GREEN, False),
    ], font=MONO)
    add_text(s, 0.75, 6.5, 5.3, 0.3,
             "Full pipeline: instruction \u2192 scene graph \u2192 global_planner \u2192 adapter \u2192 local_planner \u2192 cmd_vel",
             size=11, color=SUB_CLR)

    add_card(s, 6.55, 4.35, 6.15, 2.7, "Semantic Architecture")
    add_bullets(s, 6.75, 4.82, 5.7, 2.0, [
        ("Perception: YOLO-World + CLIP + HOV-SG 3-source fusion", ACCENT),
        "  encode_three_source: f_g(0.25) + f_l(0.50) + f_m(0.25)",
        ("Memory: ReMEmbR episodic (500-record FIFO)", ACCENT),
        "  Spatiotemporal: position + labels + room_type + timestamp",
        ("Topological: FSR-VLN Jaccard-weighted viewpoint edges", ACCENT),
        "  VLingMem region summaries per TopoNode",
        ("Failure: LERa 3-step recovery (retry/expand/requery/abort)", ACCENT),
    ], size=12)

    # ═══════════════ S13a: TSG TOPOLOGY EXPLORATION (NEW) ═══════════════
    s = ns(prs)
    add_header(s, "Topology-Aware Semantic Exploration (TSG)",
               "topology_graph.py (1129 LOC) | IG + Dijkstra + Traversal Memory | 33 tests PASS")

    # ── Left: TSG structure ──
    add_card(s, 0.55, 1.3, 4.0, 3.0, "TSG Graph Definition [NaviMind \u00a73.7]")
    add_formula_box(s, 0.75, 1.78, 3.55, 2.3, [
        "T = (N, E_T)  # Room-level abstraction",
        "",
        "Room node n_k:",
        "  (visited, visit_count, last_visit_time,",
        "   object_labels, predicted_room_type)",
        "",
        "Frontier node n_j:",
        "  (direction [dx,dy], size, predicted_room)",
        "",
        "Topology edge e_kl:",
        "  (type, traversal_count, last_traversed, conf)",
    ])

    # ── Middle: IG formula ──
    add_card(s, 4.75, 1.3, 4.0, 3.0, "Information Gain Scoring")
    add_formula_box(s, 4.95, 1.78, 3.55, 2.3, [
        "# Core formula:",
        "IG(n) = S_sem(n) x N(n) x U(n)",
        "",
        "# Novelty with temporal recovery:",
        "N(n) = 0.1 + 0.9*exp(-0.5*c_n)",
        "       * (1 - exp(-dt/120s))   if visited",
        "N(n) = 1.0                      if unvisited",
        "",
        "# Exploration target selection:",
        "n* = argmax IG(n) / (1 + 0.3*d_dijkstra)",
    ])

    # ── Right: Room connectivity detection ──
    add_card(s, 8.95, 1.3, 3.75, 3.0, "Room Connectivity Detection")
    add_bullets(s, 9.15, 1.78, 3.3, 2.3, [
        ("3 edge strategies [Hydra+SPADE]:", ACCENT),
        "  1. Door-mediated: door near 2 rooms",
        "     \u2192 high confidence edge",
        "  2. Proximity: centers < 2\u03b5_DBSCAN",
        "     \u2192 medium confidence edge",
        "  3. Passage: corridor connects all",
        ("2 frontier detection methods:", ACCENT),
        "  1. Door-outward: door at boundary",
        "  2. Sparse-sector: angular gaps",
        ("Traversal memory: actual paths", GREEN),
    ], size=12)

    # ── Bottom Left: Algorithm 2 ──
    add_card(s, 0.55, 4.55, 5.5, 2.5, "Algorithm 2: IG Exploration [NaviMind \u00a73.7.2]")
    add_formula_box(s, 0.75, 5.02, 5.05, 1.75, [
        "1: Update T from scene graph G",
        "2: Update frontier nodes from scene boundary",
        "3: Record robot position -> detect room transitions",
        "4: for each node n in T:",
        "5:   S_sem <- K.score(n.room_type, instruction)",
        "6:   IG(n) <- S_sem * novelty(n) * uncertainty(n)",
        "7: for each node n: score(n) = IG(n) / (1 + 0.3*d_dijkstra)",
        "8: return argmax score(n)  # ~1ms, zero LLM cost",
    ])

    # ── Bottom Right: Dual-layer + papers ──
    add_card(s, 6.25, 4.55, 6.45, 2.5, "Dual-Layer Strategy + Paper References")
    add_bullets(s, 6.45, 5.02, 3.1, 1.75, [
        ("Layer 1 (TSG, ~1ms):", GREEN),
        "  \u22652 rooms + edges \u2192 IG selects target",
        "  Covers ~70% exploration decisions",
        ("Layer 2 (LLM, ~2s fallback):", ACCENT),
        "  TSG insufficient \u2192 LLM with topo context",
        "  Room status + connections in prompt",
    ], size=12)
    add_lines(s, 9.65, 5.02, 2.85, 1.75, [
        ("Paper            | Contribution", 10, TITLE_CLR, True),
        ("SPADE [IROS'25]  | SG path planning", 9, TEXT_CLR, False),
        ("SG-Nav [NeurIPS] | Subgraph scoring", 9, TEXT_CLR, False),
        ("L3MVN [IROS'23]  | LLM frontier sel.", 9, TEXT_CLR, False),
        ("DovSG [RA-L'25]  | Dynamic SG update", 9, TEXT_CLR, False),
        ("Hydra [RSS'22]   | 5-level SG struct", 9, TEXT_CLR, False),
        ("EmbodiedRAG      | SG RAG retrieval", 9, TEXT_CLR, False),
    ], font=MONO)

    # ═══════════════ S13b: MEMORY + VoI (NEW) ═══════════════
    s = ns(prs)
    add_header(s, "Memory Systems & Adaptive Scheduling",
               "topological_memory.py (687 LOC) + episodic_memory.py + voi_scheduler.py | 5 paper impl.")

    # ── Left: Topological Memory ──
    add_card(s, 0.55, 1.3, 4.0, 3.0, "Topological Memory [FSR-VLN + VLingMem]")
    add_formula_box(s, 0.75, 1.78, 3.55, 2.3, [
        "# FSR-VLN Viewpoint Edges:",
        "for each new node n:",
        "  for neighbor m within 4m:",
        "    shared = labels(n) & labels(m)",
        "    jaccard = |shared| / |union|",
        "    weight = max(jaccard, 0.1)",
        "    add_edge(n, m, weight)",
        "",
        "# VLingMem Region Summary:",
        "node.region_summary =",
        '  f"objects:{labels}, room:{type}, {status}"',
    ])

    # ── Middle: Episodic Memory ──
    add_card(s, 4.75, 1.3, 4.0, 3.0, "Episodic Memory [ReMEmbR]")
    add_formula_box(s, 4.95, 1.78, 3.55, 2.3, [
        "# 500-record spatiotemporal FIFO",
        "record = {",
        "  position: [x, y, z],",
        "  labels: ['chair', 'desk', ...],",
        "  room_type: 'office',",
        "  timestamp: 1710234567.89,",
        "  clip_embedding: [512-dim]  # optional",
        "}",
        "",
        "# 3 retrieval modes:",
        "#   keyword, spatial (KDTree), temporal",
    ])

    # ── Right: VoI Scheduler ──
    add_card(s, 8.95, 1.3, 3.75, 3.0, "VoI Adaptive Scheduler")
    add_formula_box(s, 9.15, 1.78, 3.3, 2.3, [
        "# Utility function:",
        "U(a) = dE[S] - lam_t*dT",
        "       - lam_e*dE - lam_d*dd",
        "",
        "# dE[S]: expected score improvement",
        "# dT: time cost of action a",
        "# dE: energy cost",
        "# dd: deviation from current path",
        "",
        "# Result: 77% continue, 23% reperceive",
        "# Replaces heuristic fixed-interval",
    ])

    # ── Bottom Left: Text query with viewpoint boost ──
    add_card(s, 0.55, 4.55, 4.0, 2.5, "Topological Query Pipeline")
    add_bullets(s, 0.75, 5.02, 3.55, 1.75, [
        "1. Keyword match against node labels",
        "2. Score each TopoNode by relevance",
        ("3. 1-hop viewpoint neighbor boost:", ACCENT),
        "   For top nodes, check Jaccard neighbors",
        "   Propagate score through edges",
        ("4. Return top-K with positions", GREEN),
        "5. Feed into goal resolver or explorer",
    ], size=12)

    # ── Bottom Middle: Room coverage ──
    add_card(s, 4.75, 4.55, 4.0, 2.5, "Room Coverage & Backtracking")
    add_bullets(s, 4.95, 5.02, 3.55, 1.75, [
        ("get_room_coverage():", ACCENT),
        "  Per room_type: object count, visit count",
        "  \u2192 LLM prompt: explored areas summary",
        ("get_backtrack_position(steps_back=N):", ACCENT),
        "  Navigate to N-th previous position",
        "  Used by LERa recovery step 1 (LOOK)",
        ("get_least_visited_direction():", ACCENT),
        "  Bias exploration toward unvisited areas",
    ], size=12)

    # ── Bottom Right: Paper comparison ──
    add_card(s, 8.95, 4.55, 3.75, 2.5, "vs Prior Memory Systems")
    add_lines(s, 9.15, 5.02, 3.3, 1.75, [
        ("System      | Memory | Topo | VoI | Decay", 10, TITLE_CLR, True),
        ("SG-Nav      | None   | No   | No  | No", 9, SUB_CLR, False),
        ("LOVON       | None   | No   | No  | No", 9, SUB_CLR, False),
        ("Odin        | None   | No   | No  | No", 9, SUB_CLR, False),
        ("FSR-VLN     | View   | HMSG | No  | No", 9, TEXT_CLR, False),
        ("Mem4Nav     | Dual   | Yes  | No  | No", 9, TEXT_CLR, False),
        ("NaviMind    | Triple | Yes  | Yes | Yes", 9, GREEN, True),
    ], font=MONO)

    # ═══════════════ S14: RVIZ & VISUALIZATION ═══════════════
    s = ns(prs)
    add_header(s, "Navigation Visualization",
               "RViz frames + E2E trajectory analysis + terrain/path comparison")

    add_image(s, FIG_RVIZ_START, 0.55, 1.3, 3.0, 2.3)
    add_caption(s, 0.55, 3.65, 3.0, "Start position (RViz)", PP_ALIGN.CENTER)
    add_image(s, FIG_RVIZ_STAIR, 3.75, 1.3, 3.0, 2.3)
    add_caption(s, 3.75, 3.65, 3.0, "Staircase navigation", PP_ALIGN.CENTER)
    add_image(s, FIG_RVIZ_RF, 6.95, 1.3, 3.0, 2.3)
    add_caption(s, 6.95, 3.65, 3.0, "Roof level arrival", PP_ALIGN.CENTER)
    add_image(s, FIG_RVIZ_GOAL, 10.15, 1.3, 2.5, 2.3)
    add_caption(s, 10.15, 3.65, 2.5, "Goal reached", PP_ALIGN.CENTER)

    add_image(s, FIG_E2E_CMP_STR, 0.55, 4.1, 4.0, 2.8)
    add_caption(s, 0.55, 6.95, 4.0, "Straight vs A* path comparison", PP_ALIGN.CENTER)
    add_image(s, FIG_NAV_PCD, 4.75, 4.1, 4.0, 2.8)
    add_caption(s, 4.75, 6.95, 4.0, "Navigation on point cloud map", PP_ALIGN.CENTER)
    add_image(s, FIG_SIM_BUILD2V, 8.95, 4.1, 3.75, 2.8)
    add_caption(s, 8.95, 6.95, 3.75, "Building scenario trajectory", PP_ALIGN.CENTER)

    # ═══════════════ S13: QUANTITATIVE SUMMARY ═══════════════
    s = ns(prs)
    add_header(s, "Quantitative Results Summary",
               "All test modules \u2014 navigation, SLAM, integration, semantic")

    add_card(s, 0.55, 1.3, 7.0, 2.7, "Complete Test Matrix")
    add_lines(s, 0.75, 1.78, 6.55, 2.0, [
        ("Category              | Tests | Passed | Rate", 12, TITLE_CLR, True),
        ("Integration (T1-T8)   | 36    | 36     | 100%", 11, GREEN, False),
        ("Navigation E2E        | 6     | 6      | 100%", 11, GREEN, False),
        ("Semantic VLN          | 3     | 3      | 100%", 11, GREEN, False),
        ("SLAM vibration        | 4     | 3+1*   | 100%", 11, GREEN, False),
        ("Unit tests (nav_core) | 73    | 73     | 100%", 11, GREEN, False),
        ("Planning stub         | 22    | 22     | 100%", 11, GREEN, False),
        ("TOTAL                 | 144   | 143+1* | 99%+", 11, GREEN, True),
    ], font=MONO)

    add_card(s, 7.75, 1.3, 4.95, 2.7, "Key Performance Metrics")
    add_formula_box(s, 7.95, 1.78, 4.5, 2.0, [
        "SLAM static drift:     0.006-0.008m",
        "SLAM vibration Z:      <2m (220m trajectory)",
        "SLAM output rate:      10Hz stable",
        "Fast Path latency:     <200ms",
        "Fast Path hit rate:    ~80%",
        "Navigation E2E:        6/6 PASS",
        "T6 full chain:         18s goal-to-arrival",
        "Semantic VLN:          10.9-107s (3/3 PASS)",
    ])

    # Dashboard figure
    add_image(s, FIG_BAR_CMP, 0.55, 4.25, 4.0, 2.8)
    add_caption(s, 0.55, 7.1, 4.0, "Quantitative metrics bar chart", PP_ALIGN.CENTER)
    add_image(s, FIG_COMPARE_DASH, 4.75, 4.25, 4.0, 2.8)
    add_caption(s, 4.75, 7.1, 4.0, "Cross-scenario comparison dashboard", PP_ALIGN.CENTER)
    add_image(s, FIG_OVERLAY, 8.95, 4.25, 3.75, 2.8)
    add_caption(s, 8.95, 7.1, 3.75, "Overlay: corridor + slope + grass", PP_ALIGN.CENTER)

    # ═══════════════ S14: PLATFORM & DEPLOYMENT ═══════════════
    s = ns(prs)
    add_header(s, "Platform & Deployment Architecture",
               "RDK S100P (aarch64) | Dual-board: Nav Board + Dog Board | gRPC + Protobuf")

    add_card(s, 0.55, 1.3, 4.0, 3.0, "Hardware Platform")
    add_bullets(s, 0.75, 1.78, 3.55, 2.3, [
        ("RDK S100P (aarch64)", ACCENT),
        "  Horizon X5 processor",
        "  ROS2 Humble, Ubuntu 22.04",
        "Livox Mid-360 LiDAR",
        "Orbbec RGB-D camera (optional)",
        ("Dual-board architecture:", ACCENT),
        "  Nav Board: SLAM + planning + gRPC",
        "  Dog Board: RL policy + motor control",
    ], size=13)

    add_card(s, 4.75, 1.3, 4.0, 3.0, "Software Stack")
    add_bullets(s, 4.95, 1.78, 3.55, 2.3, [
        "SLAM: C++ (Fast-LIO2 / Point-LIO)",
        "Planning: C++ (local) + Python (global)",
        "Semantic: Python (YOLO-World + CLIP)",
        "Driver: Python (han_dog_bridge, gRPC)",
        "Control: Dart (brainstem, 50Hz loop)",
        "Client: Flutter (Android/Windows/iOS)",
        ("Protocol: gRPC + Protobuf (port 50051)", ACCENT),
    ], size=13)

    add_card(s, 8.95, 1.3, 3.75, 3.0, "Deployment Options")
    add_bullets(s, 9.15, 1.78, 3.3, 2.3, [
        "Systemd: 7 services (bare-metal)",
        "Docker: Supervisord chain",
        "  nav-lidar \u2192 nav-slam \u2192 nav-autonomy",
        "  \u2192 nav-planning \u2192 nav-grpc",
        "OTA updates (Ed25519 signed)",
        "CycloneDDS (UDP, no shared mem)",
        ("Health check: 13-point script", GREEN),
    ], size=13)

    add_card(s, 0.55, 4.55, 6.0, 2.5, "ROS2 Topic Contract")
    add_formula_box(s, 0.75, 5.02, 5.55, 1.75, [
        "/nav/odometry        Odometry    SLAM position output",
        "/nav/map_cloud       PointCloud2 World-frame map cloud",
        "/nav/terrain_map     PointCloud2 Traversability map",
        "/nav/global_path     Path        PCT planner output",
        "/nav/way_point       PointStamped Waypoint to local planner",
        "/nav/cmd_vel         TwistStamped Velocity to robot",
        "/nav/planner_status  String       IDLE/PLANNING/SUCCESS/STUCK",
        "/nav/adapter_status  String(JSON)  Waypoint tracking events",
    ])

    add_card(s, 6.75, 4.55, 5.95, 2.5, "gRPC Gateway Services")
    add_bullets(s, 6.95, 5.02, 5.5, 1.75, [
        ("TelemetryService: FastState(10Hz) + SlowState(1Hz)", ACCENT),
        "  pose, velocity, battery, SLAM status, navigation status",
        ("ControlService: lease + mode + e-stop + tasks", ACCENT),
        "  AcquireLease, SetMode, StartTask, StreamTeleop",
        ("MapService: save/load/list/delete maps", ACCENT),
        ("SystemService: health, config, OTA, logs", ACCENT),
    ], size=12)

    # ═══════════════ S15: CONTRIBUTIONS & FUTURE ═══════════════
    s = ns(prs)
    add_header(s, "Contributions & Future Work",
               "From simulation validation to production deployment")

    add_card(s, 0.55, 1.3, 4.0, 3.0, "Key Contributions")
    add_bullets(s, 0.75, 1.78, 3.55, 2.3, [
        ("Full-stack MuJoCo SITL with ONNX policy", GREEN),
        "  16-DOF locomotion + LiDAR + navigation",
        ("SLAM vibration robustness validated", GREEN),
        "  Point-LIO: 445s, 220m, quadruped gait",
        ("Semantic VLN with Chinese instructions", GREEN),
        "  Fast-Slow + AdaNav + LERa recovery",
        ("T1-T8 node-level integration coverage", GREEN),
        "  36/36 checks on real ROS2 nodes",
    ], size=12)

    add_card(s, 4.75, 1.3, 4.0, 3.0, "Current Limitations")
    add_bullets(s, 4.95, 1.78, 3.55, 2.3, [
        ("SLAM: no loop closure in vibration test", ORANGE),
        "  7.1m residual on 220m trajectory",
        ("Semantic: Fast Path rule-based only", ORANGE),
        "  Not learned policies (keyword match)",
        ("ONNX sim: MuJoCo freejoint, not real contact", ORANGE),
        "  MuJoCo <-> real terrain gap",
        ("Outdoor grass: Z-drift -19.4m", RED),
        "  Vegetation penetration unsolved",
    ], size=12)

    add_card(s, 8.95, 1.3, 3.75, 3.0, "Next Steps")
    add_bullets(s, 9.15, 1.78, 3.3, 2.3, [
        "1. Real-world deployment on Thunder",
        "   (pending motor chain validation)",
        "2. PGO loop closure for long trajectory",
        "3. Learned Fast Path (CLIP embedding)",
        "4. Multi-floor elevator integration",
        ("5. First paying customer (priority)", GREEN),
        "6. VLA navigation experiments",
    ], size=12)

    # Highlight metrics
    add_card(s, 0.55, 4.55, 12.15, 2.5, "System Readiness Dashboard")
    metrics = [
        ("144", "Total Tests", GREEN),
        ("99%+", "Pass Rate", GREEN),
        ("10Hz", "SLAM Output", ACCENT),
        ("<200ms", "Fast Path", ACCENT),
        ("v1.7.5", "Stable Release", GREEN),
        ("8/8", "T1-T8 PASS", GREEN),
    ]
    for i, (val, label, clr) in enumerate(metrics):
        cx = 0.75 + i * 2.0
        add_shape(s, cx, 5.0, 1.7, 1.7, WHITE, CARD_LINE, rounded=True)
        add_text(s, cx, 5.15, 1.7, 0.45, val, size=26, color=clr,
                 bold=True, align=PP_ALIGN.CENTER)
        add_text(s, cx, 5.65, 1.7, 0.7, label, size=11, color=SUB_CLR,
                 align=PP_ALIGN.CENTER)

    # ═══════════════ S16: THANK YOU ═══════════════
    s = ns(prs); set_bg(s, BG_DARK)
    add_text(s, 0.8, 2.0, 11.5, 0.6, "Thank You", size=40,
             color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, 0.8, 2.8, 11.5, 0.4, "Questions & Discussion",
             size=22, color=RGBColor(0x88, 0xCC, 0xEE), align=PP_ALIGN.CENTER)
    add_shape(s, 3.0, 3.8, 7.3, 0.02, RGBColor(0x33, 0x44, 0x66))
    add_text(s, 1.5, 4.2, 10.3, 0.35,
             "LingTu v1.7.5  |  brain/lingtu  |  ROS2 Humble  |  RDK S100P",
             size=14, color=RGBColor(0xAA, 0xAA, 0xBB), align=PP_ALIGN.CENTER)
    add_lines(s, 2.5, 5.0, 8.3, 1.5, [
        ("config/robot_config.yaml    \u2014 Robot parameters (single source of truth)", 12, RGBColor(0xAA, 0xAA, 0xBB), False),
        ("config/topic_contract.yaml  \u2014 ROS2 topic interface contract", 12, RGBColor(0xAA, 0xAA, 0xBB), False),
        ("src/semantic_planner/       \u2014 Fast-Slow VLN planner (9,645 LOC)", 12, RGBColor(0xAA, 0xAA, 0xBB), False),
        ("src/remote_monitoring/      \u2014 gRPC gateway (44 files, C++)", 12, RGBColor(0xAA, 0xAA, 0xBB), False),
        ("tests/integration/          \u2014 T1-T8 node-level verification", 12, RGBColor(0xAA, 0xAA, 0xBB), False),
    ], font=MONO)
    add_text(s, 0.8, 6.8, 11.5, 0.2,
             "Shanghai Qiongpei Technology  |  Generated 2026-03-11",
             size=11, color=RGBColor(0x66, 0x66, 0x77), align=PP_ALIGN.CENTER)

    # ──── Save ────
    prs.save(str(OUT))
    print(f"Saved: {OUT}")
    print(f"Slides: {len(prs.slides)}")
    return OUT


if __name__ == "__main__":
    build()
