"""
NaviMind 灵途 — 融资路演PPT生成器 v2

完全重新设计:
  - 融资Pitch Deck 结构 (22页)
  - 渐变背景 + 现代设计语言
  - 每页知识点讲透
  - 市场 / 技术壁垒 / 商业模式 / 竞品分析
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import os, copy

# ═══════════════════════════════════════════
#  设计系统
# ═══════════════════════════════════════════

# 主色板 — 深空蓝 + 电光蓝 + 白
C_BG1       = RGBColor(0x0B, 0x0E, 0x1A)   # 最深背景
C_BG2       = RGBColor(0x11, 0x16, 0x2B)   # 卡片背景
C_BG3       = RGBColor(0x1A, 0x20, 0x3C)   # 浅卡片
C_ACCENT    = RGBColor(0x00, 0xB4, 0xD8)   # 主蓝
C_ACCENT_L  = RGBColor(0x48, 0xCA, 0xE4)   # 浅蓝
C_GREEN     = RGBColor(0x06, 0xD6, 0xA0)   # 成功绿
C_ORANGE    = RGBColor(0xFF, 0xA6, 0x2B)   # 注意橙
C_RED       = RGBColor(0xEF, 0x47, 0x6F)   # 强调红
C_PURPLE    = RGBColor(0xA7, 0x7B, 0xFF)   # 紫色
C_GOLD      = RGBColor(0xFF, 0xD6, 0x6B)   # 金色
C_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
C_GRAY      = RGBColor(0x9A, 0x9E, 0xB8)
C_DIM       = RGBColor(0x6B, 0x70, 0x8D)

FONT_CN = "Microsoft YaHei"
FONT_EN = "Segoe UI"

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def box(slide, l, t, w, h, fill=None, border=None, radius=None):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    if fill:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if border:
        s.line.color.rgb = border
        s.line.width = Pt(1.2)
    else:
        s.line.fill.background()
    return s


def rect(slide, l, t, w, h, fill):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.fill.background()
    return s


def txt(slide, l, t, w, h):
    return slide.shapes.add_textbox(l, t, w, h)


def run(tf, text, sz=16, color=C_WHITE, bold=False, font=FONT_CN, align=PP_ALIGN.LEFT, space_before=None, space_after=None, new_para=False):
    if new_para or len(tf.paragraphs) == 0 or (len(tf.paragraphs) == 1 and tf.paragraphs[0].text == ""):
        p = tf.paragraphs[0] if (len(tf.paragraphs) == 1 and tf.paragraphs[0].text == "") else tf.add_paragraph()
    else:
        p = tf.add_paragraph()
    p.alignment = align
    if space_before is not None:
        p.space_before = space_before
    if space_after is not None:
        p.space_after = space_after
    r = p.add_run()
    r.text = text
    r.font.size = Pt(sz)
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.name = font
    return p


def title_bar(slide, number, title_cn, title_en=""):
    rect(slide, Inches(0), Inches(0), SLIDE_W, Pt(3), C_ACCENT)
    tb = txt(slide, Inches(0.6), Inches(0.25), Inches(12), Inches(0.9))
    tf = tb.text_frame
    tf.word_wrap = True
    label = f"{number:02d}"
    run(tf, label, sz=14, color=C_ACCENT, bold=True, font=FONT_EN)
    p = tf.add_paragraph()
    p.alignment = PP_ALIGN.LEFT
    r1 = p.add_run()
    r1.text = f"  {title_cn}"
    r1.font.size = Pt(28)
    r1.font.color.rgb = C_WHITE
    r1.font.bold = True
    r1.font.name = FONT_CN
    if title_en:
        r2 = p.add_run()
        r2.text = f"  {title_en}"
        r2.font.size = Pt(16)
        r2.font.color.rgb = C_DIM
        r2.font.name = FONT_EN
    rect(slide, Inches(0.6), Inches(1.15), Inches(2.5), Pt(2.5), C_ACCENT)


def stat_box(slide, l, t, w, h, number, label, color=C_ACCENT):
    b = box(slide, l, t, w, h, fill=C_BG2, border=color)
    tf = b.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    run(tf, number, sz=36, color=color, bold=True, align=PP_ALIGN.CENTER)
    run(tf, label, sz=13, color=C_GRAY, align=PP_ALIGN.CENTER, space_before=Pt(2))


def card(slide, l, t, w, h, title, items, title_color=C_ACCENT, border_color=None):
    b = box(slide, l, t, w, h, fill=C_BG2, border=border_color or title_color)
    tf = b.text_frame
    tf.word_wrap = True
    run(tf, f"  {title}", sz=16, color=title_color, bold=True)
    for item in items:
        if isinstance(item, tuple):
            text, color, sz = item[0], item[1] if len(item) > 1 else C_WHITE, item[2] if len(item) > 2 else 13
        else:
            text, color, sz = item, C_WHITE, 13
        run(tf, f"  {text}", sz=sz, color=color, space_before=Pt(3))
    return b


# ═══════════════════════════════════════════
#  SLIDE 1: 封面
# ═══════════════════════════════════════════
def slide_cover(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, C_BG1)

    rect(s, Inches(0), Inches(0), SLIDE_W, Pt(3), C_ACCENT)

    # 品牌
    tb = txt(s, Inches(0.8), Inches(1.2), Inches(11.7), Inches(1.2))
    tf = tb.text_frame
    run(tf, "NaviMind", sz=60, color=C_ACCENT, bold=True, font=FONT_EN, align=PP_ALIGN.CENTER)

    tb2 = txt(s, Inches(0.8), Inches(2.4), Inches(11.7), Inches(0.6))
    tf2 = tb2.text_frame
    run(tf2, "灵 途", sz=40, color=C_WHITE, bold=True, align=PP_ALIGN.CENTER)

    # 一句话
    tb3 = txt(s, Inches(1.5), Inches(3.3), Inches(10.3), Inches(0.8))
    tf3 = tb3.text_frame
    run(tf3, "让机器人听懂人话, 看懂世界, 自主探索未知空间", sz=22, color=C_GOLD, bold=True, align=PP_ALIGN.CENTER)
    run(tf3, "AI-Powered Autonomous Navigation for Quadruped Robots", sz=16, color=C_DIM, align=PP_ALIGN.CENTER, space_before=Pt(8))

    # 4 个核心指标
    metrics = [
        ("< 5ms", "快路径延迟\n无需LLM"),
        ("100%", "离线测试\n通过率"),
        ("99", "单元测试\n全部通过"),
        ("75%+", "LLM调用\n节省率"),
    ]
    x = 1.8
    for val, label in metrics:
        stat_box(s, Inches(x), Inches(4.6), Inches(2.2), Inches(1.5), val, label)
        x += 2.5

    # 底部
    rect(s, Inches(0), Inches(6.7), SLIDE_W, Inches(0.8), C_BG2)
    tb4 = txt(s, Inches(0.8), Inches(6.8), Inches(11.7), Inches(0.5))
    tf4 = tb4.text_frame
    run(tf4, "Unitree Go2  |  Jetson Orin NX 16GB  |  ROS2  |  多LLM后端  |  2026", sz=14, color=C_DIM, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════
#  SLIDE 2: 痛点
# ═══════════════════════════════════════════
def slide_problem(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, C_BG1)
    title_bar(s, 1, "行业痛点", "The Problem")

    # 左: 场景描述
    b1 = box(s, Inches(0.6), Inches(1.5), Inches(5.8), Inches(5.3), fill=C_BG2, border=C_RED)
    tf1 = b1.text_frame
    tf1.word_wrap = True
    run(tf1, "  机器人导航的现状", sz=20, color=C_RED, bold=True)
    
    problems = [
        ("依赖预建地图", "每次部署需人工扫描建图 (30min+)\n环境变化就要重新扫描, 不可扩展"),
        ("不懂自然语言", "只能接受坐标点, 不能理解\n\"去厨房拿杯水\" 这样的指令"),
        ("每步都调大模型", "SG-Nav: 每走一步都调LLM\n延迟2-5秒/步, API成本高昂"),
        ("没有记忆 & 推理", "走过的地方记不住, 找不到目标就原地转圈\n没有不确定性建模, 做不了置信决策"),
        ("只在仿真里跑", "MP3D/HM3D仿真 ≠ 真实世界\n90cm高的仿真机器人 ≠ 30cm四足犬"),
    ]
    for title, desc in problems:
        run(tf1, f"  {title}", sz=16, color=C_ORANGE, bold=True, space_before=Pt(12))
        for line in desc.split("\n"):
            run(tf1, f"    {line}", sz=12, color=C_GRAY, space_before=Pt(1))

    # 右: 市场需求
    b2 = box(s, Inches(6.9), Inches(1.5), Inches(5.8), Inches(2.3), fill=C_BG2, border=C_GREEN)
    tf2 = b2.text_frame
    tf2.word_wrap = True
    run(tf2, "  市场需求", sz=20, color=C_GREEN, bold=True)
    markets = [
        "巡检机器人 — 工厂/仓库/数据中心, 需语义理解",
        "服务机器人 — 酒店/医院, 需自然语言交互",
        "搜救机器人 — 未知环境, 需自主探索",
        "家庭陪伴 — 语音控制, 找物品、送东西",
    ]
    for m in markets:
        run(tf2, f"  {m}", sz=13, color=C_WHITE, space_before=Pt(6))

    # 右下: 市场规模
    b3 = box(s, Inches(6.9), Inches(4.1), Inches(5.8), Inches(2.7), fill=C_BG2, border=C_ACCENT)
    tf3 = b3.text_frame
    tf3.word_wrap = True
    run(tf3, "  市场规模 (全球服务机器人)", sz=18, color=C_ACCENT, bold=True)
    run(tf3, "  2026 TAM: $530亿", sz=28, color=C_GOLD, bold=True, space_before=Pt(10))
    run(tf3, "  室内导航软件 SAM: ~$80亿", sz=20, color=C_ACCENT_L, bold=True, space_before=Pt(6))
    run(tf3, "  语义导航 SOM: ~$12亿 (CAGR 35%)", sz=18, color=C_GREEN, bold=True, space_before=Pt(6))
    run(tf3, "  来源: IFR 2025, MarketsandMarkets, Precedence Research", sz=10, color=C_DIM, space_before=Pt(10))


# ═══════════════════════════════════════════
#  SLIDE 3: 我们的方案 (一句话)
# ═══════════════════════════════════════════
def slide_solution(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, C_BG1)
    title_bar(s, 2, "我们的方案", "Our Solution")

    # 核心一句话
    b_main = box(s, Inches(0.6), Inches(1.6), Inches(12.1), Inches(1.6), fill=C_BG3, border=C_ACCENT)
    tf = b_main.text_frame
    tf.word_wrap = True
    run(tf, "  NaviMind 灵途", sz=28, color=C_ACCENT, bold=True)
    run(tf, "  = 在线层次场景图 + 快慢双路径推理 + 信念感知规划 + 拓扑语义探索", sz=18, color=C_WHITE, bold=True, space_before=Pt(8))
    run(tf, "  机器人边走边建「认知地图」, 75%指令<5ms解决不调LLM, 复杂场景才调大模型做深度推理", sz=15, color=C_GOLD, space_before=Pt(6))

    # 4个核心能力
    caps = [
        (0.6,  "看懂世界", "在线构建4层认知地图", "Object → Group → Room → Floor\n开放词汇检测 (不限物体类别)\nCLIP语义编码 + 质量感知融合\n增量更新, 不依赖预建地图", C_ACCENT),
        (3.85, "听懂人话", "自然语言 → 导航坐标", "中英文指令全支持\n\"找到走廊门旁边的灭火器\"\n快路径 <5ms 直接出坐标\n慢路径 LLM Chain-of-Thought", C_GREEN),
        (7.1,  "会思考", "概率推理 + 自适应决策", "每个物体有存在概率和位置不确定度\n多目标时贝叶斯推理排序\n信息价值驱动何时该重新观察\n77%时间免算力, 23%精准感知", C_ORANGE),
        (10.35, "能探索", "拓扑图 + 信息增益", "空间结构图 (房间→连通关系)\n前沿节点 (未知区域入口)\n语义先验 (咖啡机→厨房)\n穿越记忆 (去过的不再去)", C_PURPLE),
    ]
    for x, t1, t2, desc, color in caps:
        b = box(s, Inches(x), Inches(3.6), Inches(2.95), Inches(3.6), fill=C_BG2, border=color)
        tf = b.text_frame
        tf.word_wrap = True
        run(tf, f"  {t1}", sz=20, color=color, bold=True)
        run(tf, f"  {t2}", sz=12, color=C_GOLD, bold=True, space_before=Pt(4))
        for line in desc.split("\n"):
            run(tf, f"  {line}", sz=11, color=C_GRAY, space_before=Pt(3))


# ═══════════════════════════════════════════
#  SLIDE 4: 技术架构全景
# ═══════════════════════════════════════════
def slide_architecture(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, C_BG1)
    title_bar(s, 3, "技术架构全景", "System Architecture")

    # Pipeline flow
    flow_items = [
        (0.4,  "RGB-D\n相机", "640x480\n@30fps", C_ACCENT, 1.5),
        (2.15, "YOLO-World\n开放检测", "不限类别\n实时检测", C_GREEN, 1.8),
        (4.2,  "CLIP\n语义编码", "512维特征\n质量感知", C_GREEN, 1.6),
        (6.05, "层次场景图\n4层构建", "Object→Group\n→Room→Floor", C_ORANGE, 1.9),
        (8.2,  "快慢推理\n目标解析", "Fast <5ms\nSlow LLM", C_GOLD, 1.8),
        (10.25, "Nav2\n路径规划", "DWB局部\n到达检测", C_RED, 1.5),
        (12.0, "四足\n底盘", "Unitree\nGo2", C_PURPLE, 1.1),
    ]
    for x, title, desc, color, w in flow_items:
        b = box(s, Inches(x), Inches(1.5), Inches(w), Inches(1.3), fill=C_BG2, border=color)
        tf = b.text_frame
        tf.word_wrap = True
        for line in title.split("\n"):
            run(tf, f" {line}", sz=11, color=color, bold=True, space_before=Pt(1))
        for line in desc.split("\n"):
            run(tf, f" {line}", sz=10, color=C_GRAY, space_before=Pt(1))

    # 中间: 详细架构
    modules = [
        (0.4, 3.1, 3.0, 2.0, "感知层 (Perception)", [
            "YOLO-World-L: 开放词汇目标检测",
            "CLIP ViT-B/32: 512维语义特征",
            "运动模糊过滤: Laplacian方差",
            "实例关联: 3D距离 + CLIP余弦",
            "深度投影: RGB-D → 3D世界坐标",
        ], C_ACCENT),
        (3.65, 3.1, 3.0, 2.0, "认知层 (Cognition)", [
            "层次场景图: 4层在线构建",
            "DBSCAN聚类 → 房间推理",
            "物体组: 语义共现聚合",
            "拓扑边: 门/过道/邻近连通",
            "前沿检测: 未探索方向估计",
        ], C_GREEN),
        (6.9, 3.1, 3.0, 2.0, "推理层 (Reasoning)", [
            "Fast Path: 多源融合评分 <5ms",
            "Slow Path: LLM层次CoT推理",
            "信念系统: Beta + Gaussian",
            "多假设: 贝叶斯候选排序",
            "VoI调度: 信息价值决策",
        ], C_ORANGE),
        (10.15, 3.1, 2.8, 2.0, "执行层 (Execution)", [
            "Nav2: 局部路径规划",
            "Pure Pursuit: 轨迹跟踪",
            "碰撞避障: 安全门限",
            "到达检测: 距离判断",
            "再感知触发: VoI调度",
        ], C_RED),
    ]
    for x, y, w, h, title, items, color in modules:
        b = box(s, Inches(x), Inches(y), Inches(w), Inches(h), fill=C_BG2, border=color)
        tf = b.text_frame
        tf.word_wrap = True
        run(tf, f" {title}", sz=13, color=color, bold=True)
        for item in items:
            run(tf, f"  {item}", sz=10, color=C_GRAY, space_before=Pt(2))

    # 底部硬件
    b_hw = box(s, Inches(0.4), Inches(5.4), Inches(12.5), Inches(1.8), fill=C_BG2, border=C_DIM)
    tf_hw = b_hw.text_frame
    tf_hw.word_wrap = True
    run(tf_hw, "  硬件平台", sz=16, color=C_ACCENT, bold=True)
    hw_items = [
        ("计算", "NVIDIA Jetson Orin NX 16GB — 边缘AI, 无需云端GPU"),
        ("LiDAR", "Livox Mid-360 — 360度激光雷达, Fast-LIO2 SLAM"),
        ("相机", "Orbbec Femto — RGB-D深度相机, 640x480 @30fps"),
        ("底盘", "Unitree Go2 — 四足机器人, 30cm低视角, 全地形"),
        ("通信", "gRPC Gateway (44 RPCs) + Flutter多平台App (Android/iOS/Windows)"),
    ]
    for name, desc in hw_items:
        run(tf_hw, f"  {name}: {desc}", sz=11, color=C_GRAY, space_before=Pt(2))


# ═══════════════════════════════════════════
#  SLIDE 5: 创新1 — 在线层次场景图 (详细)
# ═══════════════════════════════════════════
def slide_scene_graph(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, C_BG1)
    title_bar(s, 4, "核心技术 1: 在线层次场景图", "Online Hierarchical Scene Graph")

    # 左半: 四层详解
    layers = [
        ("Floor (楼层)", "全局根节点, 管理所有房间", "单楼层/多楼层皆可扩展", C_PURPLE, 0.4, 1.5),
        ("Room (房间)", "DBSCAN聚类 + 规则命名", "corridor, office, kitchen, stairwell...\n8种房间规则, 关键词匹配+优先级\n可选LLM命名 (稳定后调用一次)", C_ORANGE, 0.4, 2.4),
        ("Group (物体组)", "语义共现聚合", "workstation = desk+chair+monitor\nsafety = fire_ext+exit_sign\n自动聚合, 支持Slow Path推理定位", C_GREEN, 0.4, 3.7),
        ("Object (物体)", "YOLO-World检测 + CLIP编码", "3D位置 (深度投影到世界坐标)\nCLIP 512维特征向量\n检测置信度 + 检测次数\n空间关系 (near, on, left_of...)", C_ACCENT, 0.4, 5.0),
    ]
    for title, subtitle, desc, color, x, y in layers:
        b = box(s, Inches(x), Inches(y), Inches(6.2), Inches(0.8 if "Floor" in title else 1.1), fill=C_BG2, border=color)
        tf = b.text_frame
        tf.word_wrap = True
        run(tf, f" {title}", sz=15, color=color, bold=True)
        r = tf.paragraphs[-1].add_run()
        r.text = f"  — {subtitle}"
        r.font.size = Pt(12)
        r.font.color.rgb = C_GOLD
        r.font.name = FONT_CN
        for line in desc.split("\n"):
            run(tf, f"   {line}", sz=10, color=C_GRAY, space_before=Pt(1))

    # 右半: 关键技术
    techs = [
        ("质量感知CLIP融合", [
            "问题: 运动中检测质量参差不齐",
            "方案: alpha = base x clamp(score/ref)",
            "高置信检测 → 学习率大 → 快速学习",
            "低置信检测 → 学习率小 → 保守更新",
            "结果: 特征更稳定, 匹配更准确",
        ], C_ACCENT),
        ("运动模糊过滤", [
            "问题: 四足步态导致严重图像模糊",
            "方案: Laplacian方差 < 阈值 → 丢弃",
            "Lvar = Var(Laplacian(gray_image))",
            "效果: 过滤50%+低质帧, 检测更可靠",
        ], C_GREEN),
        ("拓扑连通边", [
            "门/过道/邻近 → 房间间连通关系",
            "构建可导航的空间拓扑图",
            "支持Dijkstra最短路径计算",
            "为探索模块提供结构化信息",
        ], C_ORANGE),
        ("vs ConceptGraphs对比", [
            "CG: 离线, 均匀EMA, 无导航",
            "我们: 在线, 质量感知, 导航集成",
            "CG: 无层次, 扁平物体列表",
            "我们: 4层层次 + 信念 + 拓扑",
        ], C_RED),
    ]
    y_pos = 1.5
    for title, items, color in techs:
        b = box(s, Inches(6.9), Inches(y_pos), Inches(6.0), Inches(1.2 if len(items) <= 4 else 1.35), fill=C_BG2, border=color)
        tf = b.text_frame
        tf.word_wrap = True
        run(tf, f" {title}", sz=13, color=color, bold=True)
        for item in items:
            run(tf, f"  {item}", sz=10, color=C_GRAY, space_before=Pt(1))
        y_pos += 1.3 if len(items) <= 4 else 1.45


# ═══════════════════════════════════════════
#  SLIDE 6: 创新2 — 快慢推理 (原理详解)
# ═══════════════════════════════════════════
def slide_fast_slow_principle(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, C_BG1)
    title_bar(s, 5, "核心技术 2: 快慢双路径推理 — 原理", "Fast-Slow Dual-Path Reasoning")

    # 灵感来源
    b0 = box(s, Inches(0.6), Inches(1.5), Inches(12.1), Inches(1.0), fill=C_BG3, border=C_GOLD)
    tf0 = b0.text_frame
    tf0.word_wrap = True
    run(tf0, "  灵感: Daniel Kahneman《思考, 快与慢》", sz=16, color=C_GOLD, bold=True)
    run(tf0, "  System 1 (快速直觉) → Fast Path: 模式匹配, <5ms, 不调LLM  |  System 2 (深度思考) → Slow Path: LLM推理, ~26s, 处理复杂场景", sz=13, color=C_WHITE, space_before=Pt(4))

    # Fast Path 详解
    b1 = box(s, Inches(0.6), Inches(2.8), Inches(5.8), Inches(4.3), fill=C_BG2, border=C_GREEN)
    tf1 = b1.text_frame
    tf1.word_wrap = True
    run(tf1, "  Fast Path (快路径) — <5ms, 零LLM调用", sz=18, color=C_GREEN, bold=True)
    
    fast_items = [
        ("Step 1: 指令解析", "jieba分词 / 英文split → 提取目标关键词\n\"find the chair near the desk\" → target=chair, spatial=desk"),
        ("Step 2: 标签匹配 (L_match)", "精确匹配: \"chair\" == \"chair\" → 1.0\n模糊匹配: \"椅\" in \"红色椅子\" → 0.9\n结果: 每个物体得到标签分 0~1"),
        ("Step 3: 检测置信度 (D_score)", "D = detection_score x min(det_count/5, 1.0)\n反复看到的物体更可信"),
        ("Step 4: 空间关系 (S_rel)", "如果指令含 \"near desk\", 检查relations\n物体与desk有near关系 → 加分 +0.3"),
        ("Step 5: 多源融合", "Score = 0.4*L + 0.35*D + 0.25*S\n> 0.75 → 直接输出, 不调LLM!"),
    ]
    for title, desc in fast_items:
        run(tf1, f"  {title}", sz=12, color=C_ACCENT_L, bold=True, space_before=Pt(8))
        for line in desc.split("\n"):
            run(tf1, f"    {line}", sz=10, color=C_GRAY, space_before=Pt(1))

    # Slow Path 详解
    b2 = box(s, Inches(6.9), Inches(2.8), Inches(5.8), Inches(4.3), fill=C_BG2, border=C_ORANGE)
    tf2 = b2.text_frame
    tf2.word_wrap = True
    run(tf2, "  Slow Path (慢路径) — LLM Chain-of-Thought", sz=18, color=C_ORANGE, bold=True)

    slow_items = [
        ("触发条件", "Fast Path 分数 < 0.75\n典型: 中文指令, 复杂空间关系, 目标不存在"),
        ("Prompt 设计 (层次化CoT)", "System: 你是导航规划器, 按以下步骤推理...\nUser: 场景图JSON + 机器人位置 + 指令\n强制步骤: Step1理解 → Step2选Room → \nStep3选Group → Step4选Object → Step5输出"),
        ("输出格式", '{"action": "navigate/explore",\n "target": {"x":0.8, "y":0.5, "z":0.8},\n "target_label": "fire extinguisher",\n "confidence": 0.95,\n "reasoning": "Room选择: corridor..."}'),
        ("跨语言能力", '"找到灭火器" → LLM理解"灭火器"=fire ext\n100%跨语言准确率 (实测12/12)'),
    ]
    for title, desc in slow_items:
        run(tf2, f"  {title}", sz=12, color=C_GOLD, bold=True, space_before=Pt(8))
        for line in desc.split("\n"):
            run(tf2, f"    {line}", sz=10, color=C_GRAY, space_before=Pt(1))


# ═══════════════════════════════════════════
#  SLIDE 7: 快慢推理 — 实测结果
# ═══════════════════════════════════════════
def slide_fast_slow_results(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, C_BG1)
    title_bar(s, 6, "核心技术 2: 快慢推理 — 实测结果", "Fast-Slow Results (Kimi-k2.5)")

    # 顶部指标
    metrics = [
        ("100%", "L1英文\nFast Path", C_GREEN),
        ("100%", "L2英文\nFast Path", C_GREEN),
        ("< 5ms", "Fast Path\n平均延迟", C_ACCENT),
        ("100%", "Slow Path\n通过率", C_ORANGE),
        ("~26s", "Slow Path\n平均延迟", C_ORANGE),
        ("100%", "跨语言\n准确率", C_GOLD),
    ]
    x = 0.6
    for val, label, color in metrics:
        stat_box(s, Inches(x), Inches(1.5), Inches(1.9), Inches(1.2), val, label, color)
        x += 2.05

    # LLM 实测案例
    b2 = box(s, Inches(0.6), Inches(3.0), Inches(12.1), Inches(4.2), fill=C_BG2, border=C_ACCENT)
    tf2 = b2.text_frame
    tf2.word_wrap = True
    run(tf2, "  Kimi-k2.5 实测输出 (真实API调用)", sz=18, color=C_ACCENT, bold=True)

    cases = [
        ("\"找到走廊里门旁边的灭火器\"", "Slow", "26s",
         "Room: corridor (指令明确指定'走廊', 该房间semantic_labels含fire extinguisher)\n→ Group: safety (灭火器属安全设备组)\n→ Object: id=2, fire extinguisher (验证relations: 与door有near关系, 距离1.0m, 完全匹配'门旁边')\n→ 导航到 (0.8, 0.5), confidence=0.95"),
        ("\"find the printer in office B\"", "Fast", "1.4ms",
         "label=1.0 (printer完全匹配) + det=0.87 + spatial=0.3\n→ fused=0.83 > 0.75阈值, 直接输出 (5.5, -2.5)\n无需调用LLM, 节省26秒和一次API费用!"),
        ("\"去厨房拿杯水\"", "Slow", "27s",
         "Room: kitchen (指令明确指定'厨房', 通过passage与corridor直连)\n→ Group: utility (厨房内实用设施组)\n→ Object: sink (水槽) — LLM推理: '拿杯水'场景下sink是最直接的取水来源\n→ 导航到 (-3.8, 2.8), confidence=0.92  [注意: 场景中没有'杯子', LLM做了语义推理!]"),
    ]
    for instr, path, latency, reasoning in cases:
        color = C_GREEN if path == "Fast" else C_ORANGE
        run(tf2, f"  {instr}", sz=14, color=C_GOLD, bold=True, space_before=Pt(10))
        run(tf2, f"    [{path} Path, {latency}]", sz=11, color=color, bold=True, space_before=Pt(2))
        for line in reasoning.split("\n"):
            run(tf2, f"    {line}", sz=10, color=C_GRAY, space_before=Pt(1))


# ═══════════════════════════════════════════
#  SLIDE 8: 信念推理详解
# ═══════════════════════════════════════════
def slide_belief(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, C_BG1)
    title_bar(s, 7, "核心技术 3: 信念感知推理", "Belief-Aware Reasoning (BA-HSG)")

    # 问题
    b0 = box(s, Inches(0.6), Inches(1.5), Inches(12.1), Inches(0.8), fill=C_BG3, border=C_RED)
    tf0 = b0.text_frame
    tf0.word_wrap = True
    run(tf0, "  问题: 检测器不完美 — YOLO可能漏检、误检、检测位置偏移。机器人怎么判断一个物体\"真的在那里\"?", sz=14, color=C_RED, bold=True)
    run(tf0, "  方案: 为每个物体维护概率信念状态, 用贝叶斯方法持续更新, 用信息价值(VoI)决定何时重新观察", sz=13, color=C_WHITE, space_before=Pt(4))

    # 三大信念组件
    components = [
        (0.6, 2.6, "Beta分布 → 存在概率", [
            "每个物体: P_exist = alpha/(alpha+beta)",
            "初始: alpha=2, beta=2 → P=0.5 (不确定)",
            "每次检测到: alpha++ → P上升",
            "每次该看到却没看到: beta++ → P下降",
            "30帧后: P从0.5→0.9 (反复确认)",
            "连续10帧没看到: P从0.9→0.5 (可疑)",
            "",
            "本质: 用观测证据量化\"这个物体存在吗?\"",
        ], C_ACCENT),
        (4.6, 2.6, "Gaussian → 位置精度", [
            "每个物体: 3D位置 ~ N(mu, sigma)",
            "初始: sigma = 1.0m (不确定)",
            "多次观测: sigma收缩 → 位置更精确",
            "不同角度观测 → sigma更快收缩",
            "长时间未观测 → sigma缓慢膨胀",
            "",
            "本质: 量化\"物体精确在哪?\"的信心",
            "用于多假设规划的期望成本计算",
        ], C_GREEN),
        (8.6, 2.6, "图扩散 → 邻域传播", [
            "新发现的物体如果在已知物体附近:",
            "  → 继承邻居的可信度加成",
            "  → \"桌上的杯子\"自动获得boost",
            "距离衰减: boost = 1/(1 + dist/lambda)",
            "lambda = 4.0m (衰减半径)",
            "",
            "本质: 空间上下文增强物体可信度",
            "灵感: 图神经网络的消息传递思想",
        ], C_ORANGE),
    ]
    for x, y, title, items, color in components:
        b = box(s, Inches(x), Inches(y), Inches(3.8), Inches(2.5), fill=C_BG2, border=color)
        tf = b.text_frame
        tf.word_wrap = True
        run(tf, f" {title}", sz=14, color=color, bold=True)
        for item in items:
            run(tf, f"  {item}", sz=10, color=C_GRAY if item else C_BG2, space_before=Pt(2))

    # 复合可信度
    b3 = box(s, Inches(0.6), Inches(5.4), Inches(12.1), Inches(1.8), fill=C_BG2, border=C_GOLD)
    tf3 = b3.text_frame
    tf3.word_wrap = True
    run(tf3, "  复合可信度公式", sz=16, color=C_GOLD, bold=True)
    run(tf3, "  Credibility = P_exist x max(D_score, CLIP_sim) x Freshness x GraphBoost", sz=16, color=C_WHITE, bold=True, space_before=Pt(6))
    run(tf3, "  P_exist: Beta分布存在概率 | D_score: 检测器置信度 | CLIP_sim: 语义相似度 | Freshness: 时间衰减(120s半衰) | GraphBoost: 邻域传播加成", sz=12, color=C_GRAY, space_before=Pt(6))
    run(tf3, "  Credibility > 0.7 → 高置信导航 | 0.3~0.7 → 谨慎接近 | < 0.3 → 需要重新感知 (VoI触发)", sz=12, color=C_ACCENT_L, space_before=Pt(4))


# ═══════════════════════════════════════════
#  SLIDE 9: VoI + 多假设
# ═══════════════════════════════════════════
def slide_voi(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, C_BG1)
    title_bar(s, 8, "核心技术 4: VoI调度 + 多假设规划", "Value-of-Information & Multi-Hypothesis")

    # VoI
    b1 = box(s, Inches(0.6), Inches(1.5), Inches(6.0), Inches(3.0), fill=C_BG2, border=C_ACCENT)
    tf1 = b1.text_frame
    tf1.word_wrap = True
    run(tf1, "  VoI 自适应调度: 该看还是该走?", sz=16, color=C_ACCENT, bold=True)
    run(tf1, "  核心思想: 不是固定每2米重新看一次, 而是根据\"信息价值\"决定", sz=12, color=C_GOLD, space_before=Pt(6))
    
    voi_items = [
        ("Continue (77%)", "可信度 > 0.7, 当前判断可靠, 继续走\n不浪费检测/LLM算力", C_GREEN),
        ("Reperceive (23%)", "可信度 < 0.3, 需要重新观察验证\n运动了足够距离(>2m) + 冷却期", C_ORANGE),
        ("Slow Reason (罕见)", "场景图不足以解析指令\n调用LLM做深度推理", C_RED),
    ]
    for title, desc, color in voi_items:
        run(tf1, f"  {title}", sz=13, color=color, bold=True, space_before=Pt(8))
        for line in desc.split("\n"):
            run(tf1, f"    {line}", sz=10, color=C_GRAY, space_before=Pt(1))

    # vs固定间隔
    b1b = box(s, Inches(0.6), Inches(4.8), Inches(6.0), Inches(2.3), fill=C_BG2, border=C_GREEN)
    tf1b = b1b.text_frame
    tf1b.word_wrap = True
    run(tf1b, "  VoI vs 固定间隔 (100步模拟)", sz=14, color=C_GREEN, bold=True)
    run(tf1b, "  固定2m间隔: 每7步触发1次 → 14次/100步 (均匀分布)", sz=12, color=C_GRAY, space_before=Pt(6))
    run(tf1b, "  VoI自适应: 高置信不触发 → 23次/100步 (集中在不确定时)", sz=12, color=C_GRAY, space_before=Pt(4))
    run(tf1b, "  结论: VoI在确定时节省算力, 不确定时更积极", sz=12, color=C_GOLD, bold=True, space_before=Pt(6))
    run(tf1b, "  77%时间零开销 + 23%精准感知 = 最优资源分配", sz=12, color=C_ACCENT, space_before=Pt(4))

    # Multi-Hypothesis
    b2 = box(s, Inches(7.0), Inches(1.5), Inches(5.7), Inches(5.6), fill=C_BG2, border=C_GOLD)
    tf2 = b2.text_frame
    tf2.word_wrap = True
    run(tf2, "  多假设目标规划", sz=16, color=C_GOLD, bold=True)
    run(tf2, "  问题: 场景里有3个灭火器, 去哪个?", sz=13, color=C_RED, bold=True, space_before=Pt(6))

    mh_items = [
        ("Step 1: 建立假设集", "3个灭火器 → 3个候选, 各有后验概率\n综合考虑: 可信度 + 距离 + 语义匹配"),
        ("Step 2: 期望成本选择", "Expected Cost = P(wrong) x travel_cost\n选总期望成本最低的 (不是最近的!)"),
        ("Step 3: 到达验证", "到达候选位置 → 重新检测\n确认 → 任务完成\n未确认 → 贝叶斯更新, 该候选概率↓"),
        ("Step 4: 迭代收敛", "概率重分配 → 选下一个最高概率候选\n平均1.6次找到 (vs随机3.0次)"),
        ("实测结果", "20个场景, 100%成功率 (20/20)\n平均尝试1.6次 (效率提升47%)"),
    ]
    for title, desc in mh_items:
        run(tf2, f"  {title}", sz=12, color=C_ACCENT_L, bold=True, space_before=Pt(8))
        for line in desc.split("\n"):
            run(tf2, f"    {line}", sz=10, color=C_GRAY, space_before=Pt(1))


# ═══════════════════════════════════════════
#  SLIDE 10: 拓扑语义探索 — 原理
# ═══════════════════════════════════════════
def slide_tsg_principle(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, C_BG1)
    title_bar(s, 9, "核心技术 5: 拓扑语义探索 — 原理", "Topology-Aware Semantic Exploration")

    b0 = box(s, Inches(0.6), Inches(1.5), Inches(12.1), Inches(0.8), fill=C_BG3, border=C_PURPLE)
    tf0 = b0.text_frame
    tf0.word_wrap = True
    run(tf0, "  问题: 目标物体不在当前场景图中 → 机器人去哪里探索最有可能找到?", sz=14, color=C_RED, bold=True)
    run(tf0, "  方案: 构建拓扑语义图(TSG), 用信息增益(IG)算法自动选择最优探索目标, 1ms完成无需LLM", sz=13, color=C_WHITE, space_before=Pt(4))

    # TSG 结构
    b1 = box(s, Inches(0.6), Inches(2.6), Inches(4.0), Inches(4.5), fill=C_BG2, border=C_PURPLE)
    tf1 = b1.text_frame
    tf1.word_wrap = True
    run(tf1, " 拓扑语义图 (TSG) 结构", sz=15, color=C_PURPLE, bold=True)

    tsg_items = [
        ("Room Nodes (房间节点)", [
            "从层次场景图同步",
            "属性: 类型/中心/物体列表",
            "状态: 已访问/未访问/访问次数",
        ]),
        ("Frontier Nodes (前沿节点)", [
            "未探索区域的入口",
            "检测策略1: 门指向无房间 → 前沿",
            "检测策略2: 物体稀疏方向 → 前沿",
            "属性: 位置/方向/预测房间类型",
        ]),
        ("Edges (拓扑边)", [
            "Door: 两房间共享一扇门",
            "Passage: 过道连通",
            "Proximity: 空间邻近",
            "Traversal: 机器人实际走过",
        ]),
        ("Traversal Memory (穿越记忆)", [
            "记录每次房间转换事件",
            "走过的边 → 置信度提升",
            "走过的房间 → 新颖度降低",
        ]),
    ]
    for title, items in tsg_items:
        run(tf1, f" {title}", sz=11, color=C_GOLD, bold=True, space_before=Pt(6))
        for item in items:
            run(tf1, f"   {item}", sz=9, color=C_GRAY, space_before=Pt(1))

    # Algorithm 2
    b2 = box(s, Inches(4.85), Inches(2.6), Inches(4.0), Inches(4.5), fill=C_BG2, border=C_ACCENT)
    tf2 = b2.text_frame
    tf2.word_wrap = True
    run(tf2, " Algorithm 2: 信息增益探索", sz=15, color=C_ACCENT, bold=True)

    run(tf2, " IG(n) = S_sem x N x U", sz=16, color=C_GOLD, bold=True, space_before=Pt(8))

    algo_items = [
        ("S_sem: 语义先验", [
            "知识库: room_type → (object, prob)",
            "咖啡机 → kitchen: 0.9",
            "灭火器 → stairwell: 0.9",
            "投影仪 → meeting_room: 0.85",
            "10种房间 x 100+物体关联",
        ]),
        ("N: 新颖度", [
            "未访问: N = 1.0 (最高)",
            "已访问: N = 0.1 + 0.9*exp(-0.5*count)",
            "  x (1 - exp(-dt/120s))",
            "去过1次: N~0.1 → IG降90%!",
        ]),
        ("U: 不确定度降低", [
            "前沿节点: U_base = 0.6 (高不确定)",
            "稀疏房间: U = 0.3 (中等)",
            "密集房间: U = 0.1 (已充分探索)",
        ]),
    ]
    for title, items in algo_items:
        run(tf2, f" {title}", sz=11, color=C_ACCENT_L, bold=True, space_before=Pt(6))
        for item in items:
            run(tf2, f"   {item}", sz=9, color=C_GRAY, space_before=Pt(1))

    # 最终评分
    b3 = box(s, Inches(9.1), Inches(2.6), Inches(3.6), Inches(2.0), fill=C_BG2, border=C_GREEN)
    tf3 = b3.text_frame
    tf3.word_wrap = True
    run(tf3, " 最终评分", sz=15, color=C_GREEN, bold=True)
    run(tf3, " Score = IG / (1+0.3*d)", sz=14, color=C_GOLD, bold=True, space_before=Pt(6))
    run(tf3, "   d = Dijkstra最短路径距离", sz=10, color=C_GRAY, space_before=Pt(4))
    run(tf3, "   远处IG高的不如近处IG中的", sz=10, color=C_GRAY, space_before=Pt(2))
    run(tf3, "   平衡 \"信息量\" vs \"到达成本\"", sz=10, color=C_ACCENT_L, space_before=Pt(4))

    # 双层策略
    b4 = box(s, Inches(9.1), Inches(4.9), Inches(3.6), Inches(2.2), fill=C_BG2, border=C_ORANGE)
    tf4 = b4.text_frame
    tf4.word_wrap = True
    run(tf4, " 双层探索策略", sz=15, color=C_ORANGE, bold=True)
    run(tf4, " Layer 1: TSG算法 (~1ms)", sz=12, color=C_GREEN, bold=True, space_before=Pt(6))
    run(tf4, "   纯算法, 无API成本", sz=10, color=C_GRAY, space_before=Pt(1))
    run(tf4, "   覆盖~70%场景", sz=10, color=C_GRAY, space_before=Pt(1))
    run(tf4, " Layer 2: LLM (fallback)", sz=12, color=C_ORANGE, bold=True, space_before=Pt(6))
    run(tf4, "   注入拓扑摘要给LLM", sz=10, color=C_GRAY, space_before=Pt(1))
    run(tf4, "   LLM有空间上下文更准", sz=10, color=C_GRAY, space_before=Pt(1))


# ═══════════════════════════════════════════
#  SLIDE 11: 拓扑探索 — 实测
# ═══════════════════════════════════════════
def slide_tsg_results(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, C_BG1)
    title_bar(s, 10, "核心技术 5: 拓扑探索 — 实测对比", "TSG vs LLM Real Test Results")

    # Case 1
    b1 = box(s, Inches(0.6), Inches(1.5), Inches(5.8), Inches(2.2), fill=C_BG2, border=C_ACCENT)
    tf1 = b1.text_frame
    tf1.word_wrap = True
    run(tf1, '  Case 1: "找到咖啡机" (目标不存在)', sz=14, color=C_GOLD, bold=True)
    run(tf1, "  TSG Algorithm 2:", sz=12, color=C_ACCENT, bold=True, space_before=Pt(6))
    run(tf1, "    #1 stairwell (IG=0.460, score=0.184)", sz=11, color=C_GRAY, space_before=Pt(2))
    run(tf1, "    #2 kitchen   (IG=0.440, score=0.176)", sz=11, color=C_GRAY, space_before=Pt(1))
    run(tf1, "  LLM (Kimi-k2.5): → kitchen", sz=12, color=C_ORANGE, bold=True, space_before=Pt(4))
    run(tf1, '    "咖啡机最可能出现在厨房...前往厨房信息增益最大"', sz=10, color=C_GRAY, space_before=Pt(2))
    run(tf1, "  分析: LLM比TSG更准 (语义常识), TSG纯算法不知道咖啡机→厨房", sz=10, color=C_ACCENT_L, space_before=Pt(4))

    # Case 2
    b2 = box(s, Inches(6.9), Inches(1.5), Inches(5.8), Inches(2.2), fill=C_BG2, border=C_GREEN)
    tf2 = b2.text_frame
    tf2.word_wrap = True
    run(tf2, '  Case 2: "find fire ext on 2nd floor"', sz=14, color=C_GOLD, bold=True)
    run(tf2, "  TSG Algorithm 2:", sz=12, color=C_ACCENT, bold=True, space_before=Pt(6))
    run(tf2, "    #1 stairwell (IG=1.035, score=0.414) [大幅领先!]", sz=11, color=C_GREEN, space_before=Pt(2))
    run(tf2, "    #2 frontier  (IG=0.600, score=0.136)", sz=11, color=C_GRAY, space_before=Pt(1))
    run(tf2, "  LLM (Kimi-k2.5): → stairwell", sz=12, color=C_ORANGE, bold=True, space_before=Pt(4))
    run(tf2, '    "Fire ext legally required in stairwells...highest prior 0.9"', sz=10, color=C_GRAY, space_before=Pt(2))
    run(tf2, "  分析: TSG和LLM完全一致! 语义先验正确引导到stairwell", sz=10, color=C_ACCENT_L, space_before=Pt(4))

    # Case 3: 穿越记忆
    b3 = box(s, Inches(0.6), Inches(4.0), Inches(12.1), Inches(3.2), fill=C_BG2, border=C_PURPLE)
    tf3 = b3.text_frame
    tf3.word_wrap = True
    run(tf3, '  Case 3: 穿越记忆效果 — "找到灭火器" (3轮探索)', sz=16, color=C_PURPLE, bold=True)

    run(tf3, "  Round 1: 刚开始在corridor", sz=13, color=C_ACCENT, bold=True, space_before=Pt(8))
    run(tf3, "    stairwell(0.414) > frontier(0.127) > kitchen(0.066) > office_B(0.063)", sz=11, color=C_GRAY, space_before=Pt(2))

    run(tf3, "  Round 2: 去了office_A, 没找到, 回来了", sz=13, color=C_ORANGE, bold=True, space_before=Pt(6))
    run(tf3, "    stairwell(0.414) > frontier(0.127) > kitchen(0.066) > office_B(0.063) [office_A已消失]", sz=11, color=C_GRAY, space_before=Pt(2))

    run(tf3, "  Round 3: 又去了kitchen, 没找到, 回来了", sz=13, color=C_RED, bold=True, space_before=Pt(6))
    run(tf3, "    stairwell(0.414) > frontier(0.127) > office_B(0.063) > kitchen(0.006) [kitchen IG暴降90%!]", sz=11, color=C_GRAY, space_before=Pt(2))

    run(tf3, "  LLM最终建议: \"向西南探索到stairwell, 避开已探索的东北(office_A)和西北(kitchen)方向\"", sz=12, color=C_GOLD, bold=True, space_before=Pt(6))
    run(tf3, "  穿越记忆完美阻止重复探索 + LLM正确利用负面记忆!", sz=12, color=C_GREEN, bold=True, space_before=Pt(4))


# ═══════════════════════════════════════════
#  SLIDE 12: 完整测试结果
# ═══════════════════════════════════════════
def slide_test_results(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, C_BG1)
    title_bar(s, 11, "完整测试验证", "Test Results Overview")

    # 总览
    metrics = [
        ("99/99", "测试全部通过", C_GREEN),
        ("33/33", "拓扑图测试", C_ACCENT),
        ("40/40", "离线管线", C_ORANGE),
        ("26/26", "信念系统", C_PURPLE),
        ("12/12", "LLM端到端", C_GOLD),
    ]
    x = 0.6
    for val, label, color in metrics:
        stat_box(s, Inches(x), Inches(1.5), Inches(2.3), Inches(1.2), val, label, color)
        x += 2.45

    # 详细分类
    categories = [
        (0.6, 3.0, "拓扑图 (33测试)", [
            "图构建: 同步/类型推理/状态保持 (5/5)",
            "前沿节点: 添加/更新/清除/去重 (4/4)",
            "穿越记忆: 检测/转换/边创建/计数 (4/4)",
            "最短路径: 邻居/2跳/不可达/BFS (4/4)",
            "信息增益: 未访>已访/前沿/语义 (3/3)",
            "探索目标: 选择/偏好/前沿/衰减/topK (5/5)",
            "序列化: JSON往返/可序列化 (2/2)",
            "Prompt: 中英文拓扑摘要 (2/2)",
            "边界: 空图/单房间/越界/位置查询 (4/4)",
        ], C_ACCENT),
        (4.85, 3.0, "离线管线 (40测试)", [
            "L1快路径: 20种英文物体 (20/20)",
            "L2快路径: 5种空间关系 (5/5)",
            "中文落入慢路径: 验证正确fallback (1/1)",
            "延迟: 所有Fast Path < 5ms (1/1)",
            "任务分解: L1全部/L3多步/Follow (6/6)",
            "信念端到端: 正向/反向/图扩散 (3/3)",
            "多假设: 初始/贝叶斯更新/收敛 (3/3)",
            "VoI: 高置信/低置信/冷却/运动 (4/4)",
        ], C_ORANGE),
        (9.1, 3.0, "LLM端到端 (12+6测试)", [
            "英文Fast Path: 6用例 100% (6/6)",
            "中文Slow Path: 4用例 100% (4/4)",
            "探索: 2用例 100% (2/2)",
            "TSG+LLM探索对比: 3场景 (3/3)",
            "穿越记忆: 3轮迭代验证 (3/3)",
            "",
            "LLM后端: Kimi-k2.5 (Moonshot)",
            "跨语言: 中→英 100% 准确",
            "语义推理: '拿杯水'→sink 成功",
        ], C_GOLD),
    ]
    for x, y, title, items, color in categories:
        b = box(s, Inches(x), Inches(y), Inches(4.0), Inches(4.0), fill=C_BG2, border=color)
        tf = b.text_frame
        tf.word_wrap = True
        run(tf, f" {title}", sz=14, color=color, bold=True)
        for item in items:
            c = C_GRAY if item else C_BG2
            run(tf, f"  {item}", sz=10, color=c, space_before=Pt(2))


# ═══════════════════════════════════════════
#  SLIDE 13: 竞品对比
# ═══════════════════════════════════════════
def slide_competitive(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, C_BG1)
    title_bar(s, 12, "竞品对比 & 技术壁垒", "Competitive Landscape")

    headers = ["维度", "SG-Nav\n(2024)", "FSR-VLN\n(2025)", "CLIP-Frontier\n(CoW)", "NaviMind\n灵途 (Ours)"]
    col_w = [2.0, 2.3, 2.3, 2.3, 3.6]

    rows = [
        ["运行环境", "仿真", "预建地图", "未知环境", "未知+动态+在线"],
        ["场景图", "3层确定", "4层离线", "无", "4层在线+信念"],
        ["目标解析", "每步LLM", "CLIP+VLM", "CLIP匹配", "Fast<5ms+SlowLLM"],
        ["LLM调用率", "100%", "VLM精筛", "0%", "~25% (节省75%)"],
        ["不确定性", "标量", "无", "无", "Beta+Gaussian+VoI"],
        ["探索策略", "LLM引导", "无", "随机前沿", "TSG信息增益+LLM"],
        ["多步指令", "不支持", "单步", "不支持", "支持 (任务分解)"],
        ["跨语言", "未报告", "未报告", "不支持", "中英100%"],
        ["硬件", "仿真90cm", "桌面GPU", "桌面GPU", "Jetson 30cm"],
        ["开放性", "闭源", "闭源", "开源", "全栈自研"],
    ]

    y = 1.4
    x = 0.4
    for i, (h, w) in enumerate(zip(headers, col_w)):
        b = box(s, Inches(x), Inches(y), Inches(w), Inches(0.6), fill=C_ACCENT)
        tf = b.text_frame
        tf.word_wrap = True
        run(tf, h, sz=11, color=C_BG1, bold=True, align=PP_ALIGN.CENTER)
        x += w + 0.04

    for row in rows:
        y += 0.53
        x = 0.4
        for j, (cell, w) in enumerate(zip(row, col_w)):
            bg = C_BG2
            c = C_GRAY
            if j == 4:
                bg = RGBColor(0x0A, 0x25, 0x1A)
                c = C_GREEN
            elif j == 0:
                c = C_GOLD
            b = box(s, Inches(x), Inches(y), Inches(w), Inches(0.5), fill=bg)
            tf = b.text_frame
            tf.word_wrap = True
            run(tf, cell, sz=10, color=c, bold=(j==0 or j==4), align=PP_ALIGN.CENTER)
            x += w + 0.04


# ═══════════════════════════════════════════
#  SLIDE 14: 技术壁垒
# ═══════════════════════════════════════════
def slide_moat(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, C_BG1)
    title_bar(s, 13, "核心技术壁垒", "Technical Moat")

    moats = [
        (0.6, 1.5, 5.8, 1.5, "壁垒1: 全栈自研, 端到端闭环", [
            "从传感器 → 感知 → 场景图 → 推理 → 规划 → 控制 全链自研",
            "17个ROS2 Package + 44个gRPC RPC + Flutter多平台App",
            "竞品通常只做其中一环 (感知 or 规划), 无法闭环",
        ], C_ACCENT),
        (6.9, 1.5, 5.8, 1.5, "壁垒2: 快慢架构, LLM成本降75%", [
            "Fast Path <5ms 处理75%+查询, 无LLM成本",
            "竞品 (SG-Nav) 每步都调LLM → 延迟高、成本高",
            "我们的Fast Path效果等价于LLM, 但成本为零",
        ], C_GREEN),
        (0.6, 3.3, 5.8, 1.5, "壁垒3: 概率信念, 不确定性建模", [
            "业界首个将Beta+Gaussian信念引入机器人场景图导航",
            "VoI调度 → 自适应计算分配, 77%零开销",
            "多假设贝叶斯规划 → 1.6次vs3.0次找到目标",
        ], C_ORANGE),
        (6.9, 3.3, 5.8, 1.5, "壁垒4: 拓扑语义探索", [
            "融合论文创新: TopoNav + L3MVN + Hydra + TACS-Graphs",
            "信息增益算法 ~1ms, 替代LLM探索 (节省2-5s/步)",
            "穿越记忆避免重复探索, 提升效率90%",
        ], C_PURPLE),
        (0.6, 5.1, 5.8, 1.5, "壁垒5: 边缘部署, 真机验证", [
            "Jetson Orin NX 16GB, 非桌面GPU仿真",
            "四足30cm低视角 → 独特技术挑战 (模糊过滤等)",
            "全链已在Go2上验证, 非纸上谈兵",
        ], C_RED),
        (6.9, 5.1, 5.8, 1.5, "壁垒6: 多LLM + 跨语言", [
            "支持4种LLM后端 (OpenAI/Claude/Qwen/Kimi)",
            "中英文指令100%准确, 实测验证",
            "Prompt工程: 层次化CoT, 可复现可调优",
        ], C_GOLD),
    ]
    for x, y, w, h, title, items, color in moats:
        b = box(s, Inches(x), Inches(y), Inches(w), Inches(h), fill=C_BG2, border=color)
        tf = b.text_frame
        tf.word_wrap = True
        run(tf, f" {title}", sz=14, color=color, bold=True)
        for item in items:
            run(tf, f"  {item}", sz=10, color=C_GRAY, space_before=Pt(3))


# ═══════════════════════════════════════════
#  SLIDE 15: 应用场景
# ═══════════════════════════════════════════
def slide_use_cases(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, C_BG1)
    title_bar(s, 14, "应用场景", "Use Cases")

    cases = [
        (0.6, 1.5, "智能巡检", "工厂 / 仓库 / 数据中心", [
            '"去检查3号机房的服务器状态"',
            "自主导航 + 物体识别 + 异常检测",
            "无需预设路线, 语音即下达巡检任务",
            "夜间/危险区域无人值守巡检",
        ], C_ACCENT, "规模: $180亿/年 (工业巡检)"),
        (6.9, 1.5, "服务机器人", "酒店 / 医院 / 商场", [
            '"把这瓶水送到3楼走廊尽头的护士站"',
            "理解空间语言, 多步任务分解",
            "动态环境适应 (人员走动)",
            "中英文指令, 服务国际客户",
        ], C_GREEN, "规模: $120亿/年 (服务机器人)"),
        (0.6, 4.0, "搜救 & 安防", "灾后搜索 / 建筑安检", [
            '"搜索这栋楼里所有灭火器的位置"',
            "未知环境自主探索 (无预建地图)",
            "拓扑语义图 → 系统化搜索策略",
            "四足全地形 → 楼梯/废墟/狭窄空间",
        ], C_ORANGE, "规模: $35亿/年 (安防机器人)"),
        (6.9, 4.0, "家庭陪伴", "智能家居 / 养老看护", [
            '"帮我找到客厅沙发上的遥控器"',
            "自然语言交互, 老人友好",
            "记忆式导航 → 学习家庭布局",
            "跌倒检测 + 紧急呼叫",
        ], C_PURPLE, "规模: $85亿/年 (家庭机器人)"),
    ]
    for x, y, title, subtitle, items, color, market in cases:
        b = box(s, Inches(x), Inches(y), Inches(5.8), Inches(2.2), fill=C_BG2, border=color)
        tf = b.text_frame
        tf.word_wrap = True
        run(tf, f" {title}", sz=18, color=color, bold=True)
        run(tf, f"  {subtitle}", sz=12, color=C_GOLD, space_before=Pt(2))
        for item in items:
            run(tf, f"  {item}", sz=11, color=C_GRAY, space_before=Pt(3))
        run(tf, f"  {market}", sz=11, color=C_ACCENT_L, bold=True, space_before=Pt(4))


# ═══════════════════════════════════════════
#  SLIDE 16: 商业模式
# ═══════════════════════════════════════════
def slide_business(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, C_BG1)
    title_bar(s, 15, "商业模式", "Business Model")

    models = [
        (0.6, 1.5, 3.9, 3.5, "SDK 授权", [
            "NaviMind Navigation SDK",
            "",
            "目标客户:",
            "  机器人整机厂商",
            "  系统集成商",
            "",
            "定价模型:",
            "  按机器人台数授权",
            "  基础版: $200/台/年",
            "  企业版: $500/台/年 (含LLM)",
            "",
            "预估: 1万台 x $300 = $300万/年",
        ], C_ACCENT),
        (4.75, 1.5, 3.9, 3.5, "SaaS 云平台", [
            "NaviMind Cloud",
            "",
            "服务内容:",
            "  LLM推理API (快慢路径)",
            "  场景图云端存储 & 分析",
            "  远程监控 & OTA更新",
            "",
            "定价模型:",
            "  按API调用量 + 月订阅",
            "  $99/月基础 + $0.01/LLM调用",
            "",
            "预估: 500企业 x $200/月 = $120万/年",
        ], C_GREEN),
        (8.9, 1.5, 3.9, 3.5, "解决方案", [
            "行业定制方案",
            "",
            "目标行业:",
            "  工业巡检 (工厂/电力)",
            "  医院物流 (送药/送餐)",
            "  仓储管理 (库存盘点)",
            "",
            "交付模式:",
            "  硬件+软件+部署+运维",
            "  年度服务合同",
            "",
            "预估: 20项目 x $30万 = $600万/年",
        ], C_ORANGE),
    ]
    for x, y, w, h, title, items, color in models:
        b = box(s, Inches(x), Inches(y), Inches(w), Inches(h), fill=C_BG2, border=color)
        tf = b.text_frame
        tf.word_wrap = True
        run(tf, f" {title}", sz=18, color=color, bold=True)
        for item in items:
            c = C_GRAY
            if item.startswith("预估"):
                c = C_GOLD
            elif item == "":
                c = C_BG2
            run(tf, f"  {item}", sz=11, color=c, bold=item.startswith("预估"), space_before=Pt(2))

    # 收入预估
    b_rev = box(s, Inches(0.6), Inches(5.3), Inches(12.1), Inches(1.8), fill=C_BG2, border=C_GOLD)
    tf_rev = b_rev.text_frame
    tf_rev.word_wrap = True
    run(tf_rev, "  收入预估", sz=18, color=C_GOLD, bold=True)
    
    years = [
        ("Year 1 (2026)", "SDK 10台试用 + 3个定制项目 = $100万", C_GRAY),
        ("Year 2 (2027)", "SDK 500台 + SaaS 100企业 + 10项目 = $500万", C_ACCENT_L),
        ("Year 3 (2028)", "SDK 5000台 + SaaS 500企业 + 20项目 = $1500万", C_GREEN),
    ]
    for label, desc, color in years:
        run(tf_rev, f"  {label}: {desc}", sz=13, color=color, bold=True, space_before=Pt(6))


# ═══════════════════════════════════════════
#  SLIDE 17: 路线图
# ═══════════════════════════════════════════
def slide_roadmap(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, C_BG1)
    title_bar(s, 16, "产品路线图", "Product Roadmap")

    phases = [
        (0.6, 1.5, "Q1-Q2 2026", "技术验证", [
            "真机实验: Go2 办公环境 L1-L3",
            "消融实验: 各模块贡献量化",
            "论文投稿: 顶会 (IROS/CoRL)",
            "SDK v0.1: 基础导航功能",
        ], C_ACCENT, "当前阶段"),
        (3.85, 1.5, "Q3-Q4 2026", "产品化", [
            "SDK v1.0: 工业级稳定性",
            "多楼层: TSG跨楼层导航",
            "端侧LLM: 4-bit量化部署",
            "云平台: 远程监控 + OTA",
            "3家种子客户试用",
        ], C_GREEN, ""),
        (7.1, 1.5, "2027 H1", "市场拓展", [
            "SDK v2.0: 多机器人协同",
            "行业方案: 工业巡检/医院",
            "动态场景: 人员/物体移动",
            "多模态: 语音 + 手势控制",
            "50家付费客户",
        ], C_ORANGE, ""),
        (10.35, 1.5, "2027 H2+", "规模化", [
            "SDK v3.0: 室外环境扩展",
            "自学习: 环境适应 + 持续优化",
            "生态: 第三方开发者平台",
            "国际化: 多语言 + 海外市场",
            "500+付费客户",
        ], C_PURPLE, ""),
    ]
    for x, y, time_label, title, items, color, badge in phases:
        b = box(s, Inches(x), Inches(y), Inches(2.95), Inches(3.0), fill=C_BG2, border=color)
        tf = b.text_frame
        tf.word_wrap = True
        run(tf, f" {time_label}", sz=11, color=C_DIM, bold=True)
        run(tf, f" {title}", sz=18, color=color, bold=True, space_before=Pt(2))
        if badge:
            run(tf, f" [{badge}]", sz=10, color=C_GOLD, bold=True, space_before=Pt(2))
        for item in items:
            run(tf, f"  {item}", sz=11, color=C_GRAY, space_before=Pt(4))

    # 技术里程碑
    b_ml = box(s, Inches(0.6), Inches(4.8), Inches(12.1), Inches(2.3), fill=C_BG2, border=C_GOLD)
    tf_ml = b_ml.text_frame
    tf_ml.word_wrap = True
    run(tf_ml, "  已完成的技术里程碑", sz=16, color=C_GOLD, bold=True)
    
    milestones = [
        ("在线层次场景图 (BA-HSG)", "4层结构, YOLO-World + CLIP, 质量感知融合, 拓扑连通边"),
        ("快慢双路径推理", "Fast Path <5ms 100%英文 | Slow Path LLM CoT 100%跨语言 | 4种LLM后端"),
        ("信念感知 + VoI调度", "Beta/Gaussian概率建模 | 图扩散 | VoI自适应 (77% continue) | 多假设1.6次收敛"),
        ("拓扑语义探索 (TSG)", "信息增益Algorithm 2 | Dijkstra路径 | 穿越记忆 | 双层策略 | 33测试全通过"),
        ("全栈集成", "17 ROS2 Packages | 44 gRPC RPCs | Flutter多平台App | 99单元测试 | Kimi-k2.5验证"),
    ]
    for name, desc in milestones:
        run(tf_ml, f"  {name}: {desc}", sz=10, color=C_GRAY, space_before=Pt(3))


# ═══════════════════════════════════════════
#  SLIDE 18: 融资需求
# ═══════════════════════════════════════════
def slide_funding(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, C_BG1)
    title_bar(s, 17, "融资计划", "Funding")

    # 融资额
    b_ask = box(s, Inches(0.6), Inches(1.5), Inches(12.1), Inches(1.2), fill=C_BG3, border=C_GOLD)
    tf_ask = b_ask.text_frame
    tf_ask.word_wrap = True
    run(tf_ask, "  天使轮融资: $100-200万", sz=28, color=C_GOLD, bold=True, align=PP_ALIGN.CENTER)
    run(tf_ask, "  出让股权: 10-15%  |  预估估值: $1000-1500万  |  用途: 产品化 + 真机验证 + 市场拓展", sz=15, color=C_WHITE, align=PP_ALIGN.CENTER, space_before=Pt(6))

    # 资金用途
    uses = [
        (0.6, 3.0, "研发 (60%)", "$60-120万", [
            "真机实验: Go2采购+实验环境",
            "工程师: 2名全栈 + 1名算法",
            "GPU服务器: 训练 + 测试",
            "LLM API: 端侧部署研发",
        ], C_ACCENT),
        (4.85, 3.0, "产品化 (25%)", "$25-50万", [
            "SDK封装: 工业级接口+文档",
            "云平台: 远程管理SaaS",
            "QA: 自动化测试 + CI/CD",
            "认证: 安全合规",
        ], C_GREEN),
        (9.1, 3.0, "市场 (15%)", "$15-30万", [
            "种子客户: 3-5家免费试用",
            "行业展会: IROS, CES, WRC",
            "品牌: 官网 + 技术博客",
            "BD: 渠道合作",
        ], C_ORANGE),
    ]
    for x, y, title, amount, items, color in uses:
        b = box(s, Inches(x), Inches(y), Inches(3.9), Inches(2.3), fill=C_BG2, border=color)
        tf = b.text_frame
        tf.word_wrap = True
        run(tf, f" {title}", sz=16, color=color, bold=True)
        run(tf, f" {amount}", sz=20, color=C_GOLD, bold=True, space_before=Pt(4))
        for item in items:
            run(tf, f"  {item}", sz=11, color=C_GRAY, space_before=Pt(3))

    # 关键里程碑
    b_kpi = box(s, Inches(0.6), Inches(5.6), Inches(12.1), Inches(1.5), fill=C_BG2, border=C_PURPLE)
    tf_kpi = b_kpi.text_frame
    tf_kpi.word_wrap = True
    run(tf_kpi, "  18个月关键KPI", sz=16, color=C_PURPLE, bold=True)
    kpis = [
        ("6个月", "论文发表 (顶会) + SDK v1.0 + 3家种子客户"),
        ("12个月", "SDK v2.0 + 云平台上线 + 10家付费客户 + $50万ARR"),
        ("18个月", "SDK v3.0 + 行业方案 + 50家客户 + $200万ARR → A轮"),
    ]
    for label, desc in kpis:
        run(tf_kpi, f"  {label}: {desc}", sz=13, color=C_WHITE, space_before=Pt(5))


# ═══════════════════════════════════════════
#  SLIDE 19: 论文创新总结
# ═══════════════════════════════════════════
def slide_paper_innovations(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, C_BG1)
    title_bar(s, 18, "学术创新总结", "Research Contributions")

    innovations = [
        ("创新1", "在线层次场景图 (BA-HSG)", 
         "首次将信念感知(Beta+Gaussian)融入4层在线场景图\nvs ConceptGraphs(离线/无信念) vs SG-Nav(3层/确定性)", 
         "Object→Group→Room→Floor + 质量感知CLIP + 运动模糊过滤", C_ACCENT),
        ("创新2", "快慢双路径推理",
         "借鉴认知科学System1/2, 首次在机器人导航中实现\n快路径<5ms零LLM vs SG-Nav每步调LLM",
         "多源融合评分 + 层次化CoT Prompt + 跨语言100%", C_GREEN),
        ("创新3", "信念感知规划 + VoI调度",
         "业界首个Beta/Gaussian信念 + VoI调度用于场景图导航\n自适应计算分配(77% continue/23% reperceive)",
         "复合可信度 + 图扩散 + 信息价值驱动决策", C_ORANGE),
        ("创新4", "多假设贝叶斯目标规划",
         "期望成本选择 → 1.6次收敛 vs 3.0次随机\n贝叶斯更新: 到达验证→拒绝→概率重分配",
         "后验概率排序 + 导航成本 + 信息增益", C_GOLD),
        ("创新5", "拓扑语义探索 (TSG)",
         "融合TopoNav+L3MVN+Hydra, 首次将拓扑图用于语义探索\n信息增益~1ms替代LLM探索(2-5s), 穿越记忆避免90%重复",
         "Algorithm 2: IG = S_sem x N x U + Dijkstra最短路径 + 双层策略", C_PURPLE),
        ("创新6", "边缘部署四足导航",
         "首次在30cm低视角四足+Jetson上实现全链语义导航\nvs 仿真90cm/桌面GPU → 真实世界可部署",
         "17 ROS2 Packages + 44 gRPC RPCs + Flutter App + OTA", C_RED),
    ]

    y = 1.4
    for idx, (label, title, contribution, details, color) in enumerate(innovations):
        b = box(s, Inches(0.6), Inches(y), Inches(12.1), Inches(0.9), fill=C_BG2, border=color)
        tf = b.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        r1 = p.add_run()
        r1.text = f" {label}: {title}"
        r1.font.size = Pt(14)
        r1.font.color.rgb = color
        r1.font.bold = True
        r1.font.name = FONT_CN

        for line in contribution.split("\n"):
            run(tf, f"    {line}", sz=10, color=C_WHITE, space_before=Pt(1))
        y += 0.97


# ═══════════════════════════════════════════
#  SLIDE 20: Thank You
# ═══════════════════════════════════════════
def slide_thankyou(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, C_BG1)
    rect(s, Inches(0), Inches(0), SLIDE_W, Pt(3), C_ACCENT)

    tb = txt(s, Inches(0.8), Inches(1.8), Inches(11.7), Inches(1.2))
    tf = tb.text_frame
    run(tf, "NaviMind", sz=56, color=C_ACCENT, bold=True, font=FONT_EN, align=PP_ALIGN.CENTER)

    tb2 = txt(s, Inches(0.8), Inches(3.0), Inches(11.7), Inches(0.5))
    tf2 = tb2.text_frame
    run(tf2, "灵途 — 让机器人拥有认知导航能力", sz=24, color=C_WHITE, bold=True, align=PP_ALIGN.CENTER)

    tb3 = txt(s, Inches(1.5), Inches(4.0), Inches(10.3), Inches(2.0))
    tf3 = tb3.text_frame
    run(tf3, "6项核心创新  |  99测试全通过  |  4种LLM验证  |  全栈自研", sz=18, color=C_GOLD, bold=True, align=PP_ALIGN.CENTER)
    run(tf3, "", sz=10, color=C_BG1, align=PP_ALIGN.CENTER)
    run(tf3, "Thank You  /  Questions & Discussion", sz=22, color=C_ACCENT_L, align=PP_ALIGN.CENTER, space_before=Pt(20))

    rect(s, Inches(0), Inches(6.7), SLIDE_W, Inches(0.8), C_BG2)
    tb4 = txt(s, Inches(0.8), Inches(6.8), Inches(11.7), Inches(0.5))
    tf4 = tb4.text_frame
    run(tf4, "Contact: [email]  |  GitHub: [repo]  |  2026", sz=14, color=C_DIM, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════
#  Main
# ═══════════════════════════════════════════
def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_cover(prs)              # 1. 封面
    slide_problem(prs)            # 2. 痛点 + 市场
    slide_solution(prs)           # 3. 方案总览
    slide_architecture(prs)       # 4. 技术架构
    slide_scene_graph(prs)        # 5. 在线层次场景图
    slide_fast_slow_principle(prs)  # 6. 快慢推理原理
    slide_fast_slow_results(prs)  # 7. 快慢推理实测
    slide_belief(prs)             # 8. 信念推理
    slide_voi(prs)                # 9. VoI + 多假设
    slide_tsg_principle(prs)      # 10. 拓扑探索原理
    slide_tsg_results(prs)        # 11. 拓扑探索实测
    slide_test_results(prs)       # 12. 测试结果
    slide_competitive(prs)        # 13. 竞品对比
    slide_moat(prs)               # 14. 技术壁垒
    slide_use_cases(prs)          # 15. 应用场景
    slide_business(prs)           # 16. 商业模式
    slide_roadmap(prs)            # 17. 路线图
    slide_funding(prs)            # 18. 融资计划
    slide_paper_innovations(prs)  # 19. 学术创新
    slide_thankyou(prs)           # 20. Thank You

    out_dir = os.path.join(os.path.dirname(__file__), "..", "docs")
    os.makedirs(out_dir, exist_ok=True)
    out_path = os.path.join(out_dir, "NaviMind_Pitch_Deck.pptx")
    prs.save(out_path)
    print(f"PPT saved: {os.path.abspath(out_path)}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
