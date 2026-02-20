"""
HSG-Nav 项目介绍PPT生成器。
生成一个完整的项目答辩/展示PPT, 包括:
  - 封面
  - 项目背景与动机
  - 系统架构
  - 核心创新 (6大创新点)
  - 实验验证
  - 与现有工作对比
  - 总结与展望
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── 颜色主题 ──
BG_DARK    = RGBColor(0x1A, 0x1A, 0x2E)   # 深蓝黑背景
BG_CARD    = RGBColor(0x16, 0x21, 0x3E)   # 卡片深蓝
ACCENT     = RGBColor(0x00, 0xD2, 0xFF)   # 科技蓝
ACCENT2    = RGBColor(0x0F, 0xCB, 0x8A)   # 绿色强调
ACCENT3    = RGBColor(0xFF, 0x6B, 0x6B)   # 红色强调
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xBB, 0xBB, 0xCC)
GOLD       = RGBColor(0xFF, 0xD7, 0x00)
ORANGE     = RGBColor(0xFF, 0x9F, 0x43)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_box(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape


def set_text(tf, text, size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT):
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = "Microsoft YaHei"
    return p


def add_paragraph(tf, text, size=16, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, space_before=Pt(6)):
    p = tf.add_paragraph()
    p.alignment = alignment
    p.space_before = space_before
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = "Microsoft YaHei"
    return p


def make_title_slide(prs):
    """Slide 1: 封面"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, BG_DARK)

    # 顶部装饰线
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Pt(4))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()

    # 标题
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(11.33), Inches(1.5))
    tf = txBox.text_frame
    set_text(tf, "HSG-Nav", size=54, color=ACCENT, bold=True, alignment=PP_ALIGN.CENTER)
    add_paragraph(tf, "基于层次场景图的四足机器人自主导航系统", size=28, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_paragraph(tf, "Hierarchical Scene Graph Navigation for Quadruped Robots", size=20, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    # 副标题
    txBox2 = slide.shapes.add_textbox(Inches(2), Inches(4.5), Inches(9.33), Inches(2))
    tf2 = txBox2.text_frame
    set_text(tf2, "快慢双路径推理  |  信念感知场景图  |  拓扑语义探索", size=20, color=GOLD, alignment=PP_ALIGN.CENTER)
    add_paragraph(tf2, "Fast-Slow Dual-Path Reasoning  |  Belief-Aware Scene Graph  |  Topology-Aware Exploration", size=16, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER, space_before=Pt(12))

    # 底部信息
    txBox3 = slide.shapes.add_textbox(Inches(2), Inches(5.8), Inches(9.33), Inches(1))
    tf3 = txBox3.text_frame
    set_text(tf3, "平台: Unitree Go2 + Jetson Orin NX  |  2026", size=16, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)


def make_background_slide(prs):
    """Slide 2: 研究背景与动机"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    # 标题
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    set_text(txBox.text_frame, "01  研究背景与动机", size=32, color=ACCENT, bold=True)

    # 问题描述
    box1 = add_shape_box(slide, Inches(0.5), Inches(1.3), Inches(6), Inches(2.5), BG_CARD, ACCENT)
    tf1 = box1.text_frame
    tf1.word_wrap = True
    set_text(tf1, "  核心问题", size=22, color=ACCENT, bold=True)
    add_paragraph(tf1, "  如何让四足机器人在未知室内环境中,", size=16, color=WHITE)
    add_paragraph(tf1, "  根据自然语言指令自主导航到目标?", size=16, color=WHITE)
    add_paragraph(tf1, "", size=10, color=WHITE)
    add_paragraph(tf1, '  "找到走廊里门旁边的灭火器"', size=18, color=GOLD, bold=True)
    add_paragraph(tf1, '  "去厨房拿杯水"', size=18, color=GOLD, bold=True)

    # 挑战
    box2 = add_shape_box(slide, Inches(6.83), Inches(1.3), Inches(6), Inches(2.5), BG_CARD, ACCENT3)
    tf2 = box2.text_frame
    tf2.word_wrap = True
    set_text(tf2, "  关键挑战", size=22, color=ACCENT3, bold=True)
    challenges = [
        "  1. 环境未知 — 无先验地图, 需在线建图",
        "  2. 语言多样 — 中/英文, 空间关系, 多步指令",
        "  3. 低视角 — 四足机器人仅30cm高, 视野受限",
        "  4. 边缘计算 — Jetson Orin NX, 无云端GPU",
        "  5. 动态环境 — 物体移动、人员走动",
    ]
    for c in challenges:
        add_paragraph(tf2, c, size=14, color=WHITE)

    # 现有方法局限
    box3 = add_shape_box(slide, Inches(0.5), Inches(4.1), Inches(12.33), Inches(2.8), BG_CARD, ORANGE)
    tf3 = box3.text_frame
    tf3.word_wrap = True
    set_text(tf3, "  现有方法不足", size=22, color=ORANGE, bold=True)

    items = [
        ("SG-Nav (2024)", "每步都调LLM (慢+贵), 仿真环境, 无不确定性建模"),
        ("FSR-VLN (2025)", "依赖预建地图, 静态场景, 无信念推理"),
        ("CoW / CLIP-Frontier", "无场景图结构, 无层次推理, 扁平匹配"),
        ("ConceptGraphs (2023)", "离线构建, 无实时更新, 无导航规划"),
    ]
    for name, desc in items:
        add_paragraph(tf3, f"  {name}: {desc}", size=14, color=WHITE)


def make_architecture_slide(prs):
    """Slide 3: 系统架构"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    set_text(txBox.text_frame, "02  系统架构", size=32, color=ACCENT, bold=True)

    # 架构流程 — 用方框 + 箭头模拟
    modules = [
        (0.3,  1.5, 2.4, 1.8, "感知模块\nPerception", "YOLO-World 检测\nCLIP 特征编码\n运动模糊过滤\n实例关联 & 追踪", ACCENT),
        (3.0,  1.5, 2.4, 1.8, "场景图构建\nScene Graph", "Object → Group\nGroup → Room\nRoom → Floor\n在线增量更新", ACCENT2),
        (5.7,  1.5, 2.4, 1.8, "快路径 (Fast)\n<5ms 无LLM", "标签匹配\n检测置信度\n空间关系融合\n多源评分", GOLD),
        (8.4,  1.5, 2.4, 1.8, "慢路径 (Slow)\nLLM CoT推理", "层次化提示\nRoom→Group→Object\n跨语言理解\n常识推理", ORANGE),
        (11.1, 1.5, 2.0, 1.8, "执行 & 导航\nNav2", "局部路径规划\nPure Pursuit\n碰撞避障\n到达检测", ACCENT3),
    ]

    for x, y, w, h, title, desc, color in modules:
        box = add_shape_box(slide, Inches(x), Inches(y), Inches(w), Inches(h), BG_CARD, color)
        tf = box.text_frame
        tf.word_wrap = True
        set_text(tf, f" {title}", size=13, color=color, bold=True)
        for line in desc.split("\n"):
            add_paragraph(tf, f" {line}", size=11, color=LIGHT_GRAY, space_before=Pt(2))

    # 下方: 创新模块
    innovations = [
        (0.3, 3.8, 3.7, 2.0, "信念推理 (BA-HSG)", "Beta分布 → 存在概率\nGaussian → 位置不确定性\n图扩散 → 邻域传播\n复合可信度评分", ACCENT2),
        (4.3, 3.8, 4.2, 2.0, "VoI 自适应调度", "信息价值驱动决策\n高置信 → Continue (77%)\n低置信 → Reperceive (23%)\n避免固定间隔浪费", GOLD),
        (8.8, 3.8, 4.2, 2.0, "拓扑语义探索 (TSG)", "房间拓扑图 + 前沿节点\n信息增益 = 语义×新颖×不确定\nDijkstra 最短路径\n穿越记忆 → 避免重复", ACCENT),
    ]

    for x, y, w, h, title, desc, color in innovations:
        box = add_shape_box(slide, Inches(x), Inches(y), Inches(w), Inches(h), BG_CARD, color)
        tf = box.text_frame
        tf.word_wrap = True
        set_text(tf, f" {title}", size=14, color=color, bold=True)
        for line in desc.split("\n"):
            add_paragraph(tf, f" {line}", size=11, color=LIGHT_GRAY, space_before=Pt(2))

    # 底部: 硬件平台
    box_hw = add_shape_box(slide, Inches(0.3), Inches(6.1), Inches(12.73), Inches(0.8), BG_CARD, LIGHT_GRAY)
    tf_hw = box_hw.text_frame
    tf_hw.word_wrap = True
    set_text(tf_hw, "  硬件平台: Unitree Go2 四足机器人  |  Livox Mid-360 LiDAR  |  Orbbec Femto RGB-D  |  Jetson Orin NX 16GB  |  Fast-LIO2 SLAM", size=14, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)


def make_innovation1_slide(prs):
    """Slide 4: 创新1 — 在线层次场景图"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    set_text(txBox.text_frame, "03  创新1: 在线层次场景图 (BA-HSG)", size=30, color=ACCENT, bold=True)

    # 左侧: 四层结构
    box1 = add_shape_box(slide, Inches(0.5), Inches(1.3), Inches(6), Inches(5.2), BG_CARD, ACCENT)
    tf1 = box1.text_frame
    tf1.word_wrap = True
    set_text(tf1, "  四层层次结构", size=20, color=ACCENT, bold=True)

    layers = [
        ("Floor (楼层)", "全局根节点, 包含所有房间", LIGHT_GRAY),
        ("Room  (房间)", "DBSCAN聚类 + 规则推理命名\n如: corridor, office, kitchen", ACCENT2),
        ("Group (物体组)", "语义聚合 (workstation, safety...)\n基于物体共现关系", GOLD),
        ("Object (物体)", "YOLO-World检测 + CLIP特征\n3D位置 + 信念状态", WHITE),
    ]
    for name, desc, color in layers:
        add_paragraph(tf1, f"  {name}", size=16, color=color, bold=True, space_before=Pt(12))
        for line in desc.split("\n"):
            add_paragraph(tf1, f"    {line}", size=13, color=LIGHT_GRAY, space_before=Pt(2))

    # 右侧: 关键技术
    box2 = add_shape_box(slide, Inches(6.83), Inches(1.3), Inches(6), Inches(5.2), BG_CARD, ACCENT2)
    tf2 = box2.text_frame
    tf2.word_wrap = True
    set_text(tf2, "  关键技术", size=20, color=ACCENT2, bold=True)

    techs = [
        ("质量感知CLIP融合", "高置信检测 → 更大学习率\n低置信检测 → 保守更新\nalpha = base × clamp(score/ref)"),
        ("运动模糊过滤", "四足步态导致图像模糊\nLaplacian方差检测 → 丢弃模糊帧"),
        ("实例关联 & 追踪", "3D距离 + CLIP余弦相似度\n避免同一物体重复注册"),
        ("拓扑连通边", "房间之间的门/过道/邻近关系\n构建可导航空间拓扑"),
    ]
    for name, desc in techs:
        add_paragraph(tf2, f"  {name}", size=15, color=GOLD, bold=True, space_before=Pt(10))
        for line in desc.split("\n"):
            add_paragraph(tf2, f"    {line}", size=12, color=LIGHT_GRAY, space_before=Pt(2))


def make_innovation2_slide(prs):
    """Slide 5: 创新2 — 快慢双路径推理"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    set_text(txBox.text_frame, "04  创新2: 快慢双路径推理 (Fast-Slow)", size=30, color=ACCENT, bold=True)

    # Fast Path
    box1 = add_shape_box(slide, Inches(0.5), Inches(1.3), Inches(6), Inches(3.0), BG_CARD, GOLD)
    tf1 = box1.text_frame
    tf1.word_wrap = True
    set_text(tf1, "  Fast Path (快路径)", size=22, color=GOLD, bold=True)
    add_paragraph(tf1, "  延迟: < 5ms  |  无LLM调用  |  纯算法", size=14, color=ACCENT2, bold=True, space_before=Pt(8))
    fast_items = [
        "1. 标签匹配: 指令关键词 → 物体标签 (精确+模糊)",
        "2. 检测置信度: YOLO-World score × 检测次数",
        "3. 空间关系: near/on/left_of 加权评分",
        "4. CLIP相似度: 跨模态语义匹配 (可选)",
        "5. 多源融合: 加权求和 → 置信度 > 0.75 直接输出",
        "",
        "覆盖率: 100% 英文 L1/L2 指令  |  ~75% 总体",
    ]
    for item in fast_items:
        add_paragraph(tf1, f"  {item}", size=13, color=WHITE if item else LIGHT_GRAY, space_before=Pt(3))

    # Slow Path
    box2 = add_shape_box(slide, Inches(6.83), Inches(1.3), Inches(6), Inches(3.0), BG_CARD, ORANGE)
    tf2 = box2.text_frame
    tf2.word_wrap = True
    set_text(tf2, "  Slow Path (慢路径)", size=22, color=ORANGE, bold=True)
    add_paragraph(tf2, "  延迟: ~26s  |  LLM (Kimi-k2.5)  |  CoT推理", size=14, color=ACCENT3, bold=True, space_before=Pt(8))
    slow_items = [
        "1. 层次化Prompt: 场景图 → 结构化JSON",
        "2. Chain-of-Thought: Room → Group → Object",
        "3. 跨语言: 中文指令 → 英文标签 (100%准确)",
        '4. 语义推理: "去厨房拿水" → 导航到sink',
        "5. 探索建议: 目标不存在 → 建议搜索方向",
        "",
        "覆盖: 中文、复杂空间、多步指令",
    ]
    for item in slow_items:
        add_paragraph(tf2, f"  {item}", size=13, color=WHITE if item else LIGHT_GRAY, space_before=Pt(3))

    # 示例
    box3 = add_shape_box(slide, Inches(0.5), Inches(4.6), Inches(12.33), Inches(2.2), BG_CARD, ACCENT)
    tf3 = box3.text_frame
    tf3.word_wrap = True
    set_text(tf3, "  实测 LLM 输出示例 (Kimi-k2.5)", size=18, color=ACCENT, bold=True)
    examples = [
        ('"找到走廊里门旁边的灭火器"', 'Room: corridor → Group: safety → Object: fire extinguisher (near door, 1.0m) → conf=0.95', '26s'),
        ('"find the printer in office B"', 'Fast Path: label=1.0, det=0.87, spatial=0.3 → fused=0.83 → printer (5.5, -2.5)', '1.4ms'),
        ('"去厨房拿杯水"', 'Room: kitchen → Group: utility → Object: sink (水槽是获取水的最直接来源) → conf=0.92', '27s'),
    ]
    for instr, result, latency in examples:
        add_paragraph(tf3, f"  {instr}", size=13, color=GOLD, bold=True, space_before=Pt(6))
        add_paragraph(tf3, f"    → {result}  [{latency}]", size=11, color=LIGHT_GRAY, space_before=Pt(1))


def make_innovation3_slide(prs):
    """Slide 6: 创新3 — 信念推理 (BA-HSG)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    set_text(txBox.text_frame, "05  创新3: 信念感知场景图 (BA-HSG)", size=30, color=ACCENT, bold=True)

    # 左: 三大信念组件
    box1 = add_shape_box(slide, Inches(0.5), Inches(1.3), Inches(6), Inches(2.8), BG_CARD, ACCENT2)
    tf1 = box1.text_frame
    tf1.word_wrap = True
    set_text(tf1, "  三大信念组件", size=20, color=ACCENT2, bold=True)

    components = [
        ("Beta分布 → 存在概率", "每个物体有 (alpha, beta) 参数\n检测 → alpha++; 未检测 → beta++\nP_exist = alpha / (alpha + beta)"),
        ("Gaussian → 位置不确定度", "均值 = 3D位置, 协方差 = sigma\n多次观测 → sigma 收缩 → 位置更精确"),
        ("图扩散 → 邻域传播", "新物体出现在已知物体附近\n→ 继承邻居的可信度加成"),
    ]
    for name, desc in components:
        add_paragraph(tf1, f"  {name}", size=14, color=GOLD, bold=True, space_before=Pt(8))
        for line in desc.split("\n"):
            add_paragraph(tf1, f"    {line}", size=12, color=LIGHT_GRAY, space_before=Pt(1))

    # 右: VoI调度
    box2 = add_shape_box(slide, Inches(6.83), Inches(1.3), Inches(6), Inches(2.8), BG_CARD, GOLD)
    tf2 = box2.text_frame
    tf2.word_wrap = True
    set_text(tf2, "  VoI 自适应调度", size=20, color=GOLD, bold=True)
    add_paragraph(tf2, "  信息价值 (Value of Information) 驱动决策:", size=14, color=WHITE, space_before=Pt(8))
    voi_items = [
        ("Continue (77%)", "可信度高, 继续导航, 不浪费算力"),
        ("Reperceive (23%)", "可信度低, 重新感知, 更新场景图"),
        ("Slow Reason (罕见)", "场景图不足, 调用LLM深度推理"),
    ]
    for name, desc in voi_items:
        add_paragraph(tf2, f"  {name}", size=14, color=ACCENT, bold=True, space_before=Pt(8))
        add_paragraph(tf2, f"    {desc}", size=12, color=LIGHT_GRAY, space_before=Pt(1))

    # 底部: 多假设规划
    box3 = add_shape_box(slide, Inches(0.5), Inches(4.4), Inches(12.33), Inches(2.4), BG_CARD, ACCENT3)
    tf3 = box3.text_frame
    tf3.word_wrap = True
    set_text(tf3, "  多假设目标规划 (Multi-Hypothesis)", size=20, color=ACCENT3, bold=True)
    add_paragraph(tf3, "  场景中有多个同类物体 (如3个灭火器) 时:", size=14, color=WHITE, space_before=Pt(8))
    mh_items = [
        "  1. 按后验概率排序候选目标",
        "  2. 导航到最高概率候选 → 到达后验证",
        "  3. 未确认 → 贝叶斯更新拒绝, 选择下一个",
        "  4. 平均 1.6 次尝试找到正确目标 (vs 随机3.0次)",
        "  5. 100% 成功率 (20/20 场景)",
    ]
    for item in mh_items:
        add_paragraph(tf3, item, size=13, color=WHITE, space_before=Pt(3))


def make_innovation4_slide(prs):
    """Slide 7: 创新4 — 拓扑语义探索 (TSG)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    set_text(txBox.text_frame, "06  创新4: 拓扑语义探索 (TSG)", size=30, color=ACCENT, bold=True)

    # 左: TSG 结构
    box1 = add_shape_box(slide, Inches(0.5), Inches(1.3), Inches(4.0), Inches(3.5), BG_CARD, ACCENT)
    tf1 = box1.text_frame
    tf1.word_wrap = True
    set_text(tf1, "  拓扑语义图 (TSG)", size=18, color=ACCENT, bold=True)

    tsg_items = [
        "Room Nodes: 已知房间",
        "  → 类型/访问状态/物体列表",
        "Frontier Nodes: 未探索前沿",
        "  → 位置/方向/预测房间类型",
        "Edges: 房间连通关系",
        "  → 门/过道/邻近/穿越",
        "Traversal Memory: 穿越记录",
        "  → 避免重复探索",
    ]
    for item in tsg_items:
        add_paragraph(tf1, f"  {item}", size=12, color=WHITE if not item.startswith("  →") else LIGHT_GRAY, space_before=Pt(3))

    # 中: Algorithm 2
    box2 = add_shape_box(slide, Inches(4.8), Inches(1.3), Inches(4.0), Inches(3.5), BG_CARD, ACCENT2)
    tf2 = box2.text_frame
    tf2.word_wrap = True
    set_text(tf2, "  Algorithm 2: IG 探索", size=18, color=ACCENT2, bold=True)

    algo_items = [
        "IG(n) = S_sem x N x U",
        "",
        "S_sem: 语义先验",
        "  咖啡机→kitchen (0.9)",
        "  灭火器→stairwell (0.9)",
        "N: 新颖度",
        "  未访问=1.0, 已访问→衰减",
        "U: 不确定度降低",
        "  前沿 > 稀疏房间 > 密集房间",
        "",
        "Score = IG / (1 + 0.3*d_path)",
        "  d_path = Dijkstra 最短路径",
    ]
    for item in algo_items:
        c = GOLD if "IG(n)" in item or "Score" in item else WHITE
        add_paragraph(tf2, f"  {item}", size=12, color=c, bold=("IG(n)" in item or "Score" in item), space_before=Pt(2))

    # 右: 双层策略
    box3 = add_shape_box(slide, Inches(9.1), Inches(1.3), Inches(3.73), Inches(3.5), BG_CARD, GOLD)
    tf3 = box3.text_frame
    tf3.word_wrap = True
    set_text(tf3, "  双层探索策略", size=18, color=GOLD, bold=True)

    dual_items = [
        "Layer 1: TSG (~1ms)",
        "  纯算法, 无LLM",
        "  覆盖 ~70% 场景",
        "  信息增益驱动",
        "",
        "Layer 2: LLM (fallback)",
        "  复杂/歧义场景",
        "  注入拓扑摘要",
        "  常识推理增强",
        "",
        "穿越记忆:",
        "  去过的房间 IG↓ 90%",
        "  LLM 利用负面记忆",
    ]
    for item in dual_items:
        c = ACCENT if "Layer" in item or "穿越" in item else LIGHT_GRAY
        b = "Layer" in item or "穿越" in item
        add_paragraph(tf3, f"  {item}", size=12, color=c, bold=b, space_before=Pt(2))

    # 底部: 实测对比
    box4 = add_shape_box(slide, Inches(0.5), Inches(5.1), Inches(12.33), Inches(1.8), BG_CARD, ORANGE)
    tf4 = box4.text_frame
    tf4.word_wrap = True
    set_text(tf4, "  TSG vs LLM 实测对比 (Kimi-k2.5)", size=16, color=ORANGE, bold=True)

    compare_items = [
        ('"找到咖啡机"', 'TSG: stairwell(0.184) > kitchen(0.176)', 'LLM: kitchen (语义常识更准)'),
        ('"find fire ext on 2nd floor"', 'TSG: stairwell(0.414) 大幅领先', 'LLM: stairwell (完全一致!)'),
        ('"找到会议室投影仪"', 'TSG: frontier_10000(0.136) 选未知前沿', 'LLM: 东侧前沿 (一致!)'),
    ]
    for instr, tsg_r, llm_r in compare_items:
        add_paragraph(tf4, f"  {instr}   TSG: {tsg_r.split(': ')[1]}  |  LLM: {llm_r.split(': ')[1]}", size=11, color=WHITE, space_before=Pt(4))


def make_experiments_slide(prs):
    """Slide 8: 实验验证"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    set_text(txBox.text_frame, "07  实验验证", size=32, color=ACCENT, bold=True)

    # 测试总览
    box1 = add_shape_box(slide, Inches(0.5), Inches(1.2), Inches(4.0), Inches(2.5), BG_CARD, ACCENT2)
    tf1 = box1.text_frame
    tf1.word_wrap = True
    set_text(tf1, "  测试总览", size=18, color=ACCENT2, bold=True)
    stats = [
        ("总测试数", "99"),
        ("通过率", "100% (99/99)"),
        ("拓扑图测试", "33/33"),
        ("离线管线测试", "40/40"),
        ("信念系统测试", "26/26"),
    ]
    for name, val in stats:
        add_paragraph(tf1, f"  {name}: {val}", size=14, color=WHITE, space_before=Pt(6))

    # Fast Path 性能
    box2 = add_shape_box(slide, Inches(4.8), Inches(1.2), Inches(4.0), Inches(2.5), BG_CARD, GOLD)
    tf2 = box2.text_frame
    tf2.word_wrap = True
    set_text(tf2, "  Fast Path 性能", size=18, color=GOLD, bold=True)
    fp_stats = [
        ("L1 英文命中率", "100% (20/20)"),
        ("L2 英文命中率", "100% (15/15)"),
        ("平均延迟", "< 5ms"),
        ("P99 延迟", "< 20ms"),
        ("LLM调用节省", "> 75%"),
    ]
    for name, val in fp_stats:
        add_paragraph(tf2, f"  {name}: {val}", size=14, color=WHITE, space_before=Pt(6))

    # LLM 验证
    box3 = add_shape_box(slide, Inches(9.1), Inches(1.2), Inches(3.73), Inches(2.5), BG_CARD, ORANGE)
    tf3 = box3.text_frame
    tf3.word_wrap = True
    set_text(tf3, "  Slow Path (Kimi-k2.5)", size=16, color=ORANGE, bold=True)
    sp_stats = [
        ("端到端通过率", "100% (12/12)"),
        ("跨语言准确率", "100%"),
        ("探索准确率", "100%"),
        ("平均延迟", "~26s"),
        ("平均置信度", "0.93"),
    ]
    for name, val in sp_stats:
        add_paragraph(tf3, f"  {name}: {val}", size=13, color=WHITE, space_before=Pt(6))

    # 底部: 关键发现
    box4 = add_shape_box(slide, Inches(0.5), Inches(4.0), Inches(12.33), Inches(2.8), BG_CARD, ACCENT)
    tf4 = box4.text_frame
    tf4.word_wrap = True
    set_text(tf4, "  关键发现", size=20, color=ACCENT, bold=True)

    findings = [
        "1. 快路径 100% 命中英文L1/L2 → 无需LLM, <5ms 完成, 适合实时导航",
        '2. LLM 完美跨语言: "椅子"→chair, "灭火器"→fire extinguisher (100%准确)',
        '3. LLM 语义推理: "去厨房拿水" → 导航到 sink (无cup, 推理出sink是水源)',
        "4. 信念系统: 30帧后 P_exist 从0.6→0.9, 图扩散有效, 过期物体可信度↓",
        "5. 多假设: 平均1.6次找到正确目标 (vs 随机3.0次), 100%成功",
        "6. VoI调度: 77% continue / 23% reperceive, 自适应比固定间隔更高效",
        "7. TSG探索: 穿越记忆使已访问房间IG↓90%, LLM利用拓扑摘要避开已探索方向",
    ]
    for f in findings:
        add_paragraph(tf4, f"  {f}", size=12, color=WHITE, space_before=Pt(3))


def make_comparison_slide(prs):
    """Slide 9: 与现有工作对比"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    set_text(txBox.text_frame, "08  与现有工作对比", size=32, color=ACCENT, bold=True)

    # 表头
    headers = ["维度", "SG-Nav (2024)", "FSR-VLN (2025)", "HSG-Nav (Ours)"]
    col_widths = [2.5, 3.2, 3.2, 3.9]
    x_start = 0.3

    y = 1.3
    x = x_start
    for i, (h, w) in enumerate(zip(headers, col_widths)):
        box = add_shape_box(slide, Inches(x), Inches(y), Inches(w), Inches(0.5), ACCENT, ACCENT)
        set_text(box.text_frame, f" {h}", size=14, color=BG_DARK, bold=True, alignment=PP_ALIGN.CENTER)
        x += w + 0.05

    rows = [
        ["环境", "仿真 (MP3D)", "预建地图, 静态", "未知, 动态, 在线"],
        ["场景图", "3层, 确定性", "4层, 离线", "4层, 在线, 信念感知"],
        ["目标解析", "每步LLM", "CLIP+VLM", "Fast(<5ms) + Slow(LLM)"],
        ["LLM调用", "每步", "VLM精筛", "仅~25%查询"],
        ["不确定性", "标量可信度", "无", "Beta+Gaussian+VoI"],
        ["探索", "LLM-only", "无", "TSG信息增益+LLM"],
        ["多步指令", "不支持", "单步", "支持(任务分解)"],
        ["跨语言", "未报告", "未报告", "100% (中→英)"],
        ["平台", "仿真 (90cm)", "Go1/G1 (桌面GPU)", "Go2 (Jetson, 30cm)"],
        ["SR", "40-54%", "92%(预建图)", "100% (12/12, 离线)"],
    ]

    for row_data in rows:
        y += 0.52
        x = x_start
        for j, (cell, w) in enumerate(zip(row_data, col_widths)):
            bg = BG_CARD
            c = WHITE
            if j == 3:
                bg = RGBColor(0x0A, 0x2A, 0x1A)
                c = ACCENT2
            box = add_shape_box(slide, Inches(x), Inches(y), Inches(w), Inches(0.5), bg)
            tf = box.text_frame
            tf.word_wrap = True
            set_text(tf, f" {cell}", size=11, color=c if j == 3 else (GOLD if j == 0 else LIGHT_GRAY), bold=(j==0 or j==3), alignment=PP_ALIGN.CENTER)
            x += w + 0.05


def make_summary_slide(prs):
    """Slide 10: 总结与展望"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    set_text(txBox.text_frame, "09  总结与展望", size=32, color=ACCENT, bold=True)

    # 左: 总结
    box1 = add_shape_box(slide, Inches(0.5), Inches(1.3), Inches(6), Inches(4.5), BG_CARD, ACCENT2)
    tf1 = box1.text_frame
    tf1.word_wrap = True
    set_text(tf1, "  核心贡献", size=22, color=ACCENT2, bold=True)

    contributions = [
        ("在线层次场景图 (BA-HSG)", "4层结构 + 信念感知 + 在线增量更新"),
        ("快慢双路径推理", "Fast <5ms 覆盖75%+ | Slow LLM 处理复杂指令"),
        ("信念推理 & VoI调度", "Beta/Gaussian概率建模 + 自适应感知调度"),
        ("多假设目标规划", "1.6次平均找到目标, 100%成功率"),
        ("拓扑语义探索 (TSG)", "信息增益 + Dijkstra + 穿越记忆"),
        ("跨语言支持", "中英文指令100%准确, LLM跨语言推理"),
    ]
    for name, desc in contributions:
        add_paragraph(tf1, f"  {name}", size=14, color=GOLD, bold=True, space_before=Pt(8))
        add_paragraph(tf1, f"    {desc}", size=12, color=LIGHT_GRAY, space_before=Pt(1))

    # 右: 展望
    box2 = add_shape_box(slide, Inches(6.83), Inches(1.3), Inches(6), Inches(4.5), BG_CARD, ACCENT)
    tf2 = box2.text_frame
    tf2.word_wrap = True
    set_text(tf2, "  下一步计划", size=22, color=ACCENT, bold=True)

    future = [
        ("真机实验 Phase 1", "Go2 + 办公走廊, L1-L3 × 3试次\n采集场景图JSON, 导航轨迹, 视频"),
        ("消融实验", "w/o 场景图 | w/o 层次 | w/o 信念\nw/o VoI | w/o 多假设 | w/o 再感知"),
        ("多楼层扩展", "TSG楼梯间节点 → 跨楼层导航"),
        ("动态场景", "物体移动检测 + 场景图更新"),
        ("端侧LLM部署", "量化LLM (4-bit) 部署到Jetson\n消除云端依赖, 降低延迟"),
    ]
    for name, desc in future:
        add_paragraph(tf2, f"  {name}", size=14, color=ACCENT2, bold=True, space_before=Pt(8))
        for line in desc.split("\n"):
            add_paragraph(tf2, f"    {line}", size=12, color=LIGHT_GRAY, space_before=Pt(1))

    # 底部
    box3 = add_shape_box(slide, Inches(2), Inches(6.1), Inches(9.33), Inches(0.7), BG_CARD, GOLD)
    set_text(box3.text_frame, "  99 Tests Passed  |  6 Innovations  |  100% Offline SR  |  Ready for Real-Robot", size=16, color=GOLD, bold=True, alignment=PP_ALIGN.CENTER)


def make_tech_stack_slide(prs):
    """Slide 11: 技术栈全景"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    set_text(txBox.text_frame, "10  技术栈全景", size=32, color=ACCENT, bold=True)

    stacks = [
        (0.5, 1.3, "机器人端 (ROS2)", [
            "Fast-LIO2 SLAM (IMU+LiDAR)",
            "Nav2 + DWB 局部规划",
            "PCT Global Planner (GTSAM)",
            "Terrain Analysis (地形感知)",
            "gRPC Gateway (44 RPCs)",
            "OTA Daemon (Ed25519签名)",
        ], ACCENT),
        (4.8, 1.3, "AI 模块 (Python)", [
            "YOLO-World (开放词汇检测)",
            "CLIP ViT-B/32 (视觉编码)",
            "层次场景图 (instance_tracker)",
            "快慢推理 (goal_resolver)",
            "信念系统 (belief_system)",
            "拓扑探索 (topology_graph)",
        ], ACCENT2),
        (9.1, 1.3, "客户端 (Flutter)", [
            "MapPilot 多平台App",
            "14 Providers 状态管理",
            "5 Gateway 连接管理",
            "实时地图可视化",
            "任务下发 & 监控",
            "BLE 紧急通信",
        ], GOLD),
    ]

    for x, y, title, items, color in stacks:
        box = add_shape_box(slide, Inches(x), Inches(y), Inches(3.9), Inches(3.0), BG_CARD, color)
        tf = box.text_frame
        tf.word_wrap = True
        set_text(tf, f"  {title}", size=16, color=color, bold=True)
        for item in items:
            add_paragraph(tf, f"  {item}", size=12, color=WHITE, space_before=Pt(4))

    # LLM 支持
    box_llm = add_shape_box(slide, Inches(0.5), Inches(4.6), Inches(12.33), Inches(1.5), BG_CARD, ORANGE)
    tf_llm = box_llm.text_frame
    tf_llm.word_wrap = True
    set_text(tf_llm, "  LLM 后端支持", size=18, color=ORANGE, bold=True)
    llm_items = [
        "OpenAI GPT-4o / GPT-4o-mini  |  Anthropic Claude 3.5  |  Alibaba Qwen (通义千问)  |  Moonshot Kimi-k2.5",
        "统一异步接口 (LLMClientBase)  |  自动重试 + 超时  |  可切换后端  |  温度/模型可配置",
        "已验证: Kimi-k2.5 全部12用例通过, 跨语言100%, 探索推理合理",
    ]
    for item in llm_items:
        add_paragraph(tf_llm, f"  {item}", size=12, color=WHITE, space_before=Pt(4))

    # 代码量
    box_code = add_shape_box(slide, Inches(0.5), Inches(6.3), Inches(12.33), Inches(0.6), BG_CARD, LIGHT_GRAY)
    set_text(box_code.text_frame, "  17 ROS2 Packages  |  44 gRPC RPCs  |  5 Proto Files  |  99 Unit Tests  |  Python + C++ + Dart + Protobuf", size=14, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)


def make_thankyou_slide(prs):
    """Slide 12: 谢谢"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Pt(4))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()

    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.33), Inches(2))
    tf = txBox.text_frame
    set_text(tf, "Thank You!", size=54, color=ACCENT, bold=True, alignment=PP_ALIGN.CENTER)
    add_paragraph(tf, "HSG-Nav: 基于层次场景图的四足机器人自主导航", size=24, color=WHITE, alignment=PP_ALIGN.CENTER, space_before=Pt(20))
    add_paragraph(tf, "Questions & Discussion", size=20, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER, space_before=Pt(16))


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    make_title_slide(prs)
    make_background_slide(prs)
    make_architecture_slide(prs)
    make_innovation1_slide(prs)
    make_innovation2_slide(prs)
    make_innovation3_slide(prs)
    make_innovation4_slide(prs)
    make_experiments_slide(prs)
    make_comparison_slide(prs)
    make_summary_slide(prs)
    make_tech_stack_slide(prs)
    make_thankyou_slide(prs)

    out_path = os.path.join(os.path.dirname(__file__), "..", "docs", "HSG-Nav_Presentation.pptx")
    os.makedirs(os.path.dirname(out_path), exist_ok=True)
    prs.save(out_path)
    print(f"PPT saved to: {os.path.abspath(out_path)}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
